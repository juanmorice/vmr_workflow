#!/usr/bin/env python
# coding: utf-8

# ## VMR Scorecard Version - (Segment definition options) (last update: 2023-11-08)
# 
# ### Developer: Juan Morice
# 
# ##### *UPDATES: This code version is allowed to:*
# 
# * Run "Reading_SharePoint_Lists.ipynb" script to import the parameters loaded by Sales Team through both Power Apps and Power Automate developed tool.
# * Define the segments you want to use for reporting the metrics: 
# 
#         1. UPC hierarchy groups
#         2.YB brand description (trademark_brand_desc)
#         3. YB category description (cat_desc)
#         4. Custom Brand's descriptions from mapping file provided. A custom file with Brand's descriptions (first column) and cat_desc (second column) fields is needed.
#         5. Retailer's descriptions from trade_item_owner_hierarchy_v table (except Kroger). A UPC list based on one specific retailer is needed.
#         6. Kroger's own descriptions from nz_hierarchy table. A UPC list based only in Kroger is needed.
#         
#         
# * Define automatically the static parameters based on the number of retailers involved.
# * Handle cases when you have NEW UPCs, so YAGO metrics cannot be shown.
# * Deliver both top 10 or all segment combinations, regarding what the analyst or Director wants to report.
# * Create two final ouputs for delivering: formatted Excel data and populated Power Point slides.


# In[1]:


# import modules
import sys
import os
from pathlib import Path
# Add src folder to path for local_modules imports
# Works both when run directly and via exec()
if '__file__' in dir():
    src_path = Path(__file__).resolve().parent.parent
else:
    # When run via exec(), use current working directory structure
    src_path = Path.cwd() / 'src'
if str(src_path) not in sys.path:
    sys.path.insert(0, str(src_path))

import pandas as pd
import numpy as np
import datetime as dt
import psycopg2 as ybconn


# In[2]:


import getpass
from local_modules.safe_password import readpw
import psycopg2 as ybconn
import pandas as pd
from datetime import date
from local_modules.yb_load import yb_load, yb_load_file


#BEFORE STARTING: DOUBLE-CHECK that sharepoint_credentials.txt file contents your own credential for accessing to Sharepoint.


dbase = 'py1usta1'                  # database on server for connection string below


# In[4]:


############################## DON'T CHANGE ################################## 

all_combos = False         

#PowerPoint final deck:
template_path = "/opt/airflow/src/templates/VMR_Scorecard_Template.pptx"   
#The code will take the PPTX template to build the final deck from this path.


#Modified from K.Mertens's J&J custom VMR code.


# connection string to server - this is quick connection to check UPCs
yb_user = os.getenv('YELLOWBRICK_USER', getpass.getuser())  # Use env var or fallback to system user
conn = ybconn.connect(user=yb_user, password=readpw("Yellowbrick"), 
                      host='orlpybvip01.catmktg.com', port='5432', 
                      database = f'{dbase}')
conn.set_session(autocommit=True)
curs = conn.cursor()

 

df_upcs_quick = pd.read_sql(f'''

    select brandnbr as brand_nbr
      ,branddesc as brand_desc
      ,count(distinct tradeitemcd) as nbr_of_upcs
    from VMR_{brand_nm}_upclmc_{analyst}
    group by 1,2
''',conn).sort_values(by=['brand_nbr']).reset_index().drop(columns=['index'])


brand_nbr_str = ",".join(str(i) for i in brand_nbr)
cat_nbr_str = ",".join(str(i) for i in cat_nbr)


#Segment Definition:

place_holder_1 = ''
place_holder_2 = ''
place_holder_3 = ''
place_holder_4 = ''
place_holder_5 = ''


if segment_type == 1:
    segment_def = 'brand_desc'

if segment_type == 2:
    segment_def = 'trademark_brand_desc'
    
if segment_type == 3:
    segment_def = 'cat_desc'
    
if segment_type == 6:
    
    upload_descr = pd.read_excel(custom_brand_desc_path)
    upload_descr = upload_descr.replace(np.nan, 'All Other Categories', regex=True)
    columns_nm = list(upload_descr.columns)
    upload_descr = upload_descr.rename(columns={columns_nm[0]:'custombranddescr',columns_nm[1]:'catdesc'})
    
    
    yb_load(Df = upload_descr,
            table_name = f'''VMR_{brand_nm}_brand_descriptions_{analyst}''',
            userid = yb_user,
            passwd = readpw("Yellowbrick"),
            append = False,
            database= f'{dbase}')
    
    place_holder_1 = f''', f.custombranddescr as custom_brand_desc '''
    place_holder_2 = f'''inner join VMR_{brand_nm}_brand_descriptions_{analyst} f ON (f.catdesc = p.cat_desc)'''
    segment_def = f'''custom_brand_desc'''
    

if segment_type == 4:
    
    lb_def = f'''l{label_level}'''
    
    place_holder_3 = f''',z.owner_nm'''
    
    place_holder_4 = f''',z.trade_item_hier_{lb_def}_desc'''
    
    place_holder_5 = f'''inner join trade_item_owner_hierarchy_v z on p.trade_item_key = z.trade_item_key'''
        
    segment_def = f'''trade_item_hier_{lb_def}_desc'''
    
    
if segment_type == 5:
    
    lb_def = f'''l{label_level}'''
    
    place_holder_4 = f''',z.tradeitemhier{lb_def}desc'''
    
    place_holder_5 = f'''inner join py1ussa1.public.nz_hierarchy z on p.trade_item_key = z.tradeitemkey'''
        
    segment_def = f'''tradeitemhier{lb_def}desc'''

    
    
## Converting "BL-..." format to "USA-BLIP-BL...":
for i in range(len(BL_CODES)):
    
    if BL_CODES[i][0] == 'B' and BL_CODES[i][1] == 'L' and BL_CODES[i][2] == '_':
        
        BL_CODES[i] = str(''.join(('USA-BLIP-',BL_CODES[i])) )

## Converting Announcement "BL-..." format to "USA-BLIP-BL...":
if Announcement != ['']:
    for i in range(len(Announcement)):
        if Announcement[i][0] == 'B' and Announcement[i][1] == 'L' and Announcement[i][2] == '_':
            Announcement[i] = str(''.join(('USA-BLIP-',Announcement[i])) )

BL_CODES = str(BL_CODES) 
BL_CODES = BL_CODES.replace("[","").replace("]","")
print("BL_CODES: " + BL_CODES)

Announcement = str(Announcement) 
Announcement = Announcement.replace("[","").replace("]","")
print("Annoucement: " + Announcement)


Min_Threshold_Value = pd.read_sql(f''' SELECT distinct cnsmr_trg_class_purch_rqmt_qty as min, rolling_wk_vld_qty as rolling
                 FROM promotion_v
                 WHERE promo_src_id_txt IN ({BL_CODES})
                 ORDER BY 1 DESC
''',conn)

print(Min_Threshold_Value)

# Calculate min_threshold_value from database
if not np.isnan(Min_Threshold_Value['min'].min()):
    min_threshold_value = int(Min_Threshold_Value['min'].min())
else:
    min_threshold_value = 0

# If calculated min_threshold_value is 0 or null and user provided a value, use user input
if min_threshold_value == 0 and user_min_threshold is not None:
    min_threshold_value = user_min_threshold
    print(f'Using user-provided min_threshold_value: {min_threshold_value}')
else:
    print(f'Using calculated min_threshold_value: {min_threshold_value}')

min_threshold_statement = f''' {threshold_unit} >= {str(min_threshold_value)} '''

print(min_threshold_statement)

# Calculate redemption_days from database
if not np.isnan(Min_Threshold_Value['rolling'].max()):
    redemption_weeks = int(Min_Threshold_Value['rolling'].max())
else:
    redemption_weeks = 2
    
calculated_redemption_days = redemption_weeks * 7

# If user provided redemption_days, use that value; otherwise use calculated value
if user_redemption_days is not None:
    redemption_days = user_redemption_days
    print(f'Using user-provided redemption_days: {redemption_days} (calculated was: {calculated_redemption_days})')
else:
    redemption_days = calculated_redemption_days
    print(f'Using calculated redemption_days: {redemption_days}')


print('UPC Quick Check: ')
print(df_upcs_quick)
print("")
print('Quick Check: ')
print('The BL codes to perfom the analysis are: ' + str(BL_CODES))
print("Focus brand(s):",brand_nbr_str)
print("Category brands:",cat_nbr_str)
print("Segment definition string:",segment_def)
print(" ")

if segment_type == 6:
    print("Custom Brand's descriptions uploaded:")
    print(upload_descr)
    


# In[5]:


#Renaming dictionary: (took from K.Mertens J&J VMR code)

renaming_dict = {
    'ord_event_key' : 'ord_event_key',
    'cnsmr_id_key' : 'cnsmr_id_key',
    'analysis_period_p1p2': 'Analysis Period', 
    'analysis_period_rec52': 'Analysis Period', 
    'cal_sun_wk_ending_dt': 'WE Date',
    'nbr_of_upcs': 'Raw # of UPCs included',
    'brand_desc': 'Brand',
    'sur_seg': 'Purchasing Segment',
    'price_seg': 'Pricing Segment',
    'fin_cmit_contract_nbr': 'Contract Nbr',
    'promo_src_id_txt': 'BL',
    'analysis_periods':'Analysis Periods',
    'start':'Start', 'end':'End',
    'count_distinct_trips':'Count distinct trips',
    'dollar_sales':'Dollar Sales',
    'dollars_per_trip':'Dollars Per Trip',
    '%_change_dollars_per_trip':'% Change Dollars Per Trip',
    'units':'Units',
    'dollars':'Dollars',
    'trips':'Trips',
    'units_per_trip':'Units Per Trip',
    'dollars_per_trip_yago_period':'Dollars Per Trip YAGO Period',
    'dollars_per_trip_vmr_period':'Dollars Per Trip VMR Period',
    'units_per_trip_yago_period':'Units Per Trip YAGO Period',
    'units_per_trip_vmr_period':'Units Per Trip VMR Period',
    'dollars_per_trip_pre_period':'Dollars Per Trip Pre Period',
    'units_per_trip_pre_period':'Units Per Trip Pre Period',
   
}


# ## **2) DIMESIONAL TABLES AND DATE CALCULATIONS:**

# 2a) Creating Basic Dimensional tables:

# In[6]:


#Modified from K.Mertens's J&J custom VMR code.

#FILTERING BY REWARD PRINTS (BLs):

curs.execute(f'''

    DROP table if exists VMR_{brand_nm}_promo_filter;
    CREATE temp table VMR_{brand_nm}_promo_filter as
    
        SELECT distinct pv.promo_src_id_txt
            ,pv.promo_src_id
            ,pv.promo_varnt_key
            ,pv.promo_key
            
        FROM promotion_variant_v pv 
        
        WHERE pv.promo_src_id_txt in ({BL_CODES})
        
        DISTRIBUTE replicate;
        
''')

#FILTERING BY TOUCH POINT FIELDS:

curs.execute(f'''

    DROP table if exists VMR_{brand_nm}_touchpoint_filter;       
    CREATE temp table VMR_{brand_nm}_touchpoint_filter as
        
        SELECT tp.touchpoint_key 
               ,tp.acct_specific_rtlr_nbr 
               ,tp.acct_specific_rtlr_nm 
               ,tp.site_key 
               ,tp.ntwk_id 
               ,tp.ntwk_nm 
               ,tp.csdb_chn_nbr
               
        FROM touchpoint_v tp 

        DISTRIBUTE ON(touchpoint_key);      
        
''')

#GETTING DETAILS ABOUT EVERY TRIP FILTERING BY BLS:

curs.execute(f''' 

    DROP table if exists VMR_{brand_nm}_print_dt_actual;
    CREATE temp table VMR_{brand_nm}_print_dt_actual as

        SELECT
             p.promo_src_id_txt
            ,p.promo_src_id
            ,a.event_typ_cd
            ,a.ord_date_key
            ,d.cal_dt
            ,tp.ntwk_id
            ,tp.ntwk_nm
            ,tp.acct_specific_rtlr_nbr
            ,tp.acct_specific_rtlr_nm
            ,tp.site_key
            ,tp.csdb_chn_nbr
            ,tp.touchpoint_key 
            ,sum(a.tot_selected_qty) as prints

        FROM ord_promo_varnt_cnsmr_ne_v a 

            INNER JOIN VMR_{brand_nm}_promo_filter p ON (a.promo_varnt_key=p.promo_varnt_key)
            INNER JOIN VMR_{brand_nm}_touchpoint_filter tp ON (a.ord_touchpoint_key=tp.touchpoint_key)
            INNER JOIN date_v d ON (d.date_key = a.ord_date_key)

        WHERE a.event_typ_cd in ('IS-SPRINT-T','IS-PRINT-T','STCNTL_CMPGN','STCNTL_SITE','IS-SPRINT-vd','LNG_TRM_CNTL')

        GROUP BY 1,2,3,4,5,6,7,8,9,10,11,12

        DISTRIBUTE ON(ord_date_key);
    
''')

#CREATING UPC TABLE:

curs.execute(f'''
    drop table if exists VMR_{brand_nm}_upc_filter;       
    create temp table VMR_{brand_nm}_upc_filter as
    
        select s.brandnbr as brand_nbr
              ,s.branddesc as brand_desc
              ,p.trademark_brand_desc
              ,p.trade_item_key
              ,p.trade_item_cd
              ,p.cat_nbr
              ,p.cat_desc
              {place_holder_3}
              {place_holder_4}
              {place_holder_1}
        from trade_item_v p 
              INNER JOIN VMR_{brand_nm}_upclmc_{analyst} s ON (p.trade_item_cd = s.tradeitemcd)
              {place_holder_2}
              {place_holder_5}
            
        distribute replicate;
''')


#GETTING PROMO DETAILS FOR SUMMARY TABLE:

curs.execute(f''' 

    DROP table if exists VMR_{brand_nm}_promo_check;
    CREATE temp table VMR_{brand_nm}_promo_check as
    
        SELECT distinct fin_cmit_contract_nbr 
               ,fin_cmit_contract_nm
               ,promo_src_id_txt
               ,promo_desc_txt
              ,promo_funding_typ_cd
              ,promo_discnt_val_amt
              ,rolling_wk_vld_qty
        
        FROM promotion_variant_v pv 
        
        WHERE pv.promo_src_id_txt in ({BL_CODES})
        
    DISTRIBUTE REPLICATE;
    
''')

#CREATING PROMO SUMMARY TABLE:

curs.execute(f''' 

    DROP table if exists VMR_{brand_nm}_promo_summary;
    CREATE temp table VMR_{brand_nm}_promo_summary as

        SELECT p.fin_cmit_contract_nm
               ,p.fin_cmit_contract_nbr
               ,p.promo_src_id_txt
               ,p.promo_desc_txt
               ,min(cal_dt) as actual_start
               ,max(cal_dt) as actual_end
               ,sum(prints) as total_reward_prints
              
        FROM VMR_{brand_nm}_print_dt_actual a
        
            INNER JOIN VMR_{brand_nm}_promo_check p on a.promo_src_id_txt = p.promo_src_id_txt
            
        GROUP BY 1,2,3,4

        DISTRIBUTE REPLICATE;
    
''')


#Dimesional table that pull-outs print dates for every chain:

curs.execute(f''' 
    DROP table if exists VMR_{brand_nm}_actual_dates_by_chain;
    CREATE temp table VMR_{brand_nm}_actual_dates_by_chain as
    
            SELECT csdb_chn_nbr
                  ,min(cal_dt) as actual_start
                  ,max(cal_dt) as actual_end
                  ,min(ord_date_key) as min_actual_date_key
                  ,max(ord_date_key) as max_actual_date_key
            
            FROM VMR_{brand_nm}_print_dt_actual a
            
            GROUP BY 1
            
            DISTRIBUTE REPLICATE;
''',conn)


#Dimensional table that pull-outs the printing stores from chains filtered in the last dimensiona table:

curs.execute(f''' 
    DROP table if exists VMR_{brand_nm}_printing_stores;
    CREATE temp table VMR_{brand_nm}_printing_stores as 
            
            SELECT distinct tp.site_key,tp.touchpoint_key, tp.acct_specific_rtlr_nm, a.min_actual_date_key, a.max_actual_date_key

            FROM touchpoint_v tp 

                INNER JOIN VMR_{brand_nm}_actual_dates_by_chain a on tp.csdb_chn_nbr = a.csdb_chn_nbr

            DISTRIBUTE ON(touchpoint_key);
''')


# 2b) Identifying number of retailers involved to define static parameters:

# In[7]:


#Getting the number of different retailers involved:

nbr_retailers = pd.read_sql(f''' SELECT count(distinct acct_specific_rtlr_nm)
                                 FROM VMR_{brand_nm}_printing_stores
''',conn)

nbr_retailers = nbr_retailers['count'].values[0]
print("Number of retailers for this analysis: " + str(nbr_retailers))


#Getting the retailers involved in the analysis:

retailers = pd.read_sql(f''' SELECT distinct acct_specific_rtlr_nm FROM VMR_{brand_nm}_printing_stores''',conn)
retailers = list(retailers['acct_specific_rtlr_nm'])
retailers = str(retailers) 
retailers = retailers.replace("[","").replace("]","")
retailers = retailers.replace("'","")
print("Retailers involved in this analysis:", retailers)


#If segment_type == 4, we are working with retailer's descriptions table,
#we need to filter by owner_nm field, which is like the acct_specific_rtlr_nm but written different, so
#we have to identify first what is the owner_nm for the retailers we are working with. 

if segment_type == 4:
    
    retail_owner_df = pd.read_sql(f'''SELECT acct_specific_rtlr_nm, owner_nm
                                      FROM ord_trd_itm_cnsmr_fact_ne_v a
                                      INNER JOIN VMR_{brand_nm}_upc_filter b ON (a.trade_item_key = b.trade_item_key)
                                      INNER JOIN VMR_{brand_nm}_printing_stores c ON (a.ord_touchpoint_key = c.touchpoint_key)
                                      GROUP BY 1,2
    ''',conn)

    owner_list = []

    #Getting the list of unique retailers (acct_specific_rtlr_nm') involved in the analysis
    retail_unique = list(retail_owner_df['acct_specific_rtlr_nm'].unique())

    for i in range(len(retail_unique)): 
        for j in range(len(retail_owner_df['owner_nm'])):
            if retail_owner_df['owner_nm'][j] in retail_owner_df['acct_specific_rtlr_nm'][i]:   
                owner_list.append(retail_owner_df['owner_nm'][j])
                
                
    owner_list = str(owner_list) 
    owner_list = owner_list.replace("[","").replace("]","")
    #owner_list = owner_list.replace("'","")  CHECK THIS!!
    
    place_holder_6 = f'''WHERE owner_nm in ({owner_list}) '''

    print("Retailers descriptions source:", owner_list)


# static variables:  static_1 --> at least X trips; static_2 --> every Y-week block; 
# static_3 --> Z number of time blocks;  static_4 --> through lag week number

if nbr_retailers < 3:
    static_1 = 1
    static_2 = 52
    print('Static: ',static_1,'trips every',static_2,'weeks')

else:
    static_1 = 2
    static_2 = 8
    print('Static: ',static_1,'trips every',static_2,'weeks')

    
if int(nbr_retailers) == 1 and retailers == '"Kroger/Roundys/Harris Teeter"':
    category_show_parameter = False
else:
    category_show_parameter = True


# 2c) Defining type of segment will be used to perform the analysis based on parameter selection made by the analyst:

# In[8]:


if segment_type == 1:
    df_segments_2 = pd.read_sql(f''' select {segment_def} as segment, count(*) as nbr_of_upcs, brand_nbr as segm_nbr from VMR_{brand_nm}_upc_filter where brand_nbr in ({brand_nbr_str}) group by 1,3 order by 3 ''',conn) 
    print(df_segments_2)
    
if segment_type == 2:
    df_segments_2 = pd.read_sql(f''' select {segment_def} as segment, count(*) as nbr_of_upcs from VMR_{brand_nm}_upc_filter where brand_nbr in ({brand_nbr_str}) group by 1 order by 1 ''',conn) 
    df_segments_2['segm_nbr'] = list(range(1,len(df_segments_2)+1))
    print(df_segments_2)
    
if segment_type == 3:
    df_segments_2 = pd.read_sql(f''' select {segment_def} as segment, count(*) as nbr_of_upcs from VMR_{brand_nm}_upc_filter where brand_nbr in ({brand_nbr_str}) group by 1 order by 1 ''',conn) 
    df_segments_2['segm_nbr'] = list(range(1,len(df_segments_2)+1))
    print(df_segments_2)    
    
if segment_type == 6:
    df_segments_2 = pd.read_sql(f''' select {segment_def} as segment, count(*) as nbr_of_upcs from VMR_{brand_nm}_upc_filter where brand_nbr in ({brand_nbr_str}) group by 1 order by 1 ''',conn) 
    df_segments_2['segm_nbr'] = list(range(1,len(df_segments_2)+1))
    print(df_segments_2)    

if segment_type == 4:
    df_segments_2 = pd.read_sql(f''' select {segment_def} as segment, count(*) as nbr_of_upcs from VMR_{brand_nm}_upc_filter where brand_nbr in ({brand_nbr_str}) and owner_nm in ({owner_list}) group by 1 order by 1 ''',conn) 
    df_segments_2['segm_nbr'] = list(range(1,len(df_segments_2)+1))
    print(df_segments_2)    
    
if segment_type == 5:
    df_segments_2 = pd.read_sql(f''' select {segment_def} as segment, count(*) as nbr_of_upcs from VMR_{brand_nm}_upc_filter where brand_nbr in ({brand_nbr_str}) group by 1 order by 1 ''',conn) 
    df_segments_2['segm_nbr'] = list(range(1,len(df_segments_2)+1))
    print(df_segments_2)    


# In[9]:


if segment_type == 4:

    curs.execute(f'''
        drop table if exists VMR_{brand_nm}_upc_filter;       
        create temp table VMR_{brand_nm}_upc_filter as
            select u.brandnbr as brand_nbr
                  ,u.branddesc as brand_desc
                  ,p.trademark_brand_desc
                  ,p.trade_item_key
                  ,p.trade_item_cd
                  ,p.cat_nbr
                  ,p.cat_desc
                  {place_holder_3}
                  {place_holder_4}
                  {place_holder_1}
            from trade_item_v p 
                  INNER JOIN VMR_{brand_nm}_upclmc_{analyst} u ON (p.trade_item_cd = u.tradeitemcd)
                  {place_holder_2}
                  {place_holder_5}
                  {place_holder_6}

            distribute replicate;
    ''')


# 2d) Displaying the offer's summary:

# In[10]:


df_promo_summary = pd.read_sql(f''' select * from VMR_{brand_nm}_promo_summary''',conn)
df_promo_summary = df_promo_summary.rename(columns = {'fin_cmit_contract_nbr': 'Contract Nbr','promo_src_id_txt': 'BL',
'fin_cmit_contract_nm': 'Fin Cmit Contract Nm', 'promo_desc_txt':'Promo Desc Txt', 'actual_start': 'Actual Start',
'actual_end':'Actual End', 'total_reward_prints':'Total Reward Prints'})
print(df_promo_summary)


# 2e) Date calculations for the analysis:

# In[11]:


#Modified code from K.Mertens's J&J custom VMR code.


#Reward print period from ord_date_key data:

reward_print_start = df_promo_summary['Actual Start'].min()
reward_print_end = df_promo_summary['Actual End'].max()

weeks_per_period_final = round(((reward_print_end - reward_print_start + dt.timedelta(days=1))/dt.timedelta(days=7)),1)


nbr_days = reward_print_end - reward_print_start
nbr_days = nbr_days.days

#number of weeks the VMR print period lasts


#VMR pre-period dates: (using the same program lenght as the reward print period)

vmr_pre_period_start = reward_print_start - dt.timedelta(days = nbr_days+1)
vmr_pre_period_end = reward_print_start - dt.timedelta(days=1)


#52 Prior-Period: (using 52 parameter)

prior_period_start = reward_print_start - dt.timedelta(weeks = pre_weeks) #pre_weeks before the reward print period starts
prior_period_end = reward_print_start - dt.timedelta(days=1) #one day before the reward print period starts


#VMR Period YAGO (same length as VMR Period):

yago_print_start = reward_print_start - dt.timedelta(weeks = 52) #52 weeks before the reward print period starts
yago_print_end = reward_print_end - dt.timedelta(weeks = 52) #52 weeks before the reward print period starts 


#Reward Period + Redemption days:

redemption_start = reward_print_start  #the BL could be redeemed the same day it was printed
redemption_end = reward_print_end + dt.timedelta(days = redemption_days)  #the last possible redemption could be the last printing day + redemption_days


#Reward Period + 4 WEEK

post_4wk_start = reward_print_end + dt.timedelta(days = 1) #the BL could be redeemed the same day it was printed
post_4wk_end = reward_print_end + dt.timedelta(days = nbr_days+1)  


#Pre weeks for PPTX Trend

vmr_pre26wk_start =  reward_print_start - dt.timedelta(weeks = 26)
vmr_pre26wk_end = reward_print_start - dt.timedelta(days = 1)


yago_26wk_start = vmr_pre26wk_start - dt.timedelta(weeks = 52)

#Lasta date of data:
last_data_dt =  (dt.date.today()- dt.timedelta(days = 2))


#Static calculations

curr_date = dt.date.today()
sun_offset = 1 + curr_date.weekday()
rec_sun = curr_date - dt.timedelta(days=sun_offset)

look_back = round(((curr_date - prior_period_start)/dt.timedelta(days=1) + 1)/7,1)

static_3 = int(round(((reward_print_end - prior_period_start)/dt.timedelta(days=1) + 1)/(7*static_2),2))
static_4 = int(round((((rec_sun - reward_print_end)/dt.timedelta(days=1) + 1)/7),0)) 


# In[12]:


print(curr_date)
print(yago_26wk_start)


# 2f) Dimensional table for dates:

# In[13]:


#Dates dimensional table

curs.execute(f'''       
    DROP table if exists VMR_{brand_nm}_date_filter;
    CREATE temp table VMR_{brand_nm}_date_filter as
    
    
        SELECT distinct date_key,cal_dt,cal_sun_wk_ending_dt,cal_sun_wk_ending_rank_nbr
              ,CASE 
                    when cal_dt between '{reward_print_start}' and '{reward_print_end}' then 'VMR Print Period ({round(weeks_per_period_final)} weeks)'
                    when cal_dt between '{vmr_pre_period_start}' and '{vmr_pre_period_end}' then 'VMR Pre-period ({round(weeks_per_period_final)} weeks)'
                    when cal_dt between '{yago_print_start}' and '{yago_print_end}' then 'Period Year Ago ({round(weeks_per_period_final)} weeks)'
                    end as analysis_period
              ,CASE when cal_dt between '{prior_period_start}' and '{prior_period_end}' then 'Prior-period ({round(pre_weeks)} weeks)'
                    end as prior_period
              ,CASE when cal_dt between '{redemption_start}' and '{redemption_end}' then 'Redemption Period ({redemption_days} days)'
                    end as redemption_period
              ,CASE when cal_dt between '{post_4wk_start}' and '{post_4wk_end}' then 'Post Period ({round(weeks_per_period_final)} weeks)'
                    end as post_period
                    
        FROM date_v
        
        WHERE cal_dt between '{yago_26wk_start}' and '{curr_date}'
        
        
    DISTRIBUTE REPLICATE;

''')


# In[14]:


#Getting summarize table with all the dates

df_date_check = pd.read_sql(f'''

    WITH union_table AS 
    
    (SELECT CASE WHEN analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN 'VMR Period - TY ({round(weeks_per_period_final)} weeks)'
                WHEN analysis_period = 'VMR Pre-period ({round(weeks_per_period_final)} weeks)' THEN 'VMR Period - Pre Period ({round(weeks_per_period_final)} weeks)'
                WHEN analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)' THEN 'VMR Period - YAGO ({round(weeks_per_period_final)} weeks)'
                ELSE null END as analysis_periods
    , min(cal_dt) as start, max(cal_dt) as end 
    FROM VMR_{brand_nm}_date_filter
    WHERE analysis_period is not null
    GROUP BY 1
    UNION
    SELECT CASE WHEN redemption_period = 'Redemption Period ({redemption_days} days)' THEN 'Reward Period + ({redemption_days} days)'
                ELSE null END as analysis_periods
    , min(cal_dt) as start, max(cal_dt) as end 
    FROM VMR_{brand_nm}_date_filter
    WHERE redemption_period = 'Redemption Period ({redemption_days} days)'
    GROUP BY 1
    UNION
    SELECT CASE WHEN prior_period = 'Prior-period ({round(pre_weeks)} weeks)' THEN '{pre_weeks} wk Prior Period'
           ELSE null END analysis_periods
    , min(cal_dt) as start, max(cal_dt) as end 
    FROM VMR_{brand_nm}_date_filter
    WHERE prior_period = 'Prior-period ({round(pre_weeks)} weeks)'
    GROUP BY 1
    UNION
    SELECT CASE WHEN post_period = 'Post Period ({round(weeks_per_period_final)} weeks)' THEN 'Post Period ({round(weeks_per_period_final)} weeks)'
           ELSE null END analysis_periods
    , min(cal_dt) as start, max(cal_dt) as end 
    FROM VMR_{brand_nm}_date_filter
    WHERE post_period = 'Post Period ({round(weeks_per_period_final)} weeks)'
    GROUP BY 1
    
    )
    
    SELECT *
    FROM union_table
    ORDER BY CASE WHEN analysis_periods = '{pre_weeks} wk Prior Period' THEN 1
               WHEN analysis_periods = 'VMR Period - YAGO ({round(weeks_per_period_final)} weeks)' THEN 2
               WHEN analysis_periods = 'VMR Period - Pre Period ({round(weeks_per_period_final)} weeks)' THEN 3
               WHEN analysis_periods = 'VMR Period - TY ({round(weeks_per_period_final)} weeks)' THEN 4
               WHEN analysis_periods = 'Redemption Period ({redemption_days} days)' THEN 5
               WHEN analysis_periods = '4wk Post Period' THEN 6
               END
    ''',conn)

df_date_check = df_date_check.rename(columns={'analysis_periods':'Analysis Periods','start':'Start', 'end':'End'})
df_date_check


# 2g) Shopper Consistent calculation parameters:

# In[15]:


import datetime as dt 
import numpy as np

Current_beg = str(yago_print_start.strftime('%Y-%m-%d'))   # The beginning date of current period; Must be 'YYYY-MM-DD' format!
Current_end = str(post_4wk_end.strftime('%Y-%m-%d')) # The ending date of current period; Must be 'YYYY-MM-DD' format!

# Convert Dates from String type to Date type
Current_beg = dt.datetime.fromisoformat(Current_beg).date()
Current_end = dt.datetime.fromisoformat(Current_end).date()

# Prepare Static Parameters
curr_date = dt.date.today()
sun_offset = 1 + curr_date.weekday()
rec_sun = curr_date - dt.timedelta(days=sun_offset)

static_3 = int(np.floor(((Current_end - Current_beg)/dt.timedelta(days=1) + 1)/(7*static_2)))
static_4 = max(int(np.floor((((rec_sun - Current_end)/dt.timedelta(days=1) + 1)/7))), 0)

# Convert Date type to Numeric type
Current_beg = str(Current_beg)[:10].replace('-','')
Current_end = str(Current_end)[:10].replace('-','')

print('Numeric Begin Date of Current Period:',Current_beg)
print('Numeric End Date of Current Period:',Current_end)

print('\n','Static Information: Based on Current_beg through Current_end', sep='')
print('Last Sunday:',rec_sun.strftime("%A %B %d %Y"))
print('Static: ',static_1,'trips every',static_2,'weeks')
print('Static: ',static_3,'time blocks ending',static_4,'weeks back','\n')


#Dimensional table for shopper-consistency

curs.execute(f''' 
DROP table if exists shopper_consistent_{analyst};
        CREATE temp table shopper_consistent_{analyst} as
                select distinct cnsmr_id_key
                from consumer_id_v
                where abuser_ind= 'N' 
                     and unident_ind = 'N'
                     and really_bad_abuser_ind='N'
                     and cnsmr_id_key>0
                     and isconsistshop(shop_ord_hst_bmap,{static_1},{static_2},{static_3},{static_4})
                DISTRIBUTE ON (cnsmr_id_key);
''')


# ## **3) GETTING ALL THE REWARD TRANSACTIONS (ORD_EVENT_KEY) AND REWARD ID'S FOR VMR PERIOD:**

# In[16]:


#ALL REWARD TRANSACTIONS FOR VMR PERIOD:

curs.execute(f'''    
    
--Getting all the reward transactions:

    DROP table if exists VMR_{brand_nm}_reward_trans;
    CREATE table VMR_{brand_nm}_reward_trans as
        
          SELECT o.ord_event_key,
                 cal_dt as reward_date,
                 sum(o.purch_qty) as units,
                 sum(o.purch_amt) as dollars
          
          FROM ord_trd_itm_cnsmr_fact_ne_v o
                 
                 INNER JOIN VMR_{brand_nm}_date_filter b ON (o.ord_date_key = b.date_key)
                 INNER JOIN VMR_{brand_nm}_upc_filter c ON (o.trade_item_key = c.trade_item_key)
                 INNER JOIN VMR_{brand_nm}_printing_stores s ON (o.ord_touchpoint_key = s.touchpoint_key)
          
          WHERE b.analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)'
                and o.purch_amt > 0
                and o.purch_qty > 0
                and brand_nbr in ({brand_nbr_str})
                and o.ord_date_key between s.min_actual_date_key and s.max_actual_date_key
        
        GROUP BY 1,2
        
        HAVING {min_threshold_statement}
        
        DISTRIBUTE RANDOM;
        
''')

curs.execute(f'''    
GRANT ALL ON TABLE VMR_{brand_nm}_reward_trans TO public
''')



print(f'''VMR_{brand_nm}_reward_trans''')

curs.execute(f'''    

--Getting the list with the all distinct reward ord_date_key:

    DROP table if exists VMR_{brand_nm}_reward_events;
    CREATE temp table VMR_{brand_nm}_reward_events as
    
        SELECT distinct ord_event_key

        FROM VMR_{brand_nm}_reward_trans

        DISTRIBUTE RANDOM;
''')

#ALL REWARD ID's FOR VMR PERIOD:

curs.execute(f'''   

--Getting all the Reward transactions and IDs:

    DROP table if exists VMR_{brand_nm}_reward_ids;
    CREATE temp table VMR_{brand_nm}_reward_ids as
        
          SELECT o.ord_event_key,
                 o.ord_designated_cnsmr_id_key, 
                 b.cal_dt as reward_date,
                 sum(o.purch_qty) as units,
                 sum(o.purch_amt) as dollars
          
          FROM ord_trd_itm_cnsmr_fact_ne_v o
                 
                 INNER JOIN VMR_{brand_nm}_date_filter b ON (o.ord_date_key = b.date_key)
                 INNER JOIN VMR_{brand_nm}_upc_filter c ON (o.trade_item_key = c.trade_item_key)
                 INNER JOIN VMR_{brand_nm}_printing_stores s ON (o.ord_touchpoint_key = s.touchpoint_key)
                 INNER JOIN consumer_id_v p ON (o.ord_designated_cnsmr_id_key = p.cnsmr_id_key)

          
          WHERE p.abuser_ind= 'N' 
                and p.unident_ind = 'N'
                and p.really_bad_abuser_ind='N'
                and p.cnsmr_id_key>0
                and b.analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)'
                and o.purch_amt > 0
                and o.purch_qty > 0
                and brand_nbr in ({brand_nbr_str})
                and o.ord_date_key between s.min_actual_date_key and s.max_actual_date_key
        
        GROUP BY 1,2,3
        
        having {min_threshold_statement}
        
        
        DISTRIBUTE ON (ord_event_key);
        
        
        
--Getting the list with the all distinct Reward ID's:

    DROP table if exists VMR_{brand_nm}_reward_ids_csmr;
    CREATE temp table VMR_{brand_nm}_reward_ids_csmr as
    
    SELECT distinct ord_designated_cnsmr_id_key
    
    FROM VMR_{brand_nm}_reward_ids
    
    DISTRIBUTE RANDOM;
        
''')


# In[17]:


total_reward_ids = pd.read_sql(f''' SELECT COUNT(ord_designated_cnsmr_id_key) FROM VMR_{brand_nm}_reward_ids_csmr''',conn)
print("The total number of trackable reward ID's is: " + str(total_reward_ids['count'].values[0]))
total_reward_ids = total_reward_ids['count'].values[0]


# In[18]:


consistent_reward_ids = pd.read_sql(f''' SELECT COUNT(ord_designated_cnsmr_id_key) 
FROM VMR_{brand_nm}_reward_ids_csmr o
INNER JOIN shopper_consistent_{analyst} f ON (o.ord_designated_cnsmr_id_key = f.cnsmr_id_key)
''',conn)
print("The total number of consistent reward ID's is: " + str(consistent_reward_ids['count'].values[0]))


# In[19]:


total_reward_trans = pd.read_sql(f''' SELECT COUNT(ord_event_key) FROM VMR_{brand_nm}_reward_events''',conn)
print("The total number of reward transactions is: " + str(total_reward_trans['count'].values[0]))

total_reward_trans_var = total_reward_trans

total_reward_trans = str(total_reward_trans['count'].values[0])


# ## **4) GETTING DOLLARS AND UNITS FOR YAGO, PRE-PERIOD AND VMR PERIOD FOR ALL TRANSACTIONS OCCURED (FOR ANY SHOPPER):**

# 4a) % dollars_per_trip change between VMR period and Yago trips (for any shopper and all transactions):

# In[20]:


curs.execute(f'''       
    DROP table if exists VMR_{brand_nm}_YAGO_trend;
    CREATE temp table VMR_{brand_nm}_YAGO_trend as
    
        SELECT  
               CASE WHEN analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN 'VMR Period'
                    ELSE 'NA' END as analysis_periods,
               count(distinct a.ord_event_key) as count_distinct_trips,
               sum(a.purch_amt) as dollar_sales,
               dollar_sales/count_distinct_trips as dollars_per_trip,
               sum(a.purch_qty) as units_sales
              
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
            
        WHERE analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)'
              and brand_nbr IN ({brand_nbr_str})
              and a.purch_amt > 0
              and a.purch_qty > 0
        
        GROUP BY 1
    
        UNION
        
        
        SELECT  
               CASE WHEN analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)' THEN 'YAGO Period'
                    ELSE 'NA' END as analysis_periods,
               count(distinct a.ord_event_key) as count_distinct_trips,
               sum(a.purch_amt) as dollar_sales,
               dollar_sales/count_distinct_trips as dollars_per_trip,
               sum(a.purch_qty) as units_sales
              
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
            
        WHERE analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)'  
              and brand_nbr IN ({brand_nbr_str})
              and a.purch_amt > 0
              and a.purch_qty > 0
        
        GROUP BY 1
        
        ORDER BY 1 DESC
        
    DISTRIBUTE REPLICATE;

''')

# In[21]:


vmr_trend_yago_summary = pd.read_sql(f'''
    SELECT analysis_periods,
           sum(count_distinct_trips) as count_distinct_trips,
           sum(dollar_sales) as dollar_sales,
           sum(dollar_sales)/sum(count_distinct_trips) as dollars_per_trip,
           sum(units_sales) as units_sales,
           round(sum(units_sales)::float/sum(count_distinct_trips)::float,1) as units_per_trip
    FROM VMR_{brand_nm}_YAGO_trend
    WHERE analysis_periods != 'NA'
    GROUP BY 1
    ORDER BY CASE WHEN analysis_periods = 'YAGO Period' THEN 1
                  WHEN analysis_periods = 'VMR Period' THEN 2
                  END
    ''',conn)


#Calculating %change to insert them in the summary table:

if vmr_trend_yago_summary.loc[0, 'analysis_periods'] != 'YAGO Period':
    
    vmr_trend_yago_summary = vmr_trend_yago_summary.rename(columns = {

    'analysis_periods':'Analysis Periods', 'count_distinct_trips':'Count Distinct Trips', 'dollar_sales':'Dollar Sales', 'dollars_per_trip':'Dollars per Trip', '%_change_dollars_per_trip':'% Change Dollars per Trip',
    'units_sales':'Units Sold', 'units_per_trip':'Units per Trip', '%_change_units_per_trip':'% Change Units per Trip'

    })
    
    dollars_moved = round(vmr_trend_yago_summary.loc[0,'Dollar Sales'])
    dollars_moved = "{:,}".format(dollars_moved)
    
    units_moved = round(vmr_trend_yago_summary.loc[0,'Units Sold'])
    units_moved = "{:,}".format(units_moved)

    dollars_per_trip_chg_yago = 0
    
    units_per_trip_chg_yago = 0
    
    dollars_chg_yago = 0
    
    units_chg_yago = 0
    
    print("NOTE: There is no YAGO data for this UPC list.")
    print("")
    print(vmr_trend_yago_summary)

    
else:
    
    change_dollars_per_trip = (vmr_trend_yago_summary.loc[1,'dollars_per_trip'] - vmr_trend_yago_summary.loc[0,'dollars_per_trip'])/vmr_trend_yago_summary.loc[0,'dollars_per_trip']
    vmr_trend_yago_summary['%_change_dollars_per_trip'] = ['-', change_dollars_per_trip]
    
    change_units_per_trip = (vmr_trend_yago_summary.loc[1,'units_per_trip'] - vmr_trend_yago_summary.loc[0,'units_per_trip'])/vmr_trend_yago_summary.loc[0,'units_per_trip']
    vmr_trend_yago_summary['%_change_units_per_trip'] = ['-', change_units_per_trip]

    vmr_trend_yago_summary = vmr_trend_yago_summary.rename(columns = {

    'analysis_periods':'Analysis Periods', 'count_distinct_trips':'Count Distinct Trips', 'dollar_sales':'Dollar Sales', 'dollars_per_trip':'Dollars per Trip', '%_change_dollars_per_trip':'% Change Dollars per Trip',
    'units_sales':'Units Sold', 'units_per_trip':'Units per Trip', '%_change_units_per_trip':'% Change Units per Trip'

    })

    #Extracting variables for slides

    dollars_moved = round(vmr_trend_yago_summary.loc[1,'Dollar Sales'])
    dollars_moved = "{:,}".format(dollars_moved)
    
    units_moved = round(vmr_trend_yago_summary.loc[0,'Units Sold'])
    units_moved = "{:,}".format(units_moved)
    
    dollars_chg_yago = round(((vmr_trend_yago_summary.loc[1,'Dollar Sales'] - vmr_trend_yago_summary.loc[0,'Dollar Sales'])/vmr_trend_yago_summary.loc[0,'Dollar Sales'])*100)
    units_chg_yago = round(((vmr_trend_yago_summary.loc[1,'Units Sold'] - vmr_trend_yago_summary.loc[0,'Units Sold'])/vmr_trend_yago_summary.loc[0,'Units Sold'])*100)

    dollars_per_trip_chg_yago = round(vmr_trend_yago_summary.loc[1,'% Change Dollars per Trip']*100)
    
    units_per_trip_chg_yago = round(vmr_trend_yago_summary.loc[1,'% Change Units per Trip']*100)

    print(vmr_trend_yago_summary)
    

# 4a) % dollars_per_trip change between VMR period and Pre-Period trips (for any shopper and all transactions):

# In[22]:


curs.execute(f'''       
    DROP table if exists VMR_{brand_nm}_vmr_pre_period_trend;
    CREATE temp table VMR_{brand_nm}_vmr_pre_period_trend as
    
        SELECT  
               CASE WHEN analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN 'VMR Period'
                    ELSE 'NA' END as analysis_periods,
               count(distinct a.ord_event_key) as count_distinct_trips,
               sum(a.purch_amt) as dollar_sales,
               dollar_sales/count_distinct_trips as dollars_per_trip,
               sum(a.purch_qty) as units_sales
              
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
            
        WHERE analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)'
              and brand_nbr IN ({brand_nbr_str})
              and a.purch_amt > 0
              and a.purch_qty > 0
        
        GROUP BY 1
        
        UNION
        
        SELECT  
               CASE WHEN analysis_period = 'VMR Pre-period ({round(weeks_per_period_final)} weeks)' THEN 'VMR Pre-Period'
                    ELSE 'NA' END as analysis_periods,
               count(distinct a.ord_event_key) as count_distinct_trips,
               sum(a.purch_amt) as dollar_sales,
               dollar_sales/count_distinct_trips as dollars_per_trip,
               sum(a.purch_qty) as units_sales
              
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
            LEFT JOIN VMR_{brand_nm}_reward_events e ON (e.ord_event_key = a.ord_event_key)
            
        WHERE analysis_period = 'VMR Pre-period ({round(weeks_per_period_final)} weeks)'
              and brand_nbr IN ({brand_nbr_str})
              and a.purch_amt > 0
              and a.purch_qty > 0
        
        GROUP BY 1
        
        ORDER BY 1 DESC
        
    DISTRIBUTE REPLICATE;

''')

# In[23]:


vmr_trend_pre_period_summary = pd.read_sql(f'''
    SELECT analysis_periods,
           sum(count_distinct_trips) as count_distinct_trips,
           sum(dollar_sales) as dollar_sales,
           sum(dollar_sales)/sum(count_distinct_trips) as dollars_per_trip,
           sum(units_sales) as units_sales,
           round(sum(units_sales)::float/sum(count_distinct_trips)::float,1) as units_per_trip
    FROM VMR_{brand_nm}_vmr_pre_period_trend
    WHERE analysis_periods != 'NA'
    GROUP BY 1
    ORDER BY CASE WHEN analysis_periods = 'VMR Pre-Period' THEN 1
                  WHEN analysis_periods = 'VMR Period' THEN 2
                  END
    ''',conn)


#Calculating %change to insert them in the summary table:

change_dollars_per_trip = ((vmr_trend_pre_period_summary.loc[1,'dollars_per_trip'] - vmr_trend_pre_period_summary.loc[0,'dollars_per_trip'])/vmr_trend_pre_period_summary.loc[0,'dollars_per_trip'])
vmr_trend_pre_period_summary['%_change_dollars_per_trip'] = ['-', change_dollars_per_trip]

change_units_per_trip = (vmr_trend_pre_period_summary.loc[1,'units_per_trip'] - vmr_trend_pre_period_summary.loc[0,'units_per_trip'])/vmr_trend_pre_period_summary.loc[0,'units_per_trip']
vmr_trend_pre_period_summary['%_change_units_per_trip'] = ['-', change_units_per_trip]


vmr_trend_pre_period_summary = vmr_trend_pre_period_summary.rename(columns = {
    
'analysis_periods':'Analysis Periods', 'count_distinct_trips':'Count Distinct Trips', 'dollar_sales':'Dollar Sales', 'dollars_per_trip':'Dollars per Trip', '%_change_dollars_per_trip':'% Change Dollars per Trip',
'units_sales':'Units Sold', 'units_per_trip':'Units per Trip', '%_change_units_per_trip':'% Change Units per Trip'
    
})

dollars_chg_pre_period = round((vmr_trend_pre_period_summary.loc[1,'Dollar Sales'] - vmr_trend_pre_period_summary.loc[0,'Dollar Sales'])/vmr_trend_pre_period_summary.loc[0,'Dollar Sales'])

dollars_per_trip_chg_pre_period = round(vmr_trend_pre_period_summary.loc[1,'% Change Dollars per Trip']*100)
    
units_per_trip_chg_pre_period = round(vmr_trend_pre_period_summary.loc[1,'% Change Units per Trip']*100)


print(vmr_trend_pre_period_summary)


# ## **5) GETTING TOTAL DOLLARS FROM REWARD TRIPS BY SEGMENT DURING VMR PERIOD (FOR ANY SHOPPER):**

# 5a) Getting the total dollars by segment for Reward trips during VMR period (for any shopper):

# In[24]:


curs.execute(f'''    
    
 --Getting Units and Dollars for all the reward transactions by segment:
        
    DROP table if exists VMR_{brand_nm}_reward_segments;
    CREATE temp table VMR_{brand_nm}_reward_segments as
        
          SELECT {segment_def},
                 sum(o.purch_qty) as units,
                 sum(o.purch_amt) as dollars
          
          FROM ord_trd_itm_cnsmr_fact_ne_v o
                 
                 INNER JOIN VMR_{brand_nm}_date_filter b ON (o.ord_date_key = b.date_key)
                 INNER JOIN VMR_{brand_nm}_upc_filter c ON (o.trade_item_key = c.trade_item_key)
                 INNER JOIN VMR_{brand_nm}_reward_events r ON (o.ord_event_key = r.ord_event_key)
          
          WHERE b.analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)'
                and o.purch_amt > 0
                and o.purch_qty > 0
                and brand_nbr in ({brand_nbr_str})
                
        
        GROUP BY 1
        
                
        DISTRIBUTE RANDOM;
    

''')


# In[25]:


total_vmr_details_by_brand = pd.read_sql(f'''

    SELECT *
    FROM VMR_{brand_nm}_reward_segments
    ORDER BY dollars DESC

    ''',conn)

print(total_vmr_details_by_brand)

total_vmr_details_by_brand["units"] = total_vmr_details_by_brand["units"].fillna(0)
total_vmr_details_by_brand["dollars"] = total_vmr_details_by_brand["dollars"].fillna(0)

total_vmr_details_by_brand["% Units"] = (total_vmr_details_by_brand["units"])/(total_vmr_details_by_brand["units"].sum())

total_vmr_details_by_brand["% Dollars"] = (total_vmr_details_by_brand["dollars"])/(total_vmr_details_by_brand["dollars"].sum())

print(total_vmr_details_by_brand)

total_vmr_details_segm_chart = total_vmr_details_by_brand[segment_def]
total_vmr_details_units_chart = total_vmr_details_by_brand["% Units"]
total_vmr_details_dollars_chart = total_vmr_details_by_brand["% Dollars"]

print(total_vmr_details_dollars_chart)

total_vmr_details_by_brand[segment_def] = total_vmr_details_by_brand[segment_def].fillna('Null')
total_vmr_details_by_brand["% Units"] = total_vmr_details_by_brand["% Units"].fillna(0)
total_vmr_details_by_brand["% Dollars"] = total_vmr_details_by_brand["% Dollars"].fillna(0)


total_vmr_details_by_brand.loc['Total']= total_vmr_details_by_brand.sum()
total_vmr_details_by_brand.reset_index(drop=True)
total_vmr_details_by_brand.at['Total',segment_def]='Total'


        
print(total_vmr_details_by_brand)

dollars_moved_vmr = total_vmr_details_by_brand.loc['Total', 'dollars']
units_moved_vmr = total_vmr_details_by_brand.loc['Total', 'units']

print(units_moved_vmr)


total_vmr_details_by_brand = total_vmr_details_by_brand.rename(columns= {segment_def: "Segment", 'units':'Units', 'dollars':'Dollars'})
print(total_vmr_details_by_brand)


# In[26]:


curs.execute(f'''

        DROP table if exists VMR_{brand_nm}_bls_levels;
        CREATE temp table VMR_{brand_nm}_bls_levels as

            select a.PROMO_SRC_ID_TXT as mclu_blip, a.promo_src_id as mclu_nbr,
                   a.promo_dist_start_dt as start_dt, a.promo_dist_stop_dt as stop_dt,
                   CASE when a.promo_discnt_val_amt>0 then a.promo_discnt_val_amt else NULL end as cpn_val_amt,
                   CASE WHEN a.cnsmr_curr_beh_purch_rqmt_qty>0 THEN a.cnsmr_curr_beh_purch_rqmt_qty 
                       ELSE a.cnsmr_TRG_CLASS_purch_rqmt_qty END as min_qty,
                   avg(b.promo_red_handling_fee_rt) as handling_fee,
                   avg(c.avg_imprsn_prc_rt) as impression_cost
    from promotion_v a
      inner join promotion_cost_v b on (a.promo_key=b.promo_key)
      inner join (select promo_key, avg_imprsn_prc_rt 
                    from promo_chanl_prfrm_cum_aggr_v 
                    where dist_chanl_cd='IN-STORE' 
                                    and imprsn_typ_cd='IN-STORE PRINT') c on (a.promo_key=c.promo_key)
    where a.PROMO_SRC_ID_TXT in ({BL_CODES})
    group by 1,2,3,4,5,6
    order by min_qty
    
    distribute random;

''')
 

bl_data = pd.read_sql(f''' select * from VMR_{brand_nm}_bls_levels''',conn)
print(bl_data)


bl_data = pd.read_sql(f''' select * from VMR_{brand_nm}_bls_levels''',conn)
print(bl_data)


min_thres = list(bl_data['min_qty'].dropna().unique())
min_thres = sorted(min_thres)
print(min_thres)

# If min_thres is empty or only contains 0, use user_min_threshold if available
if (len(min_thres) == 0 or min_thres == [0] or min_thres == [0.0]) and user_min_threshold is not None:
    min_thres = [user_min_threshold]
    print(f'Using user-provided min_thres for reward levels: {min_thres}')


if threshold_unit == 'units':
    
    for i in range(len(min_thres)):
    
        if i == 0:
            if min_thres[i] == max(min_thres):
                plchldr1 = f'''CASE WHEN sum(o.purch_qty) >= {str(min_thres[i])} then 'Buy {str(min_thres[i])}+' ''' + "\n"
            else:
                plchldr1 = f'''CASE WHEN sum(o.purch_qty) >= {str(min_thres[i])} and sum(o.purch_qty) < {str(min_thres[i+1])} then 'Buy {str(min_thres[i])}-{str(min_thres[i+1])}' ''' + "\n"
                        
                
        else:
            if min_thres[i] == max(min_thres):
                plchldr1 =  plchldr1 + f'''WHEN sum(o.purch_qty) >= {str(min_thres[i])} then 'Buy {str(min_thres[i])}+' ''' + "\n"
            else:
                plchldr1 = plchldr1 + f'''WHEN sum(o.purch_qty) >= {str(min_thres[i])} and sum(o.purch_qty) < {str(min_thres[i+1])} then 'Buy {str(min_thres[i])}-{str(min_thres[i+1])}' ''' + "\n"
            

if threshold_unit == 'dollars':
    
    for i in range(len(min_thres)):
    
        if i == 0:
            if min_thres[i] == max(min_thres):
                plchldr1 = f'''CASE WHEN sum(o.purch_amt) >= {str(min_thres[i])} then 'Buy ${str(min_thres[i])}+' ''' + "\n"
            else:
                plchldr1 = f'''CASE WHEN sum(o.purch_amt) >= {str(min_thres[i])} and sum(o.purch_amt) < {str(min_thres[i+1])} then 'Buy ${str(min_thres[i])}-${str(min_thres[i+1])}' ''' + "\n"
        else:
            if min_thres[i] == max(min_thres):
                plchldr1 =  plchldr1 + f'''WHEN sum(o.purch_amt) >= {str(min_thres[i])} then 'Buy ${str(min_thres[i])}+' ''' + "\n"
            else:
                plchldr1 = plchldr1 + f'''WHEN sum(o.purch_amt) >= {str(min_thres[i])} and sum(o.purch_amt) < {str(min_thres[i+1])} then 'Buy ${str(min_thres[i])}-${str(min_thres[i+1])}' ''' + "\n"
         
         

plchldr1 = plchldr1 + f'''ELSE 'All Other Transactions' END AS level'''

print(plchldr1)


curs.execute(f'''    
            
    DROP table if exists VMR_{brand_nm}_reward_levels;
    CREATE temp table VMR_{brand_nm}_reward_levels as
    
        
         SELECT o.ord_event_key,'reward period' as period,
          {plchldr1}, 
          sum(o.purch_qty) as units, sum(o.purch_amt) as dollars
          

          FROM ord_trd_itm_cnsmr_fact_ne_v o
                 
                 INNER JOIN VMR_{brand_nm}_date_filter b ON (o.ord_date_key = b.date_key)
                 INNER JOIN VMR_{brand_nm}_upc_filter c ON (o.trade_item_key = c.trade_item_key)
                 INNER JOIN VMR_{brand_nm}_printing_stores a ON (o.ord_touchpoint_key = a.touchpoint_key)
                 INNER JOIN VMR_{brand_nm}_reward_events r ON (o.ord_event_key = r.ord_event_key)
                 
          
          WHERE b.analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' and
                o.purch_amt > 0
                and o.purch_qty > 0
                and brand_nbr in ({brand_nbr_str})
                and cal_dt between '{reward_print_start}' and '{reward_print_end}'

            
          GROUP BY 1,2
          
           HAVING level != 'All Other Transactions'
          
          
          UNION
          
          
          SELECT o.ord_event_key,'reward period' as period,
          {plchldr1}, 
          sum(o.purch_qty) as units, sum(o.purch_amt) as dollars
          

          FROM ord_trd_itm_cnsmr_fact_ne_v o
                 
                 INNER JOIN VMR_{brand_nm}_date_filter b ON (o.ord_date_key = b.date_key)
                 INNER JOIN VMR_{brand_nm}_upc_filter c ON (o.trade_item_key = c.trade_item_key)
                 INNER JOIN VMR_{brand_nm}_printing_stores a ON (o.ord_touchpoint_key = a.touchpoint_key)                 
          
          WHERE b.analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' and
                o.purch_amt > 0
                and o.purch_qty > 0
                and brand_nbr in ({brand_nbr_str})
                and cal_dt between '{reward_print_start}' and '{reward_print_end}'

            
          GROUP BY 1,2
          
          HAVING level = 'All Other Transactions'
          
          
          UNION

          
          SELECT ord_event_key,'pre-period' as period,
          {plchldr1}, 
          sum(o.purch_qty) as units, sum(o.purch_amt) as dollars
          

          FROM ord_trd_itm_cnsmr_fact_ne_v o
                 
                 INNER JOIN VMR_{brand_nm}_date_filter b ON (o.ord_date_key = b.date_key)
                 INNER JOIN VMR_{brand_nm}_upc_filter c ON (o.trade_item_key = c.trade_item_key)
                 INNER JOIN VMR_{brand_nm}_printing_stores a ON (o.ord_touchpoint_key = a.touchpoint_key)
                 
          
          WHERE b.analysis_period = 'VMR Pre-period ({round(weeks_per_period_final)} weeks)' and
                o.purch_amt > 0
                and o.purch_qty > 0
                and brand_nbr in ({brand_nbr_str}) and
                cal_dt between '{vmr_pre_period_start}' and '{vmr_pre_period_end}' 

            
          GROUP BY 1,2
          
          
          
          UNION
          
          
          
          SELECT ord_event_key,'YAGO period' as period,
          {plchldr1}, 
          sum(o.purch_qty) as units, sum(o.purch_amt) as dollars
          

          FROM ord_trd_itm_cnsmr_fact_ne_v o
                 
                 INNER JOIN VMR_{brand_nm}_date_filter b ON (o.ord_date_key = b.date_key)
                 INNER JOIN VMR_{brand_nm}_upc_filter c ON (o.trade_item_key = c.trade_item_key)
                 INNER JOIN VMR_{brand_nm}_printing_stores a ON (o.ord_touchpoint_key = a.touchpoint_key)
                 
          
          WHERE b.analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)' and
                o.purch_amt > 0
                and o.purch_qty > 0
                and brand_nbr in ({brand_nbr_str}) and
                cal_dt between '{yago_print_start}' and '{yago_print_end}'

            
          GROUP BY 1,2
          
          
          
          UNION
          
          
          SELECT ord_event_key,'prior 52wk' as period,
          {plchldr1}, 
          sum(o.purch_qty) as units, sum(o.purch_amt) as dollars
          

          FROM ord_trd_itm_cnsmr_fact_ne_v o
                 
                 INNER JOIN VMR_{brand_nm}_date_filter b ON (o.ord_date_key = b.date_key)
                 INNER JOIN VMR_{brand_nm}_upc_filter c ON (o.trade_item_key = c.trade_item_key)
                 INNER JOIN VMR_{brand_nm}_printing_stores a ON (o.ord_touchpoint_key = a.touchpoint_key)
                 
          
          WHERE b.prior_period = 'Prior-period ({round(pre_weeks)} weeks)' and
                o.purch_amt > 0
                and o.purch_qty > 0
                and brand_nbr in ({brand_nbr_str}) and
                cal_dt between '{prior_period_start}' and '{prior_period_end}'


            
          GROUP BY 1,2
          
    
          
        
            DISTRIBUTE RANDOM;
    
''')


curs.execute(f'''    
            
    DROP table if exists VMR_{brand_nm}_reward_levels_end;
    CREATE temp table VMR_{brand_nm}_reward_levels_end as
    
         SELECT level, period,
                 sum(dollars) as dollars
                 , count(distinct ord_event_key) as trips, sum(units) as units
                 
          
          FROM VMR_{brand_nm}_reward_levels
          
          GROUP BY 1,2
          
          DISTRIBUTE RANDOM;
    
''')

# In[27]:


level_ct = pd.read_sql(f'''

    SELECT level, dollars, trips, units
    FROM VMR_{brand_nm}_reward_levels_end
    WHERE period = 'reward period'
    
    UNION
    
    SELECT 'Grand Total' as level, sum(dollars) as dollars, sum(trips) as trips, sum(units) as units
    FROM VMR_{brand_nm}_reward_levels_end
    WHERE period = 'reward period'
    
    
    ORDER BY 3 ASC


    ''',conn)

level_ct



level_ct_yago = pd.read_sql(f'''

    SELECT level, dollars, trips, units
    FROM VMR_{brand_nm}_reward_levels_end
    WHERE period = 'YAGO period'
    
    UNION
    
    SELECT 'Grand Total' as level, sum(dollars) as dollars, sum(trips) as trips, sum(units) as units
    FROM VMR_{brand_nm}_reward_levels_end
    WHERE period = 'YAGO period'
    
    
    ORDER BY 3 ASC

    ''',conn)

level_ct_yago


level_ct_prior = pd.read_sql(f'''

    SELECT level, dollars, trips, units
    FROM VMR_{brand_nm}_reward_levels_end
    WHERE period = 'prior 52wk'

    UNION
    
    SELECT 'Grand Total' as level, sum(dollars) as dollars, sum(trips) as trips, sum(units) as units
    FROM VMR_{brand_nm}_reward_levels_end
    WHERE period = 'prior 52wk'
    
    
    ORDER BY 3 ASC

    ''',conn)

level_ct_prior


level_ct_prep = pd.read_sql(f'''

    SELECT level, dollars, trips, units
    FROM VMR_{brand_nm}_reward_levels_end
    WHERE period = 'pre-period'
   
   UNION
    
    SELECT 'Grand Total' as level, sum(dollars) as dollars, sum(trips) as trips, sum(units) as units
    FROM VMR_{brand_nm}_reward_levels_end
    WHERE period = 'pre-period'
    
    
    ORDER BY 3 ASC
    ''',conn)

level_ct_prep


if level_ct['dollars'].sum() == total_vmr_details_by_brand['Dollars'].max():
    print('The Dollars moved for this campaign DO MATCH with the dollars by reward level')
else:
    print('The Dollars moved for this campaign DO NOT match with the dollars by reward level')



# ## **6) GETTING DOLLARS PER TRIP BY BRAND FOR VMR PERIOD, VMR PRE-PERIOD AND YAGO (FOR REWARD ID's):**

# 6a) Getting dollars per trip and units per trip for each period by segment (for consistent Reward ID's):

# In[28]:


analyze_start = yago_print_start.strftime('%Y-%m-%d')
analyze_end = reward_print_end.strftime('%Y-%m-%d') 

curs.execute(f'''   

    DROP table if exists VMR_{brand_nm}_participants_vmr_period_segment;
    CREATE temp table VMR_{brand_nm}_participants_vmr_period_segment as
    
    
        SELECT CASE WHEN analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN 'VMR Period'
                    WHEN analysis_period = 'VMR Pre-period ({round(weeks_per_period_final)} weeks)' THEN 'VMR Pre-Period'
                    WHEN analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)' THEN 'YAGO Period'
                    END as analysis_periods,
               {segment_def},
               sum(a.purch_amt) as dollars,
               count(distinct a.ord_event_key) as trips,
               sum(a.purch_qty) as units,
               round((dollars/trips),2)::float as dollars_per_trip,
               round((units::float/trips::float),1) as units_per_trip
               
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_reward_ids_csmr d ON (a.ord_designated_cnsmr_id_key = d.ord_designated_cnsmr_id_key)

            
        WHERE brand_nbr in ({brand_nbr_str})
              and purch_qty >0
              and purch_amt >0
              and b.cal_dt >= '{analyze_start}'
            and b.cal_dt <= '{analyze_end}'
            and analysis_period IN ('VMR Print Period ({round(weeks_per_period_final)} weeks)',
                                    'VMR Pre-period ({round(weeks_per_period_final)} weeks)',
                                    'Period Year Ago ({round(weeks_per_period_final)} weeks)')
            
        
        GROUP BY 1,2
        
        ORDER BY 2 DESC
        
    DISTRIBUTE REPLICATE;

''')


# In[29]:


vmr_period_segment_sum = pd.read_sql(f'''

    SELECT *
    FROM VMR_{brand_nm}_participants_vmr_period_segment
    WHERE analysis_periods is not null
    ORDER BY 1,2,3 DESC
    
    ''',conn)

vmr_period_segment_sum


# 6b) Getting the distinct number of trips for each period to calculte %change:

# In[30]:


unique_trips_by_period = pd.read_sql(f'''

SELECT  
   CASE WHEN analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN 'VMR Period'
        WHEN analysis_period = 'VMR Pre-period ({round(weeks_per_period_final)} weeks)' THEN 'VMR Pre-Period'
        WHEN analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)' THEN 'YAGO Period'
        ELSE 'NA' END as analysis_periods,
   count(distinct a.ord_event_key) as distinct_trips

FROM ord_trd_itm_cnsmr_fact_ne_v a 

    INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
    INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
    INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
    INNER JOIN VMR_{brand_nm}_reward_ids_csmr d ON (a.ord_designated_cnsmr_id_key = d.ord_designated_cnsmr_id_key)


WHERE analysis_period IN ('VMR Print Period ({round(weeks_per_period_final)} weeks)', 'VMR Pre-period ({round(weeks_per_period_final)} weeks)', 
      'Period Year Ago ({round(weeks_per_period_final)} weeks)')
      and brand_nbr IN ({brand_nbr_str})
      and a.purch_amt > 0
      and a.purch_qty > 0
      
GROUP BY 1

ORDER BY CASE WHEN analysis_periods = 'YAGO Period' THEN 1
              WHEN analysis_periods = 'VMR Pre-Period' THEN 2
              WHEN analysis_periods = 'VMR Period' THEN 3
              END

    ''', conn)

unique_trips_by_period


# 6c) Getting the comparison in dollars per trip and units per trip for VMR Period and YAGO period (for consistent Reward ID's):

# In[31]:


#Getting VMR and YAGO data from segment table

if unique_trips_by_period.loc[0, 'analysis_periods'] != 'YAGO Period':
    
    print("NOTE: There is no YAGO data for this UPC list.")
    
else:

    vmr_data = vmr_period_segment_sum.loc[(vmr_period_segment_sum['analysis_periods']=='VMR Period'), [segment_def, 'dollars_per_trip', 'units_per_trip']]
    vmr_data.rename(columns = {'dollars_per_trip':'dollars_per_trip_vmr_period', 'units_per_trip':'units_per_trip_vmr_period'}, inplace = True)

    vmr_yago = vmr_period_segment_sum.loc[(vmr_period_segment_sum['analysis_periods']=='YAGO Period'), [segment_def, 'dollars_per_trip', 'units_per_trip']]
    vmr_yago.rename(columns = {'dollars_per_trip':'dollars_per_trip_yago_period', 'units_per_trip':'units_per_trip_yago_period'}, inplace = True)


    #Getting comparison table VMR vs YAGO and calculating %_change for each segment

    comparison_yago = pd.merge(vmr_data, vmr_yago, how="left", on=[segment_def])
    comparison_yago = comparison_yago.reindex(sorted(comparison_yago.columns), axis=1)

    comparison_yago.insert(3, '% Change Dollars',(comparison_yago["dollars_per_trip_vmr_period"]-comparison_yago["dollars_per_trip_yago_period"])/(comparison_yago["dollars_per_trip_yago_period"]))
    comparison_yago.insert(6, '% Change Units',(comparison_yago["units_per_trip_vmr_period"]-comparison_yago["units_per_trip_yago_period"])/(comparison_yago["units_per_trip_yago_period"]))


    #Calculating total %_change dollars and units for the whole brand

    total_trips_vmr = unique_trips_by_period.loc[2, 'distinct_trips']
    total_trips_yago = unique_trips_by_period.loc[0, 'distinct_trips']

    total_dollars_vmr = vmr_period_segment_sum.loc[(vmr_period_segment_sum['analysis_periods']=='VMR Period')].dollars.sum()
    total_dollars_yago = vmr_period_segment_sum.loc[(vmr_period_segment_sum['analysis_periods']=='YAGO Period')].dollars.sum()
    change_total_dollar_trip_vmr = ((total_dollars_vmr/total_trips_vmr)-(total_dollars_yago/total_trips_yago))/(total_dollars_yago/total_trips_yago)

    total_units_vmr = vmr_period_segment_sum.loc[(vmr_period_segment_sum['analysis_periods']=='VMR Period')].units.sum()
    total_units_yago = vmr_period_segment_sum.loc[(vmr_period_segment_sum['analysis_periods']=='YAGO Period')].units.sum()
    change_total_units_trip_vmr = ((total_units_vmr/total_trips_vmr)-(total_units_yago/total_trips_yago))/(total_units_yago/total_trips_yago)

    comparison_yago = comparison_yago.rename(columns={segment_def: "Segment"})

    columns_titles = ["Segment","dollars_per_trip_yago_period","dollars_per_trip_vmr_period", "% Change Dollars", "units_per_trip_yago_period", "units_per_trip_vmr_period", "% Change Units"]
    comparison_yago= comparison_yago.reindex(columns=columns_titles)

    comparison_yago.loc['Total']= ['', (total_dollars_yago/total_trips_yago), (total_dollars_vmr/total_trips_vmr), change_total_dollar_trip_vmr, (total_units_yago/total_trips_yago), (total_units_vmr/total_trips_vmr), change_total_units_trip_vmr  ]
    comparison_yago.at['Total','Segment']='Total'

    comparison_yago = comparison_yago.reset_index(drop = True)

    comparison_yago = comparison_yago.rename(columns = {

    'dollars_per_trip_yago_period':'Dollars per Trip - YAGO Period', 
    'dollars_per_trip_vmr_period':'Campaign Dollars per Trip',
    'units_per_trip_yago_period':'Units per Trip - YAGO Period', 
    'units_per_trip_vmr_period':'Campaign Units per Trip'

    })
    
    print(comparison_yago)
    
    table_slide = comparison_yago[['Segment','Campaign Dollars per Trip', 'Dollars per Trip - YAGO Period', '% Change Dollars', 'Campaign Units per Trip', 'Units per Trip - YAGO Period','% Change Units']]
    table_slide['Campaign Dollars per Trip'] = round(table_slide['Campaign Dollars per Trip'],2)
    table_slide['Campaign Units per Trip'] = round(table_slide['Campaign Units per Trip'],1)
    x_high = table_slide
    table_slide = table_slide.T.reset_index().T.reset_index(drop=True)
    table_slide.loc[:, table_slide.columns[1:]] = table_slide.loc[:, table_slide.columns[1:]].fillna(0)

    table_slide
    
    
    
    table_slide_s = comparison_yago[['Segment','Campaign Dollars per Trip', '% Change Dollars', 'Campaign Units per Trip', '% Change Units']]
    table_slide_s['Campaign Dollars per Trip'] = round(table_slide_s['Campaign Dollars per Trip'],2)
    table_slide_s['Campaign Units per Trip'] = round(table_slide_s['Campaign Units per Trip'],1)
    table_slide_s = table_slide_s.T.reset_index().T.reset_index(drop=True)

    table_slide_s.loc[:, table_slide_s.columns[1:]] = table_slide_s.loc[:, table_slide_s.columns[1:]].fillna(0)
    table_slide_s 


# 6d) Getting the comparison in dollars per trip and units per trip for VMR Period and VMR Pre-Period period (for consistent Reward ID's):

# In[32]:


#Getting VMR and YAGO data from segment table

vmr_data = vmr_period_segment_sum.loc[(vmr_period_segment_sum['analysis_periods']=='VMR Period'), [segment_def, 'dollars_per_trip', 'units_per_trip']]
vmr_data.rename(columns = {'dollars_per_trip':'dollars_per_trip_vmr_period', 'units_per_trip':'units_per_trip_vmr_period'}, inplace = True)

vmr_preperiod = vmr_period_segment_sum.loc[(vmr_period_segment_sum['analysis_periods']=='VMR Pre-Period'), [segment_def, 'dollars_per_trip', 'units_per_trip']]
vmr_preperiod.rename(columns = {'dollars_per_trip':'dollars_per_trip_pre_period', 'units_per_trip':'units_per_trip_pre_period'}, inplace = True)


#Getting comparison table VMR vs YAGO and calculating %_change for each segment

comparison_preperiod = pd.merge(vmr_data, vmr_preperiod, how="left", on=[segment_def])
comparison_preperiod = comparison_preperiod.reindex(sorted(comparison_preperiod.columns), axis=1)

comparison_preperiod.insert(3, '% Change Dollars',(comparison_preperiod["dollars_per_trip_vmr_period"]-comparison_preperiod["dollars_per_trip_pre_period"])/(comparison_preperiod["dollars_per_trip_pre_period"]))
comparison_preperiod.insert(6, '% Change Units',(comparison_preperiod["units_per_trip_vmr_period"]-comparison_preperiod["units_per_trip_pre_period"])/(comparison_preperiod["units_per_trip_pre_period"]))


#Calculating total %_change dollars and units for the whole brand

total_trips_vmr = unique_trips_by_period.loc[(unique_trips_by_period['analysis_periods'] == 'VMR Period'), 'distinct_trips'].values[0]
total_trips_pperiod = unique_trips_by_period.loc[(unique_trips_by_period['analysis_periods'] == 'VMR Pre-Period'), 'distinct_trips'].values[0]

total_dollars_vmr = vmr_period_segment_sum.loc[(vmr_period_segment_sum['analysis_periods']=='VMR Period')].dollars.sum()
total_dollars_pperiod = vmr_period_segment_sum.loc[(vmr_period_segment_sum['analysis_periods']=='VMR Pre-Period')].dollars.sum()
change_total_dollar_trip_vmr = ((total_dollars_vmr/total_trips_vmr)-(total_dollars_pperiod/total_trips_pperiod))/(total_dollars_pperiod/total_trips_pperiod)

total_units_vmr = vmr_period_segment_sum.loc[(vmr_period_segment_sum['analysis_periods']=='VMR Period')].units.sum()
total_units_pperiod = vmr_period_segment_sum.loc[(vmr_period_segment_sum['analysis_periods']=='VMR Pre-Period')].units.sum()
change_total_units_trip_vmr = ((total_units_vmr/total_trips_vmr)-(total_units_pperiod/total_trips_pperiod))/(total_units_pperiod/total_trips_pperiod)
                       
comparison_preperiod = comparison_preperiod.rename(columns={segment_def: "Segment"})

columns_titles = ["Segment","dollars_per_trip_pre_period","dollars_per_trip_vmr_period", "% Change Dollars", "units_per_trip_pre_period", "units_per_trip_vmr_period", "% Change Units"]
comparison_preperiod= comparison_preperiod.reindex(columns=columns_titles)

comparison_preperiod.loc['Total']= ['', (total_dollars_pperiod/total_trips_pperiod), (total_dollars_vmr/total_trips_vmr), change_total_dollar_trip_vmr, (total_units_pperiod/total_trips_pperiod), (total_units_vmr/total_trips_vmr), change_total_units_trip_vmr  ]
comparison_preperiod.at['Total',"Segment"]='Total'

comparison_preperiod = comparison_preperiod.reset_index(drop = True)

comparison_preperiod = comparison_preperiod.rename(columns = {

'dollars_per_trip_pre_period':'Dollars per Trip - Pre Period', 
'dollars_per_trip_vmr_period':'Campaign Dollars per Trip',
'units_per_trip_pre_period':'Units per Trip - Pre Period', 
'units_per_trip_vmr_period':'Campaign Units per Trip'
    
})

comparison_preperiod

table_slide_2 = comparison_preperiod[['Segment','Campaign Dollars per Trip', 'Dollars per Trip - Pre Period','% Change Dollars', 'Campaign Units per Trip', 'Units per Trip - Pre Period','% Change Units']]
table_slide_2['Campaign Dollars per Trip'] = round(table_slide_2['Campaign Dollars per Trip'],2)
table_slide_2['Campaign Units per Trip'] = round(table_slide_2['Campaign Units per Trip'],1)
table_slide_2 = table_slide_2.T.reset_index().T.reset_index(drop=True)
table_slide_2.loc[:, table_slide_2.columns[1:]] = table_slide_2.loc[:, table_slide_2.columns[1:]].fillna(0)
table_slide_2


# 6e) Getting the comparison in dollars per trip and units per trip for VMR Period and VMR Pre 52 wks period (for consistent Reward ID's):


analyze_start_2 = prior_period_start.strftime('%Y-%m-%d')

curs.execute(f'''   

    DROP table if exists VMR_{brand_nm}_participants_vmr_period_segment_pre52;
    CREATE temp table VMR_{brand_nm}_participants_vmr_period_segment_pre52 as
    
    
        SELECT CASE WHEN brand_nbr in ({brand_nbr_str}) and analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN 'VMR Period'
                    WHEN brand_nbr in ({brand_nbr_str}) and prior_period = 'Prior-period ({round(pre_weeks)} weeks)' THEN '52wk Pre-Period'
                    END as analysis_periods,
               {segment_def},
               sum(a.purch_amt) as dollars,
               count(distinct a.ord_event_key) as trips,
               sum(a.purch_qty) as units,
               round((dollars/trips),2)::float as dollars_per_trip,
               round((units::float/trips::float),1) as units_per_trip
               
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_reward_ids_csmr d ON (a.ord_designated_cnsmr_id_key = d.ord_designated_cnsmr_id_key)

            
        WHERE 
              purch_qty >0
              and purch_amt >0
              and b.cal_dt >= '{analyze_start_2}'
            and b.cal_dt <= '{analyze_end}'
            and prior_period IN ('Prior-period ({round(pre_weeks)} weeks)') OR 
                analysis_period IN ('VMR Print Period ({round(weeks_per_period_final)} weeks)')            
        
        GROUP BY 1,2
        
        ORDER BY 2 DESC
        
    DISTRIBUTE REPLICATE;

''')


vmr_period_segment_sum_pre52 = pd.read_sql(f'''

    SELECT *
    FROM VMR_{brand_nm}_participants_vmr_period_segment_pre52
    WHERE analysis_periods is not null 
    ORDER BY 1,2,3 DESC
    
    ''',conn)

vmr_period_segment_sum_pre52


unique_trips_by_period_52wk = pd.read_sql(f'''

SELECT  
   CASE WHEN prior_period = 'Prior-period ({round(pre_weeks)} weeks)' THEN 'VMR 52wk PrePeriod'
        ELSE 'NA' END as analysis_periods,
   count(distinct a.ord_event_key) as distinct_trips

FROM ord_trd_itm_cnsmr_fact_ne_v a 

    INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
    INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
    INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
    INNER JOIN VMR_{brand_nm}_reward_ids_csmr d ON (a.ord_designated_cnsmr_id_key = d.ord_designated_cnsmr_id_key)


WHERE prior_period IN ('Prior-period ({round(pre_weeks)} weeks)')
      and brand_nbr IN ({brand_nbr_str})
      and a.purch_amt > 0
      and a.purch_qty > 0
      
GROUP BY 1


    ''', conn)

unique_trips_by_period_52wk


#Getting VMR and Pre 52wk data from segment table

if unique_trips_by_period_52wk.loc[0, 'analysis_periods'] != 'VMR 52wk PrePeriod':
    
    print("NOTE: There is no Pre 52 wk data for this UPC list.")
    
else:


    vmr_data = vmr_period_segment_sum_pre52.loc[(vmr_period_segment_sum_pre52['analysis_periods']=='VMR Period'), [segment_def, 'dollars_per_trip', 'units_per_trip']]
    vmr_data.rename(columns = {'dollars_per_trip':'dollars_per_trip_vmr_period', 'units_per_trip':'units_per_trip_vmr_period'}, inplace = True)

    vmr_preperiod_52wk = vmr_period_segment_sum_pre52.loc[(vmr_period_segment_sum_pre52['analysis_periods']=='52wk Pre-Period'), [segment_def, 'dollars_per_trip', 'units_per_trip']]
    vmr_preperiod_52wk.rename(columns = {'dollars_per_trip':'dollars_per_trip_pre_period', 'units_per_trip':'units_per_trip_pre_period'}, inplace = True)


    #Getting comparison table VMR vs YAGO and calculating %_change for each segment

    comparison_preperiod_52wk = pd.merge(vmr_data, vmr_preperiod_52wk, how="left", on=[segment_def])
    comparison_preperiod_52wk = comparison_preperiod_52wk.reindex(sorted(comparison_preperiod_52wk.columns), axis=1)

    comparison_preperiod_52wk.insert(3, '% Change Dollars',(comparison_preperiod_52wk["dollars_per_trip_vmr_period"]-comparison_preperiod_52wk["dollars_per_trip_pre_period"])/(comparison_preperiod_52wk["dollars_per_trip_pre_period"]))
    comparison_preperiod_52wk.insert(6, '% Change Units',(comparison_preperiod_52wk["units_per_trip_vmr_period"]-comparison_preperiod_52wk["units_per_trip_pre_period"])/(comparison_preperiod_52wk["units_per_trip_pre_period"]))


    #Calculating total %_change dollars and units for the whole brand

    total_trips_vmr = unique_trips_by_period.loc[(unique_trips_by_period['analysis_periods'] == 'VMR Period'), 'distinct_trips'].values[0]
    total_trips_pperiod = unique_trips_by_period_52wk.loc[(unique_trips_by_period_52wk['analysis_periods'] == 'VMR 52wk PrePeriod'), 'distinct_trips'].values[0]

    total_dollars_vmr = vmr_period_segment_sum_pre52.loc[(vmr_period_segment_sum_pre52['analysis_periods']=='VMR Period')].dollars.sum()
    total_dollars_pperiod = vmr_period_segment_sum_pre52.loc[(vmr_period_segment_sum_pre52['analysis_periods']=='52wk Pre-Period')].dollars.sum()
    change_total_dollar_trip_vmr = ((total_dollars_vmr/total_trips_vmr)-(total_dollars_pperiod/total_trips_pperiod))/(total_dollars_pperiod/total_trips_pperiod)

    total_units_vmr = vmr_period_segment_sum_pre52.loc[(vmr_period_segment_sum_pre52['analysis_periods']=='VMR Period')].units.sum()
    total_units_pperiod = vmr_period_segment_sum_pre52.loc[(vmr_period_segment_sum_pre52['analysis_periods']=='52wk Pre-Period')].units.sum()
    change_total_units_trip_vmr = ((total_units_vmr/total_trips_vmr)-(total_units_pperiod/total_trips_pperiod))/(total_units_pperiod/total_trips_pperiod)

    comparison_preperiod_52wk = comparison_preperiod_52wk.rename(columns={segment_def: "Segment"})

    columns_titles = ["Segment","dollars_per_trip_pre_period","dollars_per_trip_vmr_period", "% Change Dollars", "units_per_trip_pre_period", "units_per_trip_vmr_period", "% Change Units"]
    comparison_preperiod_52wk= comparison_preperiod_52wk.reindex(columns=columns_titles)

    comparison_preperiod_52wk.loc['Total']= ['', (total_dollars_pperiod/total_trips_pperiod), (total_dollars_vmr/total_trips_vmr), change_total_dollar_trip_vmr, (total_units_pperiod/total_trips_pperiod), (total_units_vmr/total_trips_vmr), change_total_units_trip_vmr  ]
    comparison_preperiod_52wk.at['Total',"Segment"]='Total'

    comparison_preperiod_52wk = comparison_preperiod_52wk.reset_index(drop = True)

    comparison_preperiod_52wk = comparison_preperiod_52wk.rename(columns = {

    'dollars_per_trip_pre_period':'Dollars per Trip - Prior 52wk', 
    'dollars_per_trip_vmr_period':'Campaign Dollars per Trip',
    'units_per_trip_pre_period':'Units per Trip - Prior 52wk', 
    'units_per_trip_vmr_period':'Campaign Units per Trip'

    })

    table_slide_1 = comparison_preperiod_52wk[['Segment','Campaign Dollars per Trip', 'Dollars per Trip - Prior 52wk', '% Change Dollars', 'Campaign Units per Trip', 'Units per Trip - Prior 52wk','% Change Units']]
    table_slide_1['Campaign Dollars per Trip'] = round(table_slide_1['Campaign Dollars per Trip'],2)
    table_slide_1['Campaign Units per Trip'] = round(table_slide_1['Campaign Units per Trip'],1)
    table_slide_1 = table_slide_1.T.reset_index().T.reset_index(drop=True)
    table_slide_1.loc[:, table_slide_1.columns[1:]] = table_slide_1.loc[:, table_slide_1.columns[1:]].fillna(0)

    print(comparison_preperiod_52wk)
    
    
vmr_data = vmr_period_segment_sum_pre52.loc[(vmr_period_segment_sum_pre52['analysis_periods']=='VMR Period'), [segment_def, 'dollars', 'units']]
vmr_data.rename(columns = {'dollars':'dollars_vmr_period', 'units':'units_vmr_period'}, inplace = True)

vmr_preperiod_52wk = vmr_period_segment_sum_pre52.loc[(vmr_period_segment_sum_pre52['analysis_periods']=='52wk Pre-Period'), [segment_def, 'dollars', 'units']]
vmr_preperiod_52wk.rename(columns = {'dollars':'dollars_vmr_preperiod', 'units':'units_vmr_preperiod'}, inplace = True)

#Getting comparison table VMR vs YAGO and calculating %_change for each segment

comparison_preperiod_52wk_chart = pd.merge(vmr_data, vmr_preperiod_52wk, how="left", on=[segment_def])
comparison_preperiod_52wk_chart = comparison_preperiod_52wk_chart.reindex(sorted(comparison_preperiod_52wk_chart.columns), axis=1)

comparison_preperiod_52wk_chart.insert(3, '% Change Dollars',(comparison_preperiod_52wk_chart["dollars_vmr_period"]-comparison_preperiod_52wk_chart["dollars_vmr_preperiod"])/(comparison_preperiod_52wk_chart["dollars_vmr_preperiod"]))
comparison_preperiod_52wk_chart.insert(6, '% Change Units',(comparison_preperiod_52wk_chart["units_vmr_period"]-comparison_preperiod_52wk_chart["units_vmr_preperiod"])/(comparison_preperiod_52wk_chart["units_vmr_preperiod"]))

total_vmr_details_segm_chart = comparison_preperiod_52wk_chart[segment_def]
total_vmr_details_units_chart = comparison_preperiod_52wk_chart["% Change Units"]
total_vmr_details_dollars_chart = comparison_preperiod_52wk_chart["% Change Dollars"]

comparison_preperiod_52wk_chart


# ## **7) SIZE BASKETS DURING VMR PERIOD AND VMR REDEMPTION PERIOD:**

# In[33]:


#Modified code from Arthur Li's previous code.

#Create New date filter to include redemption period:

curs.execute(f'''

    DROP table if exists VMR_{brand_nm}_red_date_filter;
    CREATE temp table VMR_{brand_nm}_red_date_filter as
    
        SELECT distinct date_key, 
               cal_dt, 
               cal_sun_wk_ending_dt, 
               cal_sun_wk_ending_rank_nbr
        FROM date_v
        WHERE cal_dt between '{reward_print_start}' and '{redemption_end}'
        
    DISTRIBUTE REPLICATE;
    
''')    
    
    
#Select Redemption trips:

curs.execute(f'''
    
    DROP table if exists VMR_{brand_nm}_redemption_trips; 
    CREATE temp table VMR_{brand_nm}_redemption_trips as 
        
        SELECT distinct a.ord_event_key
        
        FROM ORD_RED_CNSMR_FACT_NE_V a
        
            INNER JOIN VMR_{brand_nm}_red_date_filter d ON (a.ord_date_key = d.date_key)
            INNER JOIN VMR_{brand_nm}_promo_filter pv ON (a.promo_varnt_key=pv.promo_varnt_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
        
        DISTRIBUTE ON (ord_event_key);
''')
        
        
#Select All trips during Reward & Redemption Period with labelling Reward, Redemption and AO Trips:

curs.execute(f'''

    DROP table if exists VMR_{brand_nm}_redemption_period;
    CREATE temp table VMR_{brand_nm}_redemption_period as 
        
        SELECT a.ord_event_key, 
               a.tot_ord_amt,
               CASE WHEN r.ord_event_key is not null then 'Reward Trip'
                    WHEN r.ord_event_key is null and rd.ord_event_key is null and d.cal_dt between '{reward_print_start}' and '{reward_print_end}' then 'AO Trips - Reward Period'
                    ELSE 'NA' END as reward_p_flag,
                
               CASE WHEN r.ord_event_key is null and rd.ord_event_key is not null then 'Redemption Trip'
                    WHEN r.ord_event_key is null and rd.ord_event_key is null then 'AO Trips - Redemption Period'
                    ELSE 'NA' END as redemp_p_flag
        
        FROM ord_sum_cnsmr_fact_ne_v a
        
            INNER JOIN VMR_{brand_nm}_red_date_filter d ON (a.ord_date_key = d.date_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
            LEFT JOIN VMR_{brand_nm}_reward_events r ON (a.ord_event_key = r.ord_event_key)
            LEFT JOIN VMR_{brand_nm}_redemption_trips rd ON (rd.ord_event_key = a.ord_event_key)
            
        DISTRIBUTE RANDOM;   
       
''')


# In[34]:


nbr_reward_trips = pd.read_sql(f'''

    SELECT COUNT(distinct a.ord_event_key) as trips,
           SUM(a.tot_ord_amt) as dollars,
           dollars::float/trips::float as avg_basket_size          
    FROM ord_sum_cnsmr_fact_ne_v a
        INNER JOIN VMR_{brand_nm}_redemption_period b ON (b.ord_event_key = a.ord_event_key)
    WHERE reward_p_flag = 'Reward Trip'
        
    ''',conn)

print('The number of reward trips are: '+ str(nbr_reward_trips['trips'].values[0]))
print('The total dollars for all the reward trips is: '+ str(nbr_reward_trips['dollars'].values[0]))
print('The avg basket dollars for all the reward trips is: '+ str(nbr_reward_trips['avg_basket_size'].values[0]))


# In[35]:


nbr_reward_promoted_trips = pd.read_sql(f'''

    SELECT COUNT(distinct a.ord_event_key) as trips,
           SUM(a.purch_amt) as dollars,
           dollars::float/{nbr_reward_trips['trips'].values[0]}::float as avg_basket_size
    FROM ord_trd_itm_cnsmr_fact_ne_v a
        INNER JOIN VMR_{brand_nm}_redemption_period b ON (b.ord_event_key = a.ord_event_key)
        INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
          
        WHERE brand_nbr in ({brand_nbr_str}) and reward_p_flag = 'Reward Trip'

        
    ''',conn)

print('The number of reward trips are: '+ str(nbr_reward_promoted_trips['trips'].values[0]))
print('The total dollars of promoted product in reward trips is: '+ str(nbr_reward_promoted_trips['dollars'].values[0]))
print('The avg basket dollars of promoted product in reward trips is: '+ str(nbr_reward_promoted_trips['avg_basket_size'].values[0]))
print('The % promoted product accounted in reward baskets: '+ str(round(nbr_reward_promoted_trips['avg_basket_size'].values[0]/nbr_reward_trips['avg_basket_size'].values[0]*100)))


# In[36]:


nbr_ao_reward_trips = pd.read_sql(f'''

    SELECT count(distinct a.ord_event_key)
    FROM ord_trd_itm_cnsmr_fact_ne_v a
        INNER JOIN VMR_{brand_nm}_redemption_period b ON (b.ord_event_key = a.ord_event_key)
        INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
    WHERE reward_p_flag != 'NA'
          and c.brand_nbr IN ({brand_nbr_str})

    ''',conn)


print('The number of all brand transactions during Reward period is: '+ str(nbr_ao_reward_trips['count'].values[0]))
pct_reward_trips = round((total_reward_trans_var['count'].values[0]/nbr_ao_reward_trips['count'].values[0]),5)
print('The pct of reward trips among all brand transactions during VMR campaign is: '+ str(pct_reward_trips*100) + '%')


# In[37]:


#Same than above, just for double checking

nbr_ao_reward_trips = pd.read_sql(f'''

    SELECT count(distinct a.ord_event_key)
    
    FROM ord_trd_itm_cnsmr_fact_ne_v a
    
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
            LEFT JOIN VMR_{brand_nm}_reward_events r ON (a.ord_event_key = r.ord_event_key)
            LEFT JOIN VMR_{brand_nm}_redemption_trips rd ON (rd.ord_event_key = a.ord_event_key)
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)

    WHERE b.analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' and 
     r.ord_event_key IS NULL and rd.ord_event_key IS NULL and brand_nbr IN ({brand_nbr_str})

    ''',conn)

print('The number of reward trips during VMR Campaign is: '+ str(total_reward_trans_var['count'].values[0]))
print('The number of all brand trips during Reward period is: '+ str(nbr_ao_reward_trips['count'].values[0]))

pct_reward_trips = round((total_reward_trans_var['count'].values[0])/(nbr_ao_reward_trips['count'].values[0]+total_reward_trans_var['count'].values[0]),5)
print('The pct of reward trips among all brand transactions during VMR campaign is: '+ str(pct_reward_trips*100) + '%')


# In[38]:


nbr_redemp_trips = pd.read_sql(f'''

    SELECT COUNT(distinct a.ord_event_key) as trips,
           SUM(tot_ord_amt) as dollars,
           dollars::float/trips::float as avg_basket_size          
    FROM ord_sum_cnsmr_fact_ne_v a
        INNER JOIN VMR_{brand_nm}_redemption_trips b ON (b.ord_event_key = a.ord_event_key)
        INNER JOIN consumer_id_v p ON (a.ord_designated_cnsmr_id_key = p.cnsmr_id_key)
        
    ''',conn)

print('The number of redemption trips are: '+ str(nbr_redemp_trips['trips'].values[0]))
print('The total dollars for all the redemption trips is: '+ str(nbr_redemp_trips['dollars'].values[0]))
print('The avg basket dollars for all the redemption trips is: '+ str(nbr_redemp_trips['avg_basket_size'].values[0]))


# In[39]:


nbr_redemp_promoted = pd.read_sql(f'''

    SELECT COUNT(distinct a.ord_event_key) as trips,
           SUM(a.purch_amt) as dollars,
           dollars::float/{nbr_redemp_trips['trips'].values[0]}::float as avg_basket_size
    FROM ord_trd_itm_cnsmr_fact_ne_v a
        INNER JOIN VMR_{brand_nm}_redemption_trips b ON (b.ord_event_key = a.ord_event_key)
        INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
        INNER JOIN consumer_id_v p ON (a.ord_designated_cnsmr_id_key = p.cnsmr_id_key)
          
        WHERE brand_nbr in ({brand_nbr_str})

    ''',conn)

print('The number of redemption trips that includes the promoted product is: '+ str(nbr_redemp_promoted['trips'].values[0]))
print('The total dollars of promoted product in the redemption baskets is: '+ str(nbr_redemp_promoted['dollars'].values[0]))
print('The avg dollars of promoted product in redemption baskets is: '+ str(nbr_redemp_promoted['avg_basket_size'].values[0]))

if nbr_redemp_promoted['trips'].values[0] != 0:

    print('The % promoted product accounted in redemption baskets: '+ str(round(nbr_redemp_promoted['avg_basket_size'].values[0]/nbr_redemp_trips['avg_basket_size'].values[0]*100)))


# In[40]:


if nbr_redemp_promoted['trips'].values[0] != 0:

    #Getting summarize table:

    curs.execute(f'''

        DROP table if exists VMR_{brand_nm}_baskets_results; 
        CREATE temp table VMR_{brand_nm}_baskets_results as


        SELECT reward_p_flag as analysis_period, 
               sum(tot_ord_amt) as Dollars, 
               count(distinct ord_event_key) as Trips, 
               dollars::float/trips::float as avg_basket_size          
        FROM VMR_{brand_nm}_redemption_period
        WHERE reward_p_flag != 'NA'
        GROUP BY 1

        UNION

        SELECT redemp_p_flag as analysis_period, 
               sum(tot_ord_amt) as dollars, 
               count(distinct ord_event_key) as trips, 
               dollars::float/trips::float as avg_basket_size
        FROM VMR_{brand_nm}_redemption_period
        WHERE redemp_p_flag != 'NA'
        GROUP BY 1

        DISTRIBUTE RANDOM;   

    ''')    


    df_vmr_ttls_basket = pd.read_sql(f''' 

    SELECT * FROM VMR_{brand_nm}_baskets_results

    ORDER BY case WHEN analysis_period = 'Reward Trip' THEN 1
                  WHEN analysis_period = 'AO Trips - Reward Period' THEN 2
                  WHEN analysis_period = 'Redemption Trip' THEN 3
                  WHEN analysis_period = 'AO Trips - Redemption Period' THEN 4 END 
    ''',conn)

    df_vmr_ttls_basket.rename(columns = {'analysis_period':'Analysis period', 'avg_basket_size':'Average Basket Size', 'dollars':'Dollars', 'trips':'Trips'}, inplace = True)


    df_vmr_ttls_basket.loc[4] = ['% change reward trip', '', '', (df_vmr_ttls_basket.loc[0,'Average Basket Size']-df_vmr_ttls_basket.loc[1,'Average Basket Size'])/(df_vmr_ttls_basket.loc[1,'Average Basket Size'])]
    df_vmr_ttls_basket.loc[5]= ['% change redemption trip', '', '', (df_vmr_ttls_basket.loc[2,'Average Basket Size']-df_vmr_ttls_basket.loc[3,'Average Basket Size'])/(df_vmr_ttls_basket.loc[3,'Average Basket Size'])]

    df_vmr_ttls_basket.iloc[2], df_vmr_ttls_basket.iloc[4] = df_vmr_ttls_basket.iloc[4], df_vmr_ttls_basket.iloc[2]
    df_vmr_ttls_basket.iloc[4], df_vmr_ttls_basket.iloc[3] = df_vmr_ttls_basket.iloc[3], df_vmr_ttls_basket.iloc[4]


    print(df_vmr_ttls_basket)
    
else:
    
    curs.execute(f'''

        DROP table if exists VMR_{brand_nm}_baskets_results; 
        CREATE temp table VMR_{brand_nm}_baskets_results as


        SELECT reward_p_flag as analysis_period, 
               sum(tot_ord_amt) as Dollars, 
               count(distinct ord_event_key) as Trips, 
               dollars::float/trips::float as avg_basket_size          
        FROM VMR_{brand_nm}_redemption_period
        WHERE reward_p_flag != 'NA'
        GROUP BY 1

        DISTRIBUTE RANDOM;   

    ''')    


    df_vmr_ttls_basket = pd.read_sql(f''' 

    SELECT * FROM VMR_{brand_nm}_baskets_results

    ''',conn)

    df_vmr_ttls_basket.rename(columns = {'analysis_period':'Analysis period', 'avg_basket_size':'Average Basket Size', 'dollars':'Dollars', 'trips':'Trips'}, inplace = True)


    #df_vmr_ttls_basket.loc[2] = ['% change reward trip', '', '', (df_vmr_ttls_basket.loc[0,'Average Basket Size']-df_vmr_ttls_basket.loc[1,'Average Basket Size'])/(df_vmr_ttls_basket.loc[1,'Average Basket Size'])]

    #df_vmr_ttls_basket.iloc[2] = df_vmr_ttls_basket.iloc[4], df_vmr_ttls_basket.iloc[2]


    print(df_vmr_ttls_basket)
    
    print('The number of redemption trips is zero, so this section is not going to be shown.')


# ## **8) BRAND COMBINATIONS PURCHASED ON THE VMR REWARD TRIP (FOR TRACKABLE ID's):**

# Evaluating the limit of segments to perform the combinations:
# (It must be less than 53 segments. If not, the code will filter by those segemnts with % reward trips > 1)

# In[41]:


df_segments = df_segments_2
nbr_segments = len(total_vmr_details_by_brand[total_vmr_details_by_brand['Segment'] != 'Total'])

if nbr_segments<53:
    
    df_segments = df_segments[['segment', 'segm_nbr']]
    
    yb_load(Df = df_segments,
        table_name = f'''VMR_{brand_nm}_segments_{analyst}''',
        userid = yb_user,
        passwd = readpw("Yellowbrick"),
        append = False,
        database= f'{dbase}')

    check_upload = pd.read_sql(f"""select count(*) from VMR_{brand_nm}_segments_{analyst}""", conn)
    print('Uploaded number of records: ' + str(check_upload.values[0]))
    
    df_segments = df_segments.rename(columns={'segment': 'Segment', 'segm_nbr':'Segm Nbr'})
    print(df_segments)
    
    
else:
    
    if segment_type == 1:
    
        df_segments = pd.read_sql(f'''

        SELECT {segment_def} as segment,
               brand_nbr as segm_nbr,
               count(distinct o.ord_event_key) as reward_trips,
               {total_reward_trans} as total_reward_trips,
               round((reward_trips::float/total_reward_trips::float)*100,2) as pct_reward_trips

        FROM ord_trd_itm_cnsmr_fact_ne_v o

        INNER JOIN VMR_{brand_nm}_upc_filter b ON (o.trade_item_key = b.trade_item_key)
        INNER JOIN VMR_{brand_nm}_reward_events c ON (o.ord_event_key = c.ord_event_key)

        GROUP BY 1

        HAVING pct_reward_trips > 1

        ORDER BY 2 ASC

        ''',conn)

        
        df_segments = df_segments[['segment', 'segm_nbr']]
        
        # Ensure we don't exceed 52 segments (letter limit A-Z, a-z)
        if len(df_segments) > 52:
            print(f"Still have {len(df_segments)} segments after filtering. Keeping top 52 by segm_nbr.")
            df_segments = df_segments.head(52)

        yb_load(Df = df_segments,
            table_name = f'''VMR_{brand_nm}_segments_{analyst}''',
            userid = yb_user,
            passwd = readpw("Yellowbrick"),
            append = False,
            database= f'{dbase}')

        check_upload = pd.read_sql(f"""select count(*) from VMR_{brand_nm}_segments_{analyst}""", conn)
        print('Uploaded number of records: ' + str(check_upload.values[0]))

        df_segments = df_segments.rename(columns={"segm_nbr": "Segm Nbr", "segment": "Segment"})
        print("The number of segments is over the max limit (53). Segments were filtered by reward trips > 1%")
        print(df_segments)
        
    else:
        
        df_segments = pd.read_sql(f'''

        SELECT {segment_def} as segment,
               count(distinct o.ord_event_key) as reward_trips,
               {total_reward_trans} as total_reward_trips,
               round((reward_trips::float/total_reward_trips::float)*100,2) as pct_reward_trips

        FROM ord_trd_itm_cnsmr_fact_ne_v o

        INNER JOIN VMR_{brand_nm}_upc_filter b ON (o.trade_item_key = b.trade_item_key)
        INNER JOIN VMR_{brand_nm}_reward_events c ON (o.ord_event_key = c.ord_event_key)

        GROUP BY 1

        HAVING pct_reward_trips > 1

        ORDER BY 1

        ''',conn)

        # Ensure we don't exceed 52 segments (letter limit A-Z, a-z)
        if len(df_segments) > 52:
            print(f"Still have {len(df_segments)} segments after filtering. Keeping top 52.")
            df_segments = df_segments.head(52)

        df_segments['segm_nbr'] = list(range(1,len(df_segments)+1))
        df_segments = df_segments[['segment', 'segm_nbr']]
        

        yb_load(Df = df_segments,
            table_name = f'''VMR_{brand_nm}_segments_{analyst}''',
            userid = yb_user,
            passwd = readpw("Yellowbrick"),
            append = False,
            database= f'{dbase}')

        check_upload = pd.read_sql(f"""select count(*) from VMR_{brand_nm}_segments_{analyst}""", conn)
        print('Uploaded number of records: ' + str(check_upload.values[0]))

        df_segments = df_segments.rename(columns={"segm_nbr": "Segm Nbr", "segment": "Segment"})
        print("The number of segments is over the max limit (53). Segments were filtered by reward trips > 1%")
        print(df_segments)
        


# In[42]:


#Modified from K.Mertens's J&J custom VMR code.

#Functions to rename columns in the dataframes, Don't change it!

renaming_dict = {
    'ord_event_key' : 'ord_event_key',
    'cnsmr_id_key' : 'cnsmr_id_key',
    'analysis_period_p1p2': 'Analysis Period', 
    'analysis_period_rec52': 'Analysis Period', 
    'analysis_periods': 'Analysis Period', 
    'cal_sun_wk_ending_dt': 'WE Date',
    'nbr_of_upcs': 'Raw # of UPCs included',
    'brand_desc': 'Brand',
    'sur_seg': 'Purchasing Segment',
    'price_seg': 'Pricing Segment',
    'fin_cmit_contract_nbr': 'Contract Nbr',
    'promo_src_id_txt': 'BL',
    'segment':'Segment',
    'segm_nbr':'Segm Nbr',
    'dollars':'Dollars',
    'units':'Units',
    'trips':'Trips',
    'dollars_per_trip':'Dollars Per Trip',
    'units_per_trip':'Units Per Trip',
    'pct_of_trips' : 'Pct Of Trips',
    'pct_of_dollars' : 'Pct Of Dollars',
    'pct_of_units' : 'Pct Of Units'
    
}

format_mapping = {
    'Raw # of UPCs included' : "{:,.0f}",
    'Contract Nbr': "{:.0f}"
}



# In[43]:


#Modified code from K.Mertens's J&J custom VMR code.

###### First make the function to get the ascii letter A-Z,a-z

def get_letter_for_brand_nbr(nbr):
    '''
    Given a number 1-52, this returns A-Z, a-z.
    '''
    if nbr < 27 :
        return chr(nbr+ord("A")-1)
    elif nbr < 53 :
        return chr(nbr+ord("a")-27)
    raise IndexError(f"This function can only work with up to 52 brand_nbr items.  You gave me {nbr}")
    
###### Make a dictionary that maps from the Letter to the Brand Description

alpha_to_brand_desc = {}

idx = {name: i for i, name in enumerate(list(df_segments),1)}
for row in df_segments.itertuples():
    curr_letter = get_letter_for_brand_nbr(row[idx['Segm Nbr']])
    curr_desc = row[idx['Segment']]
    alpha_to_brand_desc[curr_letter] = curr_desc

print(alpha_to_brand_desc)


# In[44]:


#Creating a table that contains all the reward trips details: ord_event_key, cnmr_id, brand_nbr, brand_descr

curs.execute(f'''

    DROP table if exists VMR_{brand_nm}_purch_details_all_trips;
    CREATE temp table VMR_{brand_nm}_purch_details_all_trips as
        
          SELECT distinct ou.ord_event_key
                 ,d.segment as segment
                 ,d.segmnbr as segm_nbr                
                 ,b.cal_dt
                 ,b.cal_sun_wk_ending_dt
                 ,b.analysis_period
                 ,b.prior_period
                 ,c.brand_nbr
                 ,ou.ord_designated_cnsmr_id_key as cnsmr_id_key
                 ,sum(ou.purch_qty) as brand_units
                 ,sum(ou.purch_amt) as brand_dollars
          
          FROM ord_trd_itm_cnsmr_fact_ne_v ou
                 
                 INNER JOIN VMR_{brand_nm}_date_filter b ON (ou.ord_date_key = b.date_key)
                 INNER JOIN VMR_{brand_nm}_upc_filter c ON (ou.trade_item_key = c.trade_item_key)
                 INNER JOIN VMR_{brand_nm}_segments_{analyst} d ON (c.{segment_def} = d.segment)
                 
          
          WHERE ou.purch_amt > 0
                and ou.purch_qty > 0
                and c.brand_nbr in ({brand_nbr_str})
        
        GROUP BY 1,2,3,4,5,6,7,8,9
        
        DISTRIBUTE ON (ord_event_key);
        
''')


# In[45]:


#Modified from K.Mertens's J&J custom VMR code.

_skip_python_combo_processing = False  # Initialize flag - will be set to True if SQL aggregation succeeds

if nbr_segments > 1 :
    
####  Download purchase details from YB, organized by ord_event_key
####  UPDATED: Use SQL-based aggregation to compute brand combos in DB
####           This avoids downloading 5M+ rows when segment count is high

    # Step 1: Create base table with segment details per transaction
    curs.execute(f'''
    
        DROP table if exists VMR_{brand_nm}_combos_results_{analyst}; 
        CREATE table VMR_{brand_nm}_combos_results_{analyst} as
    
        SELECT a.ord_event_key,
               a.segm_nbr,
               a.segment
              ,sum(a.brand_units) as units
              ,sum(a.brand_dollars) as dollars 
              
        FROM VMR_{brand_nm}_purch_details_all_trips a
        
        INNER JOIN VMR_{brand_nm}_reward_events rd ON (rd.ord_event_key = a.ord_event_key) 
        INNER JOIN consumer_id_v p ON (a.cnsmr_id_key = p.cnsmr_id_key)

          WHERE p.abuser_ind= 'N' 
                and p.unident_ind = 'N'
                and p.really_bad_abuser_ind='N'
                and p.cnsmr_id_key>0
                and analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)'
        GROUP by 1,2,3
    
        DISTRIBUTE RANDOM;
    
    
    ''')
    
    # Step 2: Try SQL pre-aggregation (faster, but requires STRING_AGG with ORDER BY)
    # If Yellowbrick doesn't support this syntax, we fall back to Python processing
    _skip_python_combo_processing = False
    df_vmr_combos_agg = None
    
    try:
        curs.execute(f'''
        
            DROP table if exists VMR_{brand_nm}_combos_agg_{analyst};
            CREATE table VMR_{brand_nm}_combos_agg_{analyst} as
            
            WITH combo_strings AS (
                SELECT 
                    ord_event_key,
                    STRING_AGG(
                        CASE 
                            WHEN segm_nbr <= 26 THEN CHR(64 + segm_nbr::int)
                            ELSE CHR(70 + segm_nbr::int)
                        END,
                        '' ORDER BY segm_nbr
                    ) as brand_combo,
                    SUM(units) as units,
                    SUM(dollars) as dollars
                FROM VMR_{brand_nm}_combos_results_{analyst}
                GROUP BY ord_event_key
            )
            SELECT 
                brand_combo,
                COUNT(*) as trips,
                SUM(units) as units,
                SUM(dollars) as dollars
            FROM combo_strings
            GROUP BY brand_combo
            
            DISTRIBUTE RANDOM;
        
        ''')
        
        from local_modules.yb_unload import yb_unload
        
        # Try to use the pre-aggregated table
        df_vmr_combos_agg = yb_unload(yb_user, readpw("Yellowbrick"), table_name = f'''vmr_{brand_nm}_combos_agg_{analyst}''', database = f'{dbase}')
        
        if df_vmr_combos_agg is not None and len(df_vmr_combos_agg) > 0:
            print(f"Using SQL-aggregated brand combinations: {len(df_vmr_combos_agg)} unique combos")
            total_transactions = dict(zip(df_vmr_combos_agg['brand_combo'], df_vmr_combos_agg['trips']))
            total_units = dict(zip(df_vmr_combos_agg['brand_combo'], df_vmr_combos_agg['units']))
            total_dollars = dict(zip(df_vmr_combos_agg['brand_combo'], df_vmr_combos_agg['dollars']))
            _skip_python_combo_processing = True
            
    except Exception as e:
        print(f"SQL aggregation not supported (STRING_AGG ORDER BY): {e}")
        print("Falling back to Python row-by-row processing")
        _skip_python_combo_processing = False
    
    # Fall back to downloading individual rows (original approach)
    if not _skip_python_combo_processing:
        from local_modules.yb_unload import yb_unload
        print("Using Python row-by-row processing for brand combinations")
        df_vmr_period_purch = yb_unload(yb_user, readpw("Yellowbrick"), table_name = f'''vmr_{brand_nm}_combos_results_{analyst}''', database = f'{dbase}')
    
    # Only process row-by-row if SQL aggregation was NOT used
    if not _skip_python_combo_processing:
        df_vmr_period_purch.sort_values(['ord_event_key','segm_nbr']).reset_index().drop(columns=['index'])
        
        
        df_vmr_period_purch[['segm_nbr','segment']] = df_vmr_period_purch[['segm_nbr','segment']].astype("category")
        
        #sql_to_pretty_names_dict(df_vmr_period_purch)                           # update the renaming dictionary to include these column names
        df_vmr_period_purch = df_vmr_period_purch.rename(columns=renaming_dict) # rename the columns, over-writing the original DataFrame column names
        df_vmr_period_purch = df_vmr_period_purch.sort_values(['ord_event_key','Segm Nbr']).reset_index().drop(columns=['index'])



# In[46]:


if nbr_segments > 1 and not _skip_python_combo_processing:
    
    #### This creates 3 dictionaries, to accumulate the number of tranactions, dollars, and units for each brand combination

    current_ord_event_key = 0

    current_brands = ""   # empty string
    current_amount = 0.0  # no dollars accumulated yet
    current_units = 0.0

    total_transactions = {}
    total_dollars = {}
    total_units = {}
    
    idx = {name: i for i, name in enumerate(list(df_vmr_period_purch),1)}
    
    #{'ord_event_key': 1, 'Brand Nbr': 2, 'Brand': 3, 'Units': 4, 'Dollars': 5}
    
    for row in df_vmr_period_purch.itertuples():
        
    #Pandas(Index=0, ord_event_key=509252267271, _2=3, Brand='Lip', Units=1, Dollars=1.83)
        
        #if we are about to start a new transaction, add the details of the old one to our accumulation
        
        if current_ord_event_key != row[idx['ord_event_key']] and current_ord_event_key != 0:
            if not current_brands in total_transactions:
                total_transactions[current_brands] = 1
                total_dollars[current_brands] = current_amount
                total_units[current_brands] = current_units
            else:
                total_transactions[current_brands] += 1
                total_dollars[current_brands] += current_amount
                total_units[current_brands] += current_units

            current_brands = ""   # reset this to an empty string
            current_amount = 0.0  # reset this to no dollars accumulated yet
            current_units = 0.0   # reset this to no units accumulated yet

        brand_letter = get_letter_for_brand_nbr(row[idx['Segm Nbr']])
        current_brands += brand_letter
        current_amount += row[idx['Dollars']]
        current_units  += row[idx['Units']]
        current_ord_event_key = row[idx['ord_event_key']]    

    if not current_brands in total_transactions:
        total_transactions[current_brands] = 1
        total_dollars[current_brands] = current_amount
        total_units[current_brands] = current_units
    else:
        total_transactions[current_brands] += 1
        total_dollars[current_brands] += current_amount
        total_units[current_brands] += current_units

# Create df_vmr_period_combos from the dictionaries (regardless of SQL or Python path)
if nbr_segments > 1:
####  Put the dictionaries into a dataframe, sorted by transactions    

    df_vmr_period_combos = pd.DataFrame(data=[total_transactions,total_units,total_dollars],index=['Trips','Units','Dollars']).T
    df_vmr_period_combos=df_vmr_period_combos.sort_values(['Trips'],ascending = False).reset_index()
    df_vmr_period_combos.rename(columns = {'index':'Segment'},inplace=True)
    df_vmr_period_combos['No. of Segments'] = df_vmr_period_combos['Segment'].str.len()
    df_vmr_period_combos['Pct Of Trips'] = round(df_vmr_period_combos['Trips'] / df_vmr_period_combos['Trips'].sum(),5)
    df_vmr_period_combos['Pct Of Units'] = round(df_vmr_period_combos['Units'] / df_vmr_period_combos['Units'].sum(),5)
    df_vmr_period_combos['Pct Of Dollars'] = round(df_vmr_period_combos['Dollars'] / df_vmr_period_combos['Dollars'].sum(),5)
    df_vmr_period_combos['Units Per Trip'] = round(df_vmr_period_combos['Units'] / df_vmr_period_combos['Trips'],1)
    df_vmr_period_combos['Dollars Per Trip'] = round(df_vmr_period_combos['Dollars'] / df_vmr_period_combos['Trips'],2)

    #sql_to_pretty_names_dict(df_vmr_period_combos)                           # update the renaming dictionary to include these column names
    df_vmr_period_combos = df_vmr_period_combos.rename(columns=renaming_dict) # rename the columns, over-writing the original DataFrame column names
    
####  Change the A-Z,a-z code back to the brand description in the UPC hierarchy (previously downloaded from YB)    

    for i in range(len(df_vmr_period_combos)):
        curr_letters = df_vmr_period_combos.at[i,'Segment']
        curr_brands = ""
        for letter in curr_letters: 
            curr_brands += alpha_to_brand_desc[letter] + ' + '
        curr_brands = curr_brands[:-3]                            #remove the last 3 characters of the string, to delete the final ' + ' out of the brand combination
        df_vmr_period_combos.at[i,'Segment'] = curr_brands

    
####  Total by sub-brand on the VMR reward earning trip
    
    df_vmr_period_totals = pd.read_sql(f''' 
        SELECT a.segm_nbr, a.segment
              ,count(distinct a.ord_event_key) as trips
              ,sum(a.brand_units) as units
              ,sum(a.brand_dollars) as dollars
        FROM VMR_{brand_nm}_purch_details_all_trips a
        INNER JOIN VMR_{brand_nm}_reward_events rd ON (rd.ord_event_key = a.ord_event_key)
        INNER JOIN consumer_id_v p ON (a.cnsmr_id_key = p.cnsmr_id_key)
        WHERE p.abuser_ind= 'N' 
                and p.unident_ind = 'N'
                and p.really_bad_abuser_ind='N'
                and p.cnsmr_id_key>0
                and a.analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)'
        GROUP BY 1,2
        
        UNION ALL
        
        SELECT 999 as segm_nbr, 'Total' as segment
              ,count(distinct a.ord_event_key) as trips
              ,sum(a.brand_units) as units
              ,sum(a.brand_dollars) as dollars
        FROM VMR_{brand_nm}_purch_details_all_trips a
        INNER JOIN VMR_{brand_nm}_reward_events rd ON (rd.ord_event_key = a.ord_event_key)
        INNER JOIN consumer_id_v p ON (a.cnsmr_id_key = p.cnsmr_id_key)
        WHERE p.abuser_ind= 'N' 
                and p.unident_ind = 'N'
                and p.really_bad_abuser_ind='N'
                and p.cnsmr_id_key>0
                and a.analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)'
        GROUP BY 1,2
    ''',conn).sort_values(['segm_nbr']).reset_index().drop(columns=['index','segm_nbr'])
    

    ttl_trips = df_vmr_period_totals[df_vmr_period_totals['segment'] == 'Total']['trips'].item()
    ttl_units = df_vmr_period_totals[df_vmr_period_totals['segment'] == 'Total']['units'].item()
    ttl_dollars = df_vmr_period_totals[df_vmr_period_totals['segment'] == 'Total']['dollars'].item()
    
    
    # Calculate the percentage of trips based on total trips
    df_vmr_period_totals['pct_of_trips'] = df_vmr_period_totals['trips'] / ttl_trips
    df_vmr_period_totals['pct_of_trips'] = df_vmr_period_totals['pct_of_trips'].round(3)

    # Calculate the percentage of units based on total units
    df_vmr_period_totals['pct_of_units'] = df_vmr_period_totals['units'] / ttl_units
    df_vmr_period_totals['pct_of_units'] = df_vmr_period_totals['pct_of_units'].round(3)

    # Calculate the percentage of dollars based on total dollars
    df_vmr_period_totals['pct_of_dollars'] = df_vmr_period_totals['dollars'] / ttl_dollars
    df_vmr_period_totals['pct_of_dollars'] = df_vmr_period_totals['pct_of_dollars'].round(3)
    
    print("Lets see man")
    print(df_vmr_period_totals)

    #sql_to_pretty_names_dict(df_vmr_period_totals)                           # update the renaming dictionary to include these column names
    df_vmr_period_totals = df_vmr_period_totals.rename(columns=renaming_dict) # rename the columns, over-writing the original DataFrame column names
    
    ttl_row_num = df_vmr_period_totals.index[df_vmr_period_totals['Segment'] == 'Total'].tolist()
    df_vmr_period_totals.iloc[ttl_row_num,4:7] = [np.nan, np.nan, np.nan]
    
    
    df_vmr_period_combos = df_vmr_period_combos.sort_values('Pct Of Trips', ascending = False)

    
    if all_combos == False:
    
        rows = df_vmr_period_combos['Pct Of Trips'].shape[0]
        ao_rows = df_vmr_period_combos.loc[10:rows-1].sum()

    df_vmr_period_combos_all = df_vmr_period_combos
    df_vmr_period_combos_all_out = df_vmr_period_combos_all[df_vmr_period_combos_all['Pct Of Trips']>0.005]
    
    if all_combos == False:
    
        df_vmr_period_combos = df_vmr_period_combos.head(10)
        df_vmr_period_combos.loc[11] = ao_rows
        df_vmr_period_combos.at[11,'Segment']='AO combinations'

    df_vmr_period_combos = df_vmr_period_combos.rename(columns={"Segment": "Segment", "No. Of Segments": "No. of Segments" })
    df_vmr_period_combos_all = df_vmr_period_combos_all.rename(columns={"Segment": "Segment", "No. Of Segments": "No. of Segments" })
    df_vmr_period_combos_all_out = df_vmr_period_combos_all_out.rename(columns={"Segment": "Segment", "No. Of Segments": "No. of Segments" })
    
    df_vmr_period_combos = df_vmr_period_combos.reset_index(drop=True)
    df_vmr_period_combos_all = df_vmr_period_combos_all.reset_index(drop=True)
    df_vmr_period_combos_all_out = df_vmr_period_combos_all_out.reset_index(drop=True)
    df_vmr_period_totals = df_vmr_period_totals.reset_index(drop=True)
    
    combos_slide = df_vmr_period_combos[['Segment', 'Pct Of Trips']]
    combos_slide = combos_slide[combos_slide['Segment']!='AO combinations']
    
    print(df_vmr_period_combos)
    print(df_vmr_period_totals)
    
else :
    
    df_vmr_period_combos = []
    df_vmr_period_totals = []
    
    print('Only one segment / no sub-brands in promoted brand group.  Skipping Combinations Analysis.')


# ## **9) SEGMENT ANALYSIS OF PURCHASES VMR REWARD TRIP (FOR CONSISTENT REWARD ID'S):**

# In[47]:


#Modified from K.Mertens's J&J custom VMR code.

if nbr_segments > 1 :

####  Download purchase details from YB, organized by cnsmr_id_key

    curs.execute(f'''
        
        DROP table if exists VMR_{brand_nm}_segments_results_{analyst}; 
        CREATE table VMR_{brand_nm}_segments_results_{analyst} as
        
        SELECT a.cnsmr_id_key,
                CASE WHEN a.analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' and rd.ord_event_key is not null THEN 'VMR Period'
                     WHEN a.prior_period = 'Prior-period ({round(pre_weeks)} weeks)' THEN 'Prior-Period'
                     ELSE 'NA' END as analysis_periods, 
                a.segm_nbr,
                sum(a.brand_units) as units 
        FROM VMR_{brand_nm}_purch_details_all_trips a
            INNER JOIN VMR_{brand_nm}_reward_ids_csmr d ON (a.cnsmr_id_key = d.ord_designated_cnsmr_id_key)
            INNER JOIN shopper_consistent_{analyst} s ON (a.cnsmr_id_key = s.cnsmr_id_key)
            LEFT JOIN VMR_{brand_nm}_reward_events rd ON (rd.ord_event_key = a.ord_event_key)
        WHERE analysis_periods != 'NA'
        GROUP by 1,2,3
        
        DISTRIBUTE RANDOM;
        
    ''')
    
    
    from local_modules.yb_unload import yb_unload

    df_pre_vmr_cnsmrs = yb_unload(yb_user, readpw("Yellowbrick"), table_name = f'''vmr_{brand_nm}_segments_results_{analyst}''', database = f'{dbase}')
        
    df_pre_vmr_cnsmrs.sort_values(['analysis_periods','cnsmr_id_key','segm_nbr'],ascending = True).reset_index().drop(columns=['index'])
    
    
    df_pre_vmr_cnsmrs[['analysis_periods','segm_nbr']] = df_pre_vmr_cnsmrs[['analysis_periods','segm_nbr']].astype('category')
    df_pre_vmr_cnsmrs[['cnsmr_id_key']] = df_pre_vmr_cnsmrs[['cnsmr_id_key']].astype('str')

    #sql_to_pretty_names_dict(df_pre_vmr_cnsmrs)                           # update the renaming dictionary to include these column names
    df_pre_vmr_cnsmrs = df_pre_vmr_cnsmrs.rename(columns=renaming_dict)   # rename the columns, over-writing the original DataFrame column names
    df_pre_vmr_cnsmrs.sort_values(['Analysis Period', 'Segm Nbr'],ascending = True).reset_index().drop(columns=['index'])
    

#### This creates 2 dictionaries, to look at sub-brand purchasing before and at reward earning trip

    idx = {name: i for i, name in enumerate(list(df_pre_vmr_cnsmrs),1)}

    # Two dicts, accessible by analysis period
    consumer_behavior = {}
    consumer_behavior['Prior-Period'] = {}
    consumer_behavior['VMR Period'] = {}

    # Accumulate consumer behavhior before and during print
    for row in df_pre_vmr_cnsmrs.itertuples():
        # unpack the details from the row
        current_period = row[idx['Analysis Period']]
        current_consumer = row[idx['cnsmr_id_key']]
        brand_letter = get_letter_for_brand_nbr(row[idx['Segm Nbr']])

        # Assign consumer's brand purchase to the current period
        if not current_consumer in consumer_behavior[current_period]:
            consumer_behavior[current_period][current_consumer] = brand_letter
        else:
            consumer_behavior[current_period][current_consumer] += brand_letter

    N_prev_seg = 0
    N_new_seg = 0
    N_prev_new_seg = 0
    
    #creating empty lists where the consumer ids for each segment will be saved: (jmorice)
    ids_prev_seg = []
    ids_prev_new_seg = []
    ids_new_seg = []
    
    
    # Note: all consumers in pre-period are also in print period -- but sometimes, when we filter for purch_qty>0 and purch_amt>0.01, then the print period trip drops out.
    # This forces analysis to look at only existing brand buyers: the shoppers with purchases in both the pre-period and analysis period
    for consumer in consumer_behavior['Prior-Period']:
        # XXX skip if they are not in print period
        if not consumer in consumer_behavior['VMR Period']: continue

        # Construct sets for pre- and print-periods
        # Note: consider constructing these sets in the previous loop, rather than doing string stuff
        pre_set = set(list(consumer_behavior['Prior-Period'][consumer]))
        print_set = set(list(consumer_behavior['VMR Period'][consumer]))

        # If print set is entirely in previous, then this is "only previous" behavior
        if print_set.issubset(pre_set):
            N_prev_seg += 1
            ids_prev_seg.append(consumer)

        # no overlap is called "only new"
        elif print_set.isdisjoint(pre_set):
            N_new_seg += 1
            ids_new_seg.append(consumer)

        # must be combination
        else:
            N_prev_new_seg += 1
            ids_prev_new_seg.append(consumer)

            
    N_prev_seg_bump_up = (N_prev_seg)
    N_new_seg_bump_up = (N_new_seg)
    N_prev_new_seg_bump_up = (N_prev_new_seg)
    
    
    #Transforming the segment consumer ids's lists into dataframes:
    ids_prev_seg = pd.DataFrame(ids_prev_seg).rename(columns={0:'cnsmr_id_key'})
    ids_prev_new_seg = pd.DataFrame(ids_prev_new_seg).rename(columns={0:'cnsmr_id_key'})
    ids_new_seg = pd.DataFrame(ids_new_seg).rename(columns={0:'cnsmr_id_key'})

    
#### Put result into a dataframe and condense
    df_plus_one = pd.DataFrame(data=[N_prev_seg_bump_up,N_new_seg_bump_up + N_prev_new_seg_bump_up],index=['Bought the same segments as they did before','Added at least one new segment'])
    df_plus_one.rename(columns={df_plus_one.columns[0]: 'Nbr Of Shoppers'},inplace=True)
    df_plus_one['Pct Of Shoppers'] = round(df_plus_one['Nbr Of Shoppers'] / df_plus_one['Nbr Of Shoppers'].sum(),3)
    df_plus_one = df_plus_one.reset_index().rename(columns={'index': ''})
    
    #sql_to_pretty_names_dict(df_plus_one)                           # update the renaming dictionary to include these column names
    
#### Same result into another dataframe with more detail
    df_new_existing_combos = pd.DataFrame(data=[N_prev_seg_bump_up,N_prev_new_seg_bump_up,N_new_seg_bump_up],index=['Bought the same segments as they did before','Both previous and new segment(s) on this trip','Only the new segment(s) on this trip'])
    df_new_existing_combos.rename(columns={df_new_existing_combos.columns[0]: 'Nbr Of Shoppers'},inplace=True)
    df_new_existing_combos['Pct Of Shoppers'] = round(df_new_existing_combos['Nbr Of Shoppers'] / df_new_existing_combos['Nbr Of Shoppers'].sum(),3)
    df_new_existing_combos = df_new_existing_combos.reset_index().rename(columns={'index': ''})

    #sql_to_pretty_names_dict(df_new_existing_combos)                           # update the renaming dictionary to include these column names

else :
    
    df_plus_one = []
    df_new_existing_combos = []   
    
    print('Only one segment / no sub-brands in promoted brand group.  Skipping Combinations Analysis.')    


# In[48]:


if nbr_segments > 1 :
     print(df_plus_one)

if nbr_segments > 1 :
    print(df_new_existing_combos)


# ### Decomp by segment for Shoppers who bought a new segment during the reward trip:

# In[49]:


if nbr_segments > 1 :

    ids_new_seg_2 = ids_new_seg
    ids_prev_new_seg_2 = ids_prev_new_seg

    all_ids_new_seg = pd.concat([ids_new_seg, ids_prev_new_seg])
    all_ids_new_seg = all_ids_new_seg.reset_index(drop=True)

    yb_load(Df = all_ids_new_seg,
            table_name = f'''VMR_{brand_nm}_all_ids_new_seg_{analyst}''',
            userid = yb_user,
            passwd = readpw("Yellowbrick"),
            append = False,
            database= f'{dbase}')


# In[50]:


if nbr_segments > 1 :
    
    curs.execute(f'''       
        DROP table if exists VMR_{brand_nm}_new_buyer_flag;
        CREATE temp table VMR_{brand_nm}_new_buyer_flag as

           SELECT a.cnsmr_id_key,
                   a.segm_nbr,
                   a.segment,
                   count(distinct CASE WHEN a.analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' and rd.ord_event_key is not null THEN a.ord_event_key ELSE null END) AS trips_vmr_period,
                   count(distinct CASE WHEN a.prior_period = 'Prior-period ({round(pre_weeks)} weeks)' THEN a.ord_event_key ELSE null END) AS trips_prior_period

                   FROM VMR_{brand_nm}_purch_details_all_trips a
                        INNER JOIN VMR_{brand_nm}_all_ids_new_seg_{analyst} n ON (a.cnsmr_id_key = n.cnsmridkey)
                       LEFT JOIN VMR_{brand_nm}_reward_events rd ON (rd.ord_event_key = a.ord_event_key)

                    GROUP BY 1,2,3
        DISTRIBUTE REPLICATE;

    ''')


    curs.execute(f'''       
        DROP table if exists VMR_{brand_nm}_new_seg_denom;
        CREATE temp table VMR_{brand_nm}_new_seg_denom as

           SELECT count(distinct cnsmr_id_key) as buyer_denom
           FROM VMR_{brand_nm}_new_buyer_flag 

        DISTRIBUTE REPLICATE;
    ''')


    curs.execute(f'''       
        DROP table if exists VMR_{brand_nm}_new_seg_numer;
        CREATE temp table VMR_{brand_nm}_new_seg_numer as

           SELECT segm_nbr, segment,
                  count(distinct cnsmr_id_key) as buyer_numer
           FROM VMR_{brand_nm}_new_buyer_flag
           where trips_prior_period=0 and trips_vmr_period>0
           GROUP BY 1,2
        DISTRIBUTE REPLICATE;
    ''')    

    new_seg_pct = pd.read_sql(f'''select a.segm_nbr, segment, buyer_numer as new_seg_buyer, 
                                  buyer_denom as ttl_new_buyer, new_seg_buyer::float / ttl_new_buyer::float as new_seg_pct
    from VMR_{brand_nm}_new_seg_numer a
    cross join VMR_{brand_nm}_new_seg_denom b
    order by 1''',conn)    

    print(new_seg_pct)


# In[51]:


# ## **11) PRE-52 WEEKS PROFILE (FOR CONSISTENT REWARD ID'S BASED ON DOLLARS):**

# In[52]:



#Getting the consistent ID's who purchased the promoted brand during the recent 52 weeks before the reward print period.
#Getting #trips and dollars for both brand and categories

curs.execute(f'''

    DROP table if exists VMR_{brand_nm}_buyers_Pre52;
    CREATE temp table VMR_{brand_nm}_buyers_Pre52 as

        SELECT distinct o.ord_designated_cnsmr_id_key,
               count(distinct CASE WHEN c.brand_nbr in ({brand_nbr_str}) THEN o.ord_event_key else null end) as b_trips,
               count(distinct CASE WHEN c.brand_nbr in ({cat_nbr_str}) THEN o.ord_event_key else null end) as c_trips,
               sum(CASE WHEN c.brand_nbr in ({brand_nbr_str}) THEN o.purch_amt else 0 end) as bdollars,
               sum(CASE WHEN c.brand_nbr in ({cat_nbr_str}) THEN o.purch_amt else 0 end) as cdollars
               
         FROM ORD_TRD_ITM_CNSMR_FACT_NE_V o
                
                INNER JOIN shopper_consistent_{analyst} s ON (o.ord_designated_cnsmr_id_key = s.cnsmr_id_key)
                INNER JOIN VMR_{brand_nm}_upc_filter c ON (o.trade_item_key = c.trade_item_key)
                INNER JOIN date_v b ON (o.ord_date_key = b.date_key)
         
         WHERE b.cal_dt BETWEEN '{prior_period_start}' and '{prior_period_end}' 
               and purch_amt > 0
               and purch_qty > 0
               
         GROUP BY 1
         
                 
         DISTRIBUTE RANDOM; 
              
''')


# In[53]:


pre_52 = pd.read_sql(f''' SELECT * FROM VMR_{brand_nm}_buyers_Pre52 limit 10''',conn)
pre_52


# In[54]:


nbr_buyers_pre52 = pd.read_sql(f''' 
SELECT count(ord_designated_cnsmr_id_key)
FROM VMR_{brand_nm}_buyers_Pre52
WHERE b_trips > 0
''',conn)
print("Number of total promoted brand buyers during Pre-52 period: " + str(nbr_buyers_pre52['count'].values[0]))


# #### For Brand Comsuption:

# In[55]:


brand_dol50 = pd.read_sql(f''' SELECT percentile_cont(0.50) within group(order by bdollars) as pct FROM VMR_{brand_nm}_buyers_Pre52 where b_trips>1''',conn)
brand_dol50 = brand_dol50['pct'].values[0]
brand_dol75 = pd.read_sql(f''' SELECT percentile_cont(0.75) within group(order by bdollars) as pct FROM VMR_{brand_nm}_buyers_Pre52 where b_trips>1''',conn)
brand_dol75 = brand_dol75['pct'].values[0]
print("Brand Pct50: " + str(brand_dol50))
print("Brand Pct75: " + str(brand_dol75))


# In[56]:


brand_consump = pd.read_sql(f''' 

SELECT CASE 
       WHEN b_trips=0 or  b.ord_designated_cnsmr_id_key is null THEN 'Never Buyer'
       WHEN b_trips=1 AND b.ord_designated_cnsmr_id_key is not null THEN '1x Buyer'
       WHEN b_trips>1 AND bdollars>0 AND bdollars<{brand_dol50} AND b.ord_designated_cnsmr_id_key is not null THEN 'Light Buyer (1%-50%)'
       WHEN b_trips>1 AND bdollars>={brand_dol50} AND bdollars<{brand_dol75} AND b.ord_designated_cnsmr_id_key is not null THEN 'Medium Buyer (50%-75%)'
       WHEN b_trips>1 AND bdollars>={brand_dol75} AND b.ord_designated_cnsmr_id_key is not null THEN 'Heavy Buyer (75%-100%)' END as brand_grp,
       COUNT(distinct a.ord_designated_cnsmr_id_key) as brand_ids
       
FROM VMR_{brand_nm}_reward_ids_csmr a
    LEFT JOIN VMR_{brand_nm}_buyers_Pre52 b ON (b.ord_designated_cnsmr_id_key = a.ord_designated_cnsmr_id_key)
    INNER JOIN shopper_consistent_{analyst} s ON (a.ord_designated_cnsmr_id_key = s.cnsmr_id_key)
WHERE brand_grp is not null
GROUP BY 1
ORDER BY CASE WHEN brand_grp = 'Never Buyer' then 1
              WHEN brand_grp = '1x Buyer' then 2
              WHEN brand_grp = 'Light Buyer (1%-50%)' then 3
              WHEN brand_grp = 'Medium Buyer (50%-75%)' then 4
              WHEN brand_grp = 'Heavy Buyer (75%-100%)' then 5
              END

''',conn)

brand_consump['%_brand_ids'] = (brand_consump['brand_ids']/brand_consump['brand_ids'].sum())

brand_consump.loc[5] = brand_consump.sum()
brand_consump.at[5, 'brand_grp'] = 'Total'

brand_consump = brand_consump.rename(columns = {'brand_grp':'Brand Group', 'brand_ids':'Brand IDs', '%_brand_ids':'% Brand IDs'})

new_brand_buyers = "{:,}".format(brand_consump.loc[0,'Brand IDs'])

brand_consump


# In[57]:


#Saving all the categorized IDs in a table and loading it into YB:

brand_consump_ids = pd.read_sql(f''' 

SELECT a.ord_designated_cnsmr_id_key,
       CASE 
       WHEN b_trips=0 or  b.ord_designated_cnsmr_id_key is null THEN 'Never Buyer'
       WHEN b_trips=1 AND b.ord_designated_cnsmr_id_key is not null THEN '1x Buyer'
       WHEN b_trips>1 AND bdollars>0 AND bdollars<{brand_dol50} AND b.ord_designated_cnsmr_id_key is not null THEN 'Light'
       WHEN b_trips>1 AND bdollars>={brand_dol50} AND bdollars<{brand_dol75} AND b.ord_designated_cnsmr_id_key is not null THEN 'Medium'
       WHEN b_trips>1 AND bdollars>={brand_dol75} AND b.ord_designated_cnsmr_id_key is not null THEN 'Heavy' END as brand_grp
       
FROM VMR_{brand_nm}_reward_ids_csmr a
    LEFT JOIN VMR_{brand_nm}_buyers_Pre52 b ON (b.ord_designated_cnsmr_id_key = a.ord_designated_cnsmr_id_key)
WHERE brand_grp is not null
GROUP BY 1,2
ORDER BY CASE WHEN brand_grp = 'Never Buyer' then 1
              WHEN brand_grp = '1x Buyer' then 2
              WHEN brand_grp = 'Light' then 3
              WHEN brand_grp = 'Medium' then 4
              WHEN brand_grp = 'Heavy' then 5
              END

''',conn)

brand_consump_ids


# In[58]:


yb_load(Df = brand_consump_ids,
            table_name = f'''VMR_{brand_nm}_brand_buyers_ids_{analyst}''',
            userid = yb_user,
            passwd = readpw("Yellowbrick"),
            append = False,
            database= f'{dbase}')

check_upload = pd.read_sql(f"""select count(*) from VMR_{brand_nm}_brand_buyers_ids_{analyst}""", conn)
print('Uploaded number of records: ' + str(check_upload.values[0]))


# #### For Category Comsuption:

# In[59]:


cat_dol50 = pd.read_sql(f''' SELECT percentile_cont(0.50) within group(order by cdollars) as pct FROM VMR_{brand_nm}_buyers_Pre52 WHERE c_trips>1 ''',conn)
cat_dol50 = cat_dol50['pct'].values[0]
cat_dol75 = pd.read_sql(f''' SELECT percentile_cont(0.75) within group(order by cdollars) as pct FROM VMR_{brand_nm}_buyers_Pre52 WHERE c_trips>1''',conn)
cat_dol75 = cat_dol75['pct'].values[0]
print("Cat Pct50: " + str(cat_dol50))
print("Cat Pct75: " + str(cat_dol75))


# In[60]:


cat_consump = pd.read_sql(f''' 

SELECT CASE 
       WHEN b.ord_designated_cnsmr_id_key is null THEN 'Never Buyer'
       WHEN c_trips=1 AND b.ord_designated_cnsmr_id_key is not null THEN '1x Buyer'
       WHEN c_trips>1 AND cdollars>0 AND cdollars<{cat_dol50} AND b.ord_designated_cnsmr_id_key is not null THEN 'Light Buyer (1%-50%)'
       WHEN c_trips>1 AND cdollars>={cat_dol50} AND cdollars<{cat_dol75} AND b.ord_designated_cnsmr_id_key is not null THEN 'Medium Buyer (50%-75%)'
       WHEN c_trips>1 AND cdollars>={cat_dol75} AND b.ord_designated_cnsmr_id_key is not null THEN 'Heavy Buyer (75%-100%)' END as category_grp,
       COUNT(distinct a.ord_designated_cnsmr_id_key) as category_ids
       
FROM VMR_{brand_nm}_reward_ids_csmr a
    LEFT JOIN VMR_{brand_nm}_buyers_Pre52 b ON (b.ord_designated_cnsmr_id_key = a.ord_designated_cnsmr_id_key)
    INNER JOIN shopper_consistent_{analyst} s ON (a.ord_designated_cnsmr_id_key = s.cnsmr_id_key)
WHERE category_grp is not null
GROUP BY 1
ORDER BY CASE WHEN category_grp = 'Never Buyer' then 1
              WHEN category_grp = '1x Buyer' then 2
              WHEN category_grp = 'Light Buyer (1%-50%)' then 3
              WHEN category_grp = 'Medium Buyer (50%-75%)' then 4
              WHEN category_grp = 'Heavy Buyer (75%-100%)' then 5
              END
 

''',conn)

cat_consump['%_category_ids'] = (cat_consump['category_ids']/cat_consump['category_ids'].sum())

cat_consump.loc[5] = cat_consump.sum()
cat_consump.at[5, 'category_grp'] = 'Total'

cat_consump = cat_consump.rename(columns = {'category_grp':'Category Group', 'category_ids':'Category IDs', '%_category_ids':'% Category IDs'})

cat_consump


# ## **12) GETTING REPURCHASE RATE FOR VMR PARTICIPANTS:**

# In[61]:


curs.execute(f'''

    DROP table if exists VMR_{brand_nm}_repurch_details;
    CREATE temp table VMR_{brand_nm}_repurch_details as
        
          SELECT d.ord_designated_cnsmr_id_key,
                 e.brandgrp as brand_group,
                 min(d.reward_date) as reward_date
          
          FROM ord_trd_itm_cnsmr_fact_ne_v o
                 
                 INNER JOIN VMR_{brand_nm}_date_filter b ON (o.ord_date_key = b.date_key)
                 INNER JOIN VMR_{brand_nm}_upc_filter c ON (o.trade_item_key = c.trade_item_key)
                 INNER JOIN VMR_{brand_nm}_reward_ids d ON (o.ord_designated_cnsmr_id_key = d.ord_designated_cnsmr_id_key)
                 INNER JOIN VMR_{brand_nm}_brand_buyers_ids_{analyst} e ON (o.ord_designated_cnsmr_id_key = e.orddesignatedcnsmridkey)
                INNER JOIN shopper_consistent_{analyst} s ON (o.ord_designated_cnsmr_id_key = s.cnsmr_id_key)                 
          
          WHERE o.purch_amt > 0
                and o.purch_qty > 0
                and c.brand_nbr in ({brand_nbr_str})
                and b.cal_dt > reward_date               --Repurch date after the printing date
                and b.cal_dt between '{reward_print_start}' and '{post_4wk_end}'  --Repurch including N weeks after the reward purchase
        
        GROUP BY 1,2
        
        DISTRIBUTE ON (ord_designated_cnsmr_id_key);
        
''')


# In[62]:


pd.read_sql(f'''
    SELECT *
    FROM VMR_{brand_nm}_repurch_details
    ''',conn)


# In[63]:


repurch_total_buyers = pd.read_sql(f'''
    SELECT count(distinct ord_designated_cnsmr_id_key)
    FROM VMR_{brand_nm}_repurch_details
    ''',conn)



repurch_total_buyers = repurch_total_buyers['count'].values[0]
print("The total of VMR participants that repurchased during the post 4 weeks is: " + str(repurch_total_buyers))

total_reward_ids_m = pd.read_sql(f'''
    SELECT count(distinct orddesignatedcnsmridkey)
    FROM VMR_{brand_nm}_brand_buyers_ids_{analyst} a
    INNER JOIN shopper_consistent_{analyst} s ON (a.orddesignatedcnsmridkey = s.cnsmr_id_key)
    ''',conn)

pct_repurch = (repurch_total_buyers/total_reward_ids_m['count'].values[0])
print("The repurchase percentage among all the VMR participants is: " + str(round(pct_repurch*100)))


# In[64]:


repurch_new_buyers = pd.read_sql(f'''
    SELECT count(distinct ord_designated_cnsmr_id_key)
    FROM VMR_{brand_nm}_repurch_details
    WHERE brand_group = 'Never Buyer'
    ''',conn)

total_new_buy = pd.read_sql(f'''
    SELECT count(distinct orddesignatedcnsmridkey)
    FROM VMR_{brand_nm}_brand_buyers_ids_{analyst} a
    INNER JOIN shopper_consistent_{analyst} s ON (a.orddesignatedcnsmridkey = s.cnsmr_id_key)
    WHERE brandgrp = 'Never Buyer'
    ''',conn)

repurch_new_buyers = repurch_new_buyers['count'].values[0]
print("The total of VMR participants that repurchased during the post 4 weeks is: " + str(repurch_new_buyers))

pct_repurch_new_buyers = repurch_new_buyers/total_new_buy['count'].values[0]
print("The repurchase percentage among all the VMR participants is: " + str(round(pct_repurch_new_buyers*100)))


# In[65]:


repurch_existing_buyers = pd.read_sql(f'''
    SELECT count(distinct ord_designated_cnsmr_id_key)
    FROM VMR_{brand_nm}_repurch_details
    WHERE brand_group != 'Never Buyer'
    ''',conn)


total_exist_buy = pd.read_sql(f'''
    SELECT count(distinct orddesignatedcnsmridkey)
    FROM VMR_{brand_nm}_brand_buyers_ids_{analyst} a
    INNER JOIN shopper_consistent_{analyst} s ON (a.orddesignatedcnsmridkey = s.cnsmr_id_key)
    WHERE brandgrp != 'Never Buyer'
    ''',conn)

repurch_existing_buyers = repurch_existing_buyers['count'].values[0]
print("The total of existing participants that repurchased during the post 4 weeks is: " + str(repurch_existing_buyers))

pct_repurch_existing_buyers = repurch_existing_buyers/total_exist_buy['count'].values[0]
print("The repurchase percentage among the existing participants is: " + str(round(pct_repurch_existing_buyers*100)))


# ## **13) GETTING DATA FOR CATEGORY SHARE SLIDE:**

# In[66]:


curs.execute(f'''       
    DROP table if exists VMR_{brand_nm}_vmr_category_shade;
    CREATE temp table VMR_{brand_nm}_vmr_category_shade as
    
        SELECT  'VMR Period' as Period,
                SUM(CASE WHEN brand_nbr IN ({brand_nbr_str}) and analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN a.purch_amt ELSE 0 END) as dollars_promoted,
                SUM(CASE WHEN analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN a.purch_amt ELSE 0 END) as dollars_category,
                SUM(CASE WHEN brand_nbr IN ({brand_nbr_str}) and analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN a.purch_qty ELSE 0 END) as units_promoted,
                SUM(CASE WHEN analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN a.purch_qty ELSE 0 END) as units_category
                
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
            INNER JOIN VMR_{brand_nm}_reward_ids f ON (a.ord_designated_cnsmr_id_key = f.ord_designated_cnsmr_id_key)

            
        WHERE analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' and
              a.purch_amt > 0
              and a.purch_qty > 0 and
              cal_dt between '{reward_print_start}' and '{reward_print_end}'
              
        UNION
        
        SELECT  'Pre52wk Period' as Period,
                SUM(CASE WHEN brand_nbr IN ({brand_nbr_str}) and prior_period = 'Prior-period ({round(pre_weeks)} weeks)' THEN a.purch_amt ELSE 0 END) as dollars_promoted,
                SUM(CASE WHEN prior_period = 'Prior-period ({round(pre_weeks)} weeks)' THEN a.purch_amt ELSE 0 END) as dollars_category,
                SUM(CASE WHEN brand_nbr IN ({brand_nbr_str}) and prior_period = 'Prior-period ({round(pre_weeks)} weeks)' THEN a.purch_qty ELSE 0 END) as units_promoted,
                SUM(CASE WHEN prior_period = 'Prior-period ({round(pre_weeks)} weeks)' THEN a.purch_qty ELSE 0 END) as units_category
        
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
            INNER JOIN VMR_{brand_nm}_reward_ids f ON (a.ord_designated_cnsmr_id_key = f.ord_designated_cnsmr_id_key)

            
        WHERE prior_period = 'Prior-period ({round(pre_weeks)} weeks)' and
              a.purch_amt > 0
              and a.purch_qty > 0 and
              cal_dt between '{prior_period_start}' and '{prior_period_end}'
              
              
        UNION
        
        SELECT  'Post4wk Period' as Period,
                SUM(CASE WHEN brand_nbr IN ({brand_nbr_str}) and post_period = 'Post Period ({round(weeks_per_period_final)} weeks)' THEN a.purch_amt ELSE 0 END) as dollars_promoted,
                SUM(CASE WHEN post_period = 'Post Period ({round(weeks_per_period_final)} weeks)' THEN a.purch_amt ELSE 0 END) as dollars_category,
                SUM(CASE WHEN brand_nbr IN ({brand_nbr_str}) and post_period = 'Post Period ({round(weeks_per_period_final)} weeks)' THEN a.purch_qty ELSE 0 END) as units_promoted,
                SUM(CASE WHEN post_period = 'Post Period ({round(weeks_per_period_final)} weeks)' THEN a.purch_qty ELSE 0 END) as units_category
                
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
            INNER JOIN VMR_{brand_nm}_reward_ids f ON (a.ord_designated_cnsmr_id_key = f.ord_designated_cnsmr_id_key)

            
        WHERE post_period = 'Post Period ({round(weeks_per_period_final)} weeks)' and
              a.purch_amt > 0
              and a.purch_qty > 0 and
              cal_dt between '{post_4wk_start}' and '{post_4wk_end}'
              
        UNION
        
        
        SELECT  'YAGO Period' as Period,
                SUM(CASE WHEN brand_nbr IN ({brand_nbr_str}) and analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)' THEN a.purch_amt ELSE 0 END) as dollars_promoted,
                SUM(CASE WHEN analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)' THEN a.purch_amt ELSE 0 END) as dollars_category,
                SUM(CASE WHEN brand_nbr IN ({brand_nbr_str}) and analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)' THEN a.purch_qty ELSE 0 END) as units_promoted,
                SUM(CASE WHEN analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)' THEN a.purch_qty ELSE 0 END) as units_category
                
                
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
            INNER JOIN VMR_{brand_nm}_reward_ids f ON (a.ord_designated_cnsmr_id_key = f.ord_designated_cnsmr_id_key)

            
        WHERE analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)' AND
              a.purch_amt > 0
              and a.purch_qty > 0 and
              cal_dt between '{yago_print_start}' and '{yago_print_end}'
        
    DISTRIBUTE REPLICATE;

''')


# In[67]:


share_category_data = pd.read_sql(f'''
    SELECT *
    FROM VMR_{brand_nm}_vmr_category_shade
    ''',conn)
share_category_data 


# In[68]:

dollars_promoted_yago = share_category_data.loc[share_category_data['period']=='YAGO Period', 'dollars_promoted'].values[0]
dollars_category_yago = share_category_data.loc[share_category_data['period']=='YAGO Period', 'dollars_category'].values[0]
units_promoted_yago = share_category_data.loc[share_category_data['period']=='YAGO Period', 'units_promoted'].values[0]
units_category_yago = share_category_data.loc[share_category_data['period']=='YAGO Period', 'units_category'].values[0]

dollars_promoted_vmr = share_category_data.loc[share_category_data['period']=='VMR Period', 'dollars_promoted'].values[0]
dollars_category_vmr = share_category_data.loc[share_category_data['period']=='VMR Period', 'dollars_category'].values[0]
units_promoted_vmr = share_category_data.loc[share_category_data['period']=='VMR Period', 'units_promoted'].values[0]
units_category_vmr = share_category_data.loc[share_category_data['period']=='VMR Period', 'units_category'].values[0]

dollars_promoted_pre52 = share_category_data.loc[share_category_data['period']=='Pre52wk Period', 'dollars_promoted'].values[0]
dollars_category_pre52 = share_category_data.loc[share_category_data['period']=='Pre52wk Period', 'dollars_category'].values[0]
units_promoted_pre52 = share_category_data.loc[share_category_data['period']=='Pre52wk Period', 'units_promoted'].values[0]
units_category_pre52 = share_category_data.loc[share_category_data['period']=='Pre52wk Period', 'units_category'].values[0]

dollars_promoted_post = share_category_data.loc[share_category_data['period']=='Post4wk Period', 'dollars_promoted'].values[0]
dollars_category_post = share_category_data.loc[share_category_data['period']=='Post4wk Period', 'dollars_category'].values[0]
units_promoted_post = share_category_data.loc[share_category_data['period']=='Post4wk Period', 'units_promoted'].values[0]
units_category_post = share_category_data.loc[share_category_data['period']=='Post4wk Period', 'units_category'].values[0]


share_dollar_vmr = dollars_promoted_vmr/dollars_category_vmr
share_dollar_vmr = 0 if np.isnan(share_dollar_vmr) else share_dollar_vmr
print("share_dollar_vmr: " + str(share_dollar_vmr))
share_dollar_pre52 = dollars_promoted_pre52/dollars_category_pre52
share_dollar_pre52 = 0 if np.isnan(share_dollar_pre52) else share_dollar_pre52
print("share_dollar_pre52: " + str(share_dollar_pre52))
share_dollar_post_4wk = dollars_promoted_post/dollars_category_post
share_dollar_post_4wk = 0 if np.isnan(share_dollar_post_4wk) else share_dollar_post_4wk
print("share_dollar_post_4wk: " + str(share_dollar_post_4wk))
share_dollar_yago = dollars_promoted_yago/dollars_category_yago
share_dollar_yago = 0 if np.isnan(share_dollar_yago) else share_dollar_yago
print("share_dollar_yago: " + str(share_dollar_yago))

share_units_vmr = units_promoted_vmr/units_category_vmr
share_units_vmr = 0 if np.isnan(share_units_vmr) else share_units_vmr
print("share_units_vmr: " + str(share_units_vmr))
share_units_pre52 = units_promoted_pre52/units_category_pre52
share_units_pre52 = 0 if np.isnan(share_units_pre52) else share_units_pre52
print("share_units_pre52: " + str(share_units_pre52))
share_units_post_4wk = units_promoted_post/units_category_post
share_units_post_4wk = 0 if np.isnan(share_units_post_4wk) else share_units_post_4wk
print("share_units_post_4wk: " + str(share_units_post_4wk))
share_units_yago = units_promoted_yago/units_category_yago
share_units_yago = 0 if np.isnan(share_units_yago) else share_units_yago
print("share_units_yago: " + str(share_units_yago))
        

#share metrics for all IDS

curs.execute(f'''       
    DROP table if exists VMR_{brand_nm}_vmr_category_shade_all;
    CREATE temp table VMR_{brand_nm}_vmr_category_shade_all as
    
        SELECT  'VMR Period' as Period,
                SUM(CASE WHEN brand_nbr IN ({brand_nbr_str}) and analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN a.purch_amt ELSE 0 END) as dollars_promoted,
                SUM(CASE WHEN analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN a.purch_amt ELSE 0 END) as dollars_category,
                SUM(CASE WHEN brand_nbr IN ({brand_nbr_str}) and analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN a.purch_qty ELSE 0 END) as units_promoted,
                SUM(CASE WHEN analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN a.purch_qty ELSE 0 END) as units_category
                
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)

            
        WHERE analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' and
              a.purch_amt > 0
              and a.purch_qty > 0 and
              cal_dt between '{reward_print_start}' and '{reward_print_end}'
              
        UNION
        
        SELECT  'Pre52wk Period' as Period,
                SUM(CASE WHEN brand_nbr IN ({brand_nbr_str}) and prior_period = 'Prior-period ({round(pre_weeks)} weeks)' THEN a.purch_amt ELSE 0 END) as dollars_promoted,
                SUM(CASE WHEN prior_period = 'Prior-period ({round(pre_weeks)} weeks)' THEN a.purch_amt ELSE 0 END) as dollars_category,
                SUM(CASE WHEN brand_nbr IN ({brand_nbr_str}) and prior_period = 'Prior-period ({round(pre_weeks)} weeks)' THEN a.purch_qty ELSE 0 END) as units_promoted,
                SUM(CASE WHEN prior_period = 'Prior-period ({round(pre_weeks)} weeks)' THEN a.purch_qty ELSE 0 END) as units_category
        
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)

            
        WHERE prior_period = 'Prior-period ({round(pre_weeks)} weeks)' and
              a.purch_amt > 0
              and a.purch_qty > 0 and
              cal_dt between '{prior_period_start}' and '{prior_period_end}'
              
              
        UNION
        
        SELECT  'Post4wk Period' as Period,
                SUM(CASE WHEN brand_nbr IN ({brand_nbr_str}) and post_period = 'Post Period ({round(weeks_per_period_final)} weeks)' THEN a.purch_amt ELSE 0 END) as dollars_promoted,
                SUM(CASE WHEN post_period = 'Post Period ({round(weeks_per_period_final)} weeks)' THEN a.purch_amt ELSE 0 END) as dollars_category,
                SUM(CASE WHEN brand_nbr IN ({brand_nbr_str}) and post_period = 'Post Period ({round(weeks_per_period_final)} weeks)' THEN a.purch_qty ELSE 0 END) as units_promoted,
                SUM(CASE WHEN post_period = 'Post Period ({round(weeks_per_period_final)} weeks)' THEN a.purch_qty ELSE 0 END) as units_category
                
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)

            
        WHERE post_period = 'Post Period ({round(weeks_per_period_final)} weeks)' and
              a.purch_amt > 0
              and a.purch_qty > 0 and
              cal_dt between '{post_4wk_start}' and '{post_4wk_end}'
              
        UNION
        
        
        SELECT  'YAGO Period' as Period,
                SUM(CASE WHEN brand_nbr IN ({brand_nbr_str}) and analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)' THEN a.purch_amt ELSE 0 END) as dollars_promoted,
                SUM(CASE WHEN analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)' THEN a.purch_amt ELSE 0 END) as dollars_category,
                SUM(CASE WHEN brand_nbr IN ({brand_nbr_str}) and analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)' THEN a.purch_qty ELSE 0 END) as units_promoted,
                SUM(CASE WHEN analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)' THEN a.purch_qty ELSE 0 END) as units_category
                
                
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)

            
        WHERE analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)' AND
              a.purch_amt > 0
              and a.purch_qty > 0 and
              cal_dt between '{yago_print_start}' and '{yago_print_end}'
        
    DISTRIBUTE REPLICATE;
''')


# In[67]:


share_category_data_all = pd.read_sql(f'''
    SELECT *
    FROM VMR_{brand_nm}_vmr_category_shade_all
    ''',conn)
share_category_data_all 


dollars_promoted_vmr_all = share_category_data_all.loc[share_category_data_all['period']=='VMR Period', 'dollars_promoted'].values[0]
dollars_category_vmr_all = share_category_data_all.loc[share_category_data_all['period']=='VMR Period', 'dollars_category'].values[0]
units_promoted_vmr_all = share_category_data_all.loc[share_category_data_all['period']=='VMR Period', 'units_promoted'].values[0]
units_category_vmr_all = share_category_data_all.loc[share_category_data_all['period']=='VMR Period', 'units_category'].values[0]

dollars_promoted_yago_all = share_category_data_all.loc[share_category_data_all['period']=='YAGO Period', 'dollars_promoted'].values[0]
dollars_category_yago_all = share_category_data_all.loc[share_category_data_all['period']=='YAGO Period', 'dollars_category'].values[0]
units_promoted_yago_all = share_category_data_all.loc[share_category_data_all['period']=='YAGO Period', 'units_promoted'].values[0]
units_category_yago_all = share_category_data_all.loc[share_category_data_all['period']=='YAGO Period', 'units_category'].values[0]

share_vmr_dol = dollars_promoted_vmr_all/dollars_category_vmr_all
share_vmr_dol = 0 if np.isnan(share_vmr_dol) else share_vmr_dol
share_yago_dol = dollars_promoted_yago_all/dollars_category_yago_all
share_yago_dol = 0 if np.isnan(share_yago_dol) else share_yago_dol
share_vmr_unit = units_promoted_vmr_all/units_category_vmr_all
share_vmr_unit = 0 if np.isnan(share_vmr_unit) else share_vmr_unit
share_yago_unit = units_promoted_yago_all/units_category_yago_all
share_yago_unit = 0 if np.isnan(share_yago_unit) else share_yago_unit



# ## **14) GETTING DATA FOR PRIOR 52 WEEKS (PER 28 WEEKS) AND VMR CAMPAIGN PARTICIPANTS SLIDE:**

# In[69]:


vmr_length = (reward_print_end - reward_print_start).days
print(vmr_length)


# In[70]:


curs.execute(f'''

    DROP table if exists VMR_{brand_nm}_prior_table_shares;
    CREATE temp table VMR_{brand_nm}_prior_table_shares as
        
          SELECT CASE WHEN prior_period = 'Prior-period ({round(pre_weeks)} weeks)' THEN 'Prior 52 wks (per {vmr_length} days)'
                     ELSE 'NA' END as Time_Period,
          
                COUNT(distinct o.ord_designated_cnsmr_id_key) as campaign_buyers,
                COUNT(distinct o.ord_event_key) as campaign_trips,
                SUM(o.purch_amt) as campaign_dollars,
                SUM(o.purch_qty) as campaign_units,
                
                (campaign_dollars:: float/campaign_buyers:: float) :: float as dollars_per_buyer_prior,
                (campaign_units:: float/campaign_buyers:: float) :: float as units_per_buyer_prior,
                (campaign_trips:: float/campaign_buyers:: float) :: float as trips_per_buyer_prior,
                (campaign_dollars:: float/campaign_trips:: float) :: float as dollars_per_trip_prior,
                (campaign_units:: float/campaign_trips:: float) :: float as units_per_trip_prior,
                
                round(((dollars_per_buyer_prior)/364)*{vmr_length} :: float, 4) as dollars_per_buyer,
                round(((units_per_buyer_prior)/364)*{vmr_length} ::float, 4) as units_per_buyer,
                round(((trips_per_buyer_prior)/364)*{vmr_length} :: float, 4) as trips_per_buyer,
                round((dollars_per_trip_prior) :: float, 4) as dollars_per_trip,
                round((units_per_trip_prior) :: float, 4) as units_per_trip
                
                
          FROM ord_trd_itm_cnsmr_fact_ne_v o
                 
                 INNER JOIN VMR_{brand_nm}_date_filter b ON (o.ord_date_key = b.date_key)
                 INNER JOIN VMR_{brand_nm}_upc_filter c ON (o.trade_item_key = c.trade_item_key)
                INNER JOIN VMR_{brand_nm}_reward_ids_csmr t ON (t.ord_designated_cnsmr_id_key = o.ord_designated_cnsmr_id_key)                 
                
          WHERE o.purch_amt > 0
                and o.purch_qty > 0
                and prior_period = 'Prior-period ({round(pre_weeks)} weeks)' and prior_period != 'NA'
                and c.brand_nbr in ({brand_nbr_str})
                and cal_dt between '{prior_period_start}' and '{prior_period_end}'
                
        GROUP BY 1 
        
        DISTRIBUTE REPLICATE;
        
''')




curs.execute(f'''       
    DROP table if exists VMR_{brand_nm}_vmr_metrics_participants;
    CREATE temp table VMR_{brand_nm}_vmr_metrics_participants as
    
        SELECT  
                CASE WHEN analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN 'Campaign Period'
                     ELSE 'NA' END as Time_Period,
        
                COUNT(distinct a.ord_designated_cnsmr_id_key) as campaign_buyers,
                COUNT(distinct a.ord_event_key) as campaign_trips,
                SUM(a.purch_amt) as campaign_dollars,
                SUM(a.purch_qty) as campaign_units,
                
                round((campaign_dollars::float/campaign_buyers::float),4) as dollars_per_buyer,
                round((campaign_units::float/campaign_buyers::float),4) as units_per_buyer,
                round((campaign_trips::float/campaign_buyers::float),4) as trips_per_buyer,
                round((campaign_dollars::float/campaign_trips::float),4) as dollars_per_trip,
                round((campaign_units::float/campaign_trips::float),4) as units_per_trip
                
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
            INNER JOIN VMR_{brand_nm}_reward_ids_csmr t ON (t.ord_designated_cnsmr_id_key = a.ord_designated_cnsmr_id_key)
            
            
        WHERE analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)'
              AND a.purch_amt > 0
              and a.purch_qty > 0
              and brand_nbr IN ({brand_nbr_str})
              and Time_Period != 'NA'
              
        GROUP BY 1
        
   
        UNION
        
        SELECT  
                Time_Period,
        
                campaign_buyers,
                campaign_trips,
                campaign_dollars,
                campaign_units,
                
                dollars_per_buyer,
                units_per_buyer,
                trips_per_buyer,
                dollars_per_trip,
                units_per_trip
                
        FROM VMR_{brand_nm}_prior_table_shares
        
        
    DISTRIBUTE REPLICATE;

''')


# In[71]:


results_participants = pd.read_sql(f'''
    SELECT *
    FROM VMR_{brand_nm}_vmr_metrics_participants
    ORDER BY CASE WHEN Time_Period = 'Campaign Period' then 1
              WHEN Time_Period = 'Prior 52 wks (per {vmr_length} days)' then 2
              END

    ''',conn)

results_participants


# In[72]:


results_participants = results_participants[['time_period', 'campaign_buyers', 'campaign_dollars', 'campaign_units', 'dollars_per_buyer', 'units_per_buyer', 'trips_per_buyer', 'dollars_per_trip', 'units_per_trip']]
results_participants = results_participants.rename(columns = {'time_period':'Time Period', 'campaign_buyers': 'Buyers', 'campaign_dollars':'Brand Dollars', 'campaign_units': 'Brand Units', 'dollars_per_buyer':'Dollars per Buyer', 'units_per_buyer':'Units per Buyer', 'trips_per_buyer':'Trips per Buyer', 'dollars_per_trip':'Dollars per Trip', 'units_per_trip':'Units per Trip'})

change_buyers = (results_participants['Buyers'][0] - results_participants['Buyers'][1])/results_participants['Buyers'][1]
change_brand_dollars = (results_participants['Brand Dollars'][0] - results_participants['Brand Dollars'][1])/results_participants['Brand Dollars'][1]
change_brand_units = (results_participants['Brand Units'][0] - results_participants['Brand Units'][1])/results_participants['Brand Units'][1]
change_dpb = (results_participants['Dollars per Buyer'][0] - results_participants['Dollars per Buyer'][1])/results_participants['Dollars per Buyer'][1]
change_upb = (results_participants['Units per Buyer'][0] - results_participants['Units per Buyer'][1])/results_participants['Units per Buyer'][1]
change_tpb = (results_participants['Trips per Buyer'][0] - results_participants['Trips per Buyer'][1])/results_participants['Trips per Buyer'][1]
change_dpt = (results_participants['Dollars per Trip'][0] - results_participants['Dollars per Trip'][1])/results_participants['Dollars per Trip'][1]
change_upt = (results_participants['Units per Trip'][0] - results_participants['Units per Trip'][1])/results_participants['Units per Trip'][1]

if str(change_buyers) == 'inf':
    change_buyers = ''
if str(change_brand_dollars) == 'inf':
    change_brand_dollars = ''
if str(change_brand_units) == 'inf':
    change_brand_units = ''
if str(change_dpb) == 'inf':
    change_dpb = ''
if str(change_upb) == 'inf':
    change_upb = ''
if str(change_tpb) == 'inf':
    change_tpb = ''
if str(change_dpt) == 'inf':
    change_dpt = ''
if str(change_upt) == 'inf':
    change_upt = ''

    
new_row = pd.DataFrame({'Time Period' : ['% Change'], 'Buyers': [change_buyers], 'Brand Dollars': [change_brand_dollars], 'Brand Units': [change_brand_units], 'Dollars per Buyer': [change_dpb], 'Units per Buyer':[change_upb], 'Trips per Buyer': [change_tpb], 'Dollars per Trip': [change_dpt], 'Units per Trip': [change_upt]})
results_participants = pd.concat([results_participants, new_row], ignore_index=True)
results_participants = results_participants.T.reset_index().T.reset_index(drop=True)
results_participants_1 = results_participants
results_participants_01 = results_participants_1 
print(results_participants_01)


curs.execute(f'''

    DROP table if exists VMR_{brand_nm}_prior_table_shares;
    CREATE temp table VMR_{brand_nm}_prior_table_shares as
        
          SELECT CASE WHEN analysis_period = 'VMR Pre-period ({round(weeks_per_period_final)} weeks)' THEN 'VMR Pre Period ({round(weeks_per_period_final)} weeks)'
                     ELSE 'NA' END as Time_Period,
                     
                     min(cal_dt) as min,
                     max(cal_dt) as max,
          
                COUNT(distinct o.ord_designated_cnsmr_id_key) as campaign_buyers,
                COUNT(distinct o.ord_event_key) as campaign_trips,
                SUM(o.purch_amt) as campaign_dollars,
                SUM(o.purch_qty) as campaign_units,
                
                round((campaign_dollars::float/campaign_buyers::float),4) as dollars_per_buyer,
                round((campaign_units::float/campaign_buyers::float),4) as units_per_buyer,
                round((campaign_trips::float/campaign_buyers::float),4) as trips_per_buyer,
                round((campaign_dollars::float/campaign_trips::float),4) as dollars_per_trip,
                round((campaign_units::float/campaign_trips::float),4) as units_per_trip
                
                
          FROM ord_trd_itm_cnsmr_fact_ne_v o
                 
                 INNER JOIN VMR_{brand_nm}_date_filter b ON (o.ord_date_key = b.date_key)
                 INNER JOIN VMR_{brand_nm}_upc_filter c ON (o.trade_item_key = c.trade_item_key)
                INNER JOIN VMR_{brand_nm}_reward_ids_csmr t ON (t.ord_designated_cnsmr_id_key = o.ord_designated_cnsmr_id_key)                 
                
          WHERE o.purch_amt > 0
                and o.purch_qty > 0
                and analysis_period = 'VMR Pre-period ({round(weeks_per_period_final)} weeks)'
                and c.brand_nbr in ({brand_nbr_str})
                and Time_Period != 'NA'

                
        GROUP BY 1 
        
        DISTRIBUTE REPLICATE;
        
''')


curs.execute(f'''       
    DROP table if exists VMR_{brand_nm}_vmr_metrics_participants;
    CREATE temp table VMR_{brand_nm}_vmr_metrics_participants as
    
        SELECT  
                CASE WHEN analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN 'Campaign Period'
                     ELSE 'NA' END as Time_Period,
                     
                     
                     min(cal_dt) as min,
                     max(cal_dt) as max,
        
                COUNT(distinct a.ord_designated_cnsmr_id_key) as campaign_buyers,
                COUNT(distinct a.ord_event_key) as campaign_trips,
                SUM(a.purch_amt) as campaign_dollars,
                SUM(a.purch_qty) as campaign_units,
                
                round((campaign_dollars::float/campaign_buyers::float),4) as dollars_per_buyer,
                round((campaign_units::float/campaign_buyers::float),4) as units_per_buyer,
                round((campaign_trips::float/campaign_buyers::float),4) as trips_per_buyer,
                round((campaign_dollars::float/campaign_trips::float),4) as dollars_per_trip,
                round((campaign_units::float/campaign_trips::float),4) as units_per_trip
                
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
            INNER JOIN VMR_{brand_nm}_reward_ids_csmr t ON (t.ord_designated_cnsmr_id_key = a.ord_designated_cnsmr_id_key)
            --INNER JOIN shopper_consistent_{analyst} v ON (a.ord_designated_cnsmr_id_key = v.cnsmr_id_key)
            
            
        WHERE analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)'
              AND a.purch_amt > 0
              and a.purch_qty > 0
              and brand_nbr IN ({brand_nbr_str})
              and Time_Period != 'NA'
              
        GROUP BY 1
        
   
        UNION
        
        
        SELECT  
                Time_Period,
                
                min,
                max,
        
                campaign_buyers,
                campaign_trips,
                campaign_dollars,
                campaign_units,
                
                dollars_per_buyer,
                units_per_buyer,
                trips_per_buyer,
                dollars_per_trip,
                units_per_trip
                
        FROM VMR_{brand_nm}_prior_table_shares
        
        
    DISTRIBUTE REPLICATE;

''')


results_participants = pd.read_sql(f'''
    SELECT *
    FROM VMR_{brand_nm}_vmr_metrics_participants
    ORDER BY CASE WHEN Time_Period = 'Campaign Period' then 1
              END

    ''',conn)

results_participants


results_participants = results_participants[['time_period', 'campaign_buyers', 'campaign_dollars', 'campaign_units', 'dollars_per_buyer', 'units_per_buyer', 'trips_per_buyer', 'dollars_per_trip', 'units_per_trip']]
results_participants = results_participants.rename(columns = {'time_period':'Time Period', 'campaign_buyers': 'Buyers', 'campaign_dollars':'Brand Dollars', 'campaign_units': 'Brand Units', 'dollars_per_buyer':'Dollars per Buyer', 'units_per_buyer':'Units per Buyer', 'trips_per_buyer':'Trips per Buyer', 'dollars_per_trip':'Dollars per Trip', 'units_per_trip':'Units per Trip'})

change_buyers = (results_participants['Buyers'][0] - results_participants['Buyers'][1])/results_participants['Buyers'][1]
change_brand_dollars = (results_participants['Brand Dollars'][0] - results_participants['Brand Dollars'][1])/results_participants['Brand Dollars'][1]
change_brand_units = (results_participants['Brand Units'][0] - results_participants['Brand Units'][1])/results_participants['Brand Units'][1]
change_dpb = (results_participants['Dollars per Buyer'][0] - results_participants['Dollars per Buyer'][1])/results_participants['Dollars per Buyer'][1]
change_upb = (results_participants['Units per Buyer'][0] - results_participants['Units per Buyer'][1])/results_participants['Units per Buyer'][1]
change_tpb = (results_participants['Trips per Buyer'][0] - results_participants['Trips per Buyer'][1])/results_participants['Trips per Buyer'][1]
change_dpt = (results_participants['Dollars per Trip'][0] - results_participants['Dollars per Trip'][1])/results_participants['Dollars per Trip'][1]
change_upt = (results_participants['Units per Trip'][0] - results_participants['Units per Trip'][1])/results_participants['Units per Trip'][1]

if str(change_buyers) == 'inf':
    change_buyers = ''
if str(change_brand_dollars) == 'inf':
    change_brand_dollars = ''
if str(change_brand_units) == 'inf':
    change_brand_units = ''
if str(change_dpb) == 'inf':
    change_dpb = ''
if str(change_upb) == 'inf':
    change_upb = ''
if str(change_tpb) == 'inf':
    change_tpb = ''
if str(change_dpt) == 'inf':
    change_dpt = ''
if str(change_upt) == 'inf':
    change_upt = ''

    
new_row = pd.DataFrame({'Time Period' : ['% Change'], 'Buyers': [change_buyers], 'Brand Dollars': [change_brand_dollars], 'Brand Units': [change_brand_units], 'Dollars per Buyer': [change_dpb], 'Units per Buyer':[change_upb], 'Trips per Buyer': [change_tpb], 'Dollars per Trip': [change_dpt], 'Units per Trip': [change_upt]})
results_participants = pd.concat([results_participants, new_row], ignore_index=True)
results_participants = results_participants.T.reset_index().T.reset_index(drop=True)
results_participants_2 = results_participants
results_participants_02 = results_participants_2
print(results_participants_02)


curs.execute(f'''

    DROP table if exists VMR_{brand_nm}_prior_table_shares;
    CREATE temp table VMR_{brand_nm}_prior_table_shares as
        
          SELECT CASE WHEN analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)' THEN 'VMR Period - YAGO ({round(weeks_per_period_final)} weeks)'

                     ELSE 'NA' END as Time_Period,
                     
                     min(cal_dt) as min,
                     max(cal_dt) as max,
        
          
                COUNT(distinct o.ord_designated_cnsmr_id_key) as campaign_buyers,
                COUNT(distinct o.ord_event_key) as campaign_trips,
                SUM(o.purch_amt) as campaign_dollars,
                SUM(o.purch_qty) as campaign_units,
                
                round((campaign_dollars::float/campaign_buyers::float),4) as dollars_per_buyer,
                round((campaign_units::float/campaign_buyers::float),4) as units_per_buyer,
                round((campaign_trips::float/campaign_buyers::float),4) as trips_per_buyer,
                round((campaign_dollars::float/campaign_trips::float),4) as dollars_per_trip,
                round((campaign_units::float/campaign_trips::float),4) as units_per_trip
                
                
          FROM ord_trd_itm_cnsmr_fact_ne_v o
                 
                 INNER JOIN VMR_{brand_nm}_date_filter b ON (o.ord_date_key = b.date_key)
                 INNER JOIN VMR_{brand_nm}_upc_filter c ON (o.trade_item_key = c.trade_item_key)
                INNER JOIN VMR_{brand_nm}_reward_ids_csmr t ON (t.ord_designated_cnsmr_id_key = o.ord_designated_cnsmr_id_key)                 
                
          WHERE o.purch_amt > 0
                and o.purch_qty > 0
                and analysis_period = 'Period Year Ago ({round(weeks_per_period_final)} weeks)'
                and c.brand_nbr in ({brand_nbr_str})
                and Time_Period != 'NA'

                
        GROUP BY 1 
        
        DISTRIBUTE REPLICATE;
        
''')


curs.execute(f'''       
    DROP table if exists VMR_{brand_nm}_vmr_metrics_participants;
    CREATE temp table VMR_{brand_nm}_vmr_metrics_participants as
    
        SELECT  
                CASE WHEN analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)' THEN 'Campaign Period'
                     ELSE 'NA' END as Time_Period,
                     
                     min(cal_dt) as min,
                     max(cal_dt) as max,
        
        
                COUNT(distinct a.ord_designated_cnsmr_id_key) as campaign_buyers,
                COUNT(distinct a.ord_event_key) as campaign_trips,
                SUM(a.purch_amt) as campaign_dollars,
                SUM(a.purch_qty) as campaign_units,
                
                round((campaign_dollars::float/campaign_buyers::float),4) as dollars_per_buyer,
                round((campaign_units::float/campaign_buyers::float),4) as units_per_buyer,
                round((campaign_trips::float/campaign_buyers::float),4) as trips_per_buyer,
                round((campaign_dollars::float/campaign_trips::float),4) as dollars_per_trip,
                round((campaign_units::float/campaign_trips::float),4) as units_per_trip
                
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
            INNER JOIN VMR_{brand_nm}_reward_ids_csmr t ON (t.ord_designated_cnsmr_id_key = a.ord_designated_cnsmr_id_key)
            --INNER JOIN shopper_consistent_{analyst} v ON (a.ord_designated_cnsmr_id_key = v.cnsmr_id_key)
            
            
        WHERE analysis_period = 'VMR Print Period ({round(weeks_per_period_final)} weeks)'
              AND a.purch_amt > 0
              and a.purch_qty > 0
              and brand_nbr IN ({brand_nbr_str})
              and Time_Period != 'NA'
              
        GROUP BY 1
        
   
        UNION
        
        
        SELECT  
                Time_Period,
                
                min,
                max,
        
                campaign_buyers,
                campaign_trips,
                campaign_dollars,
                campaign_units,
                
                dollars_per_buyer,
                units_per_buyer,
                trips_per_buyer,
                dollars_per_trip,
                units_per_trip
                
        FROM VMR_{brand_nm}_prior_table_shares
        
        
    DISTRIBUTE REPLICATE;

''')


results_participants = pd.read_sql(f'''
    SELECT *
    FROM VMR_{brand_nm}_vmr_metrics_participants
    ORDER BY CASE WHEN Time_Period = 'Campaign Period' then 1
              END

    ''',conn)

results_participants


results_participants = results_participants[['time_period', 'campaign_buyers', 'campaign_dollars', 'campaign_units', 'dollars_per_buyer', 'units_per_buyer', 'trips_per_buyer', 'dollars_per_trip', 'units_per_trip']]
results_participants = results_participants.rename(columns = {'time_period':'Time Period', 'campaign_buyers': 'Buyers', 'campaign_dollars':'Brand Dollars', 'campaign_units': 'Brand Units', 'dollars_per_buyer':'Dollars per Buyer', 'units_per_buyer':'Units per Buyer', 'trips_per_buyer':'Trips per Buyer', 'dollars_per_trip':'Dollars per Trip', 'units_per_trip':'Units per Trip'})

change_buyers = (results_participants['Buyers'][0] - results_participants['Buyers'][1])/results_participants['Buyers'][1]
change_brand_dollars = (results_participants['Brand Dollars'][0] - results_participants['Brand Dollars'][1])/results_participants['Brand Dollars'][1]
change_brand_units = (results_participants['Brand Units'][0] - results_participants['Brand Units'][1])/results_participants['Brand Units'][1]
change_dpb = (results_participants['Dollars per Buyer'][0] - results_participants['Dollars per Buyer'][1])/results_participants['Dollars per Buyer'][1]
change_upb = (results_participants['Units per Buyer'][0] - results_participants['Units per Buyer'][1])/results_participants['Units per Buyer'][1]
change_tpb = (results_participants['Trips per Buyer'][0] - results_participants['Trips per Buyer'][1])/results_participants['Trips per Buyer'][1]
change_dpt = (results_participants['Dollars per Trip'][0] - results_participants['Dollars per Trip'][1])/results_participants['Dollars per Trip'][1]
change_upt = (results_participants['Units per Trip'][0] - results_participants['Units per Trip'][1])/results_participants['Units per Trip'][1]

if str(change_buyers) == 'inf':
    change_buyers = ''
if str(change_brand_dollars) == 'inf':
    change_brand_dollars = ''
if str(change_brand_units) == 'inf':
    change_brand_units = ''
if str(change_dpb) == 'inf':
    change_dpb = ''
if str(change_upb) == 'inf':
    change_upb = ''
if str(change_tpb) == 'inf':
    change_tpb = ''
if str(change_dpt) == 'inf':
    change_dpt = ''
if str(change_upt) == 'inf':
    change_upt = ''

    
new_row = pd.DataFrame({'Time Period' : ['% Change'], 'Buyers': [change_buyers], 'Brand Dollars': [change_brand_dollars], 'Brand Units': [change_brand_units], 'Dollars per Buyer': [change_dpb], 'Units per Buyer':[change_upb], 'Trips per Buyer': [change_tpb], 'Dollars per Trip': [change_dpt], 'Units per Trip': [change_upt]})
results_participants = pd.concat([results_participants, new_row], ignore_index=True)
results_participants = results_participants.T.reset_index().T.reset_index(drop=True)
results_participants_3 = results_participants
results_participants_03 = results_participants_3
print(results_participants_03)




# ## **15) GETTING THE TREND FOR ALL BRAND TRANSACTIONS (65 WEEKS PRIOR PERIOD + VMR PERIOD):**

# In[107]:



vmr_preperiod_start = int(str(vmr_pre_period_start.strftime('%Y-%m-%d')).replace('-',''))
vmr_preperiod_end = int(str(vmr_pre_period_end.strftime('%Y-%m-%d')).replace('-',''))

yago_26wk_start_var = int(str(yago_26wk_start.strftime('%Y-%m-%d')).replace('-',''))
vmr_pre26wk_start_var = int(str(vmr_pre26wk_start.strftime('%Y-%m-%d')).replace('-',''))
yago_start = int(str(yago_print_start.strftime('%Y-%m-%d')).replace('-',''))
reward_start = int(str(reward_print_start.strftime('%Y-%m-%d')).replace('-',''))
reward_end = int(str(reward_print_end.strftime('%Y-%m-%d')).replace('-',''))
yago_start = int(str(yago_print_start.strftime('%Y-%m-%d')).replace('-',''))

#if we want to use redemption end as a max date limit for the trend:
#redempt_end = int(str(redemption_end.strftime('%Y-%m-%d')).replace('-','')) #redemption_end defined at the beggining of the code
#redemption_end_yago = redemption_end - dt.timedelta(weeks = 52)
#yago_end = int(str(redemption_end_yago.strftime('%Y-%m-%d')).replace('-',''))

#if we want to use post 4wk redemption period as a max date limit for the trend
redempt_end = int(str(post_4wk_end.strftime('%Y-%m-%d')).replace('-','')) #redemption_end defined at the beggining of the code
redemption_end_yago = post_4wk_end - dt.timedelta(weeks = 52)
yago_end = int(str(redemption_end_yago.strftime('%Y-%m-%d')).replace('-',''))


curs.execute(f'''       
    DROP table if exists VMR_{brand_nm}_brand_trend;
    CREATE temp table VMR_{brand_nm}_brand_trend as
    
    
        SELECT  
               cal_sun_wk_ending_dt,
               CASE WHEN ord_date_key between {yago_start} and {yago_end} THEN 'YAGO Period'
                    WHEN ord_date_key between {yago_end}+1 and {reward_start} THEN 'Prior Period in between'
                    WHEN ord_date_key between {reward_start}+1 and {reward_end} THEN 'VMR Period'
                    ELSE 'NA'
                    END as analysis_periods,
               CASE WHEN ord_date_key between '{vmr_preperiod_start}' and '{vmr_preperiod_end}' THEN 'VMR Prior Period'
                    WHEN ord_date_key between {reward_start}+1 and {reward_end} THEN 'VMR Period'
                    WHEN ord_date_key between {yago_start} and {yago_end} THEN 'YAGO Period'
                    ELSE 'NA'
                    END as analysis_periods_second_chart,
                CASE WHEN ord_date_key between {vmr_pre26wk_start_var} and {redempt_end} THEN 'Recent Trend Period'
                     WHEN ord_date_key between {yago_26wk_start_var} and {yago_end} THEN 'YAGO Trend Period'
                    ELSE 'NA'
                    END as trend_chart,
               count(distinct a.ord_event_key) as count_distinct_trips,
               sum(a.purch_amt) as dollar_sales,
               dollar_sales/count_distinct_trips as dollars_per_trip,
               sum(a.purch_qty) as units_sales,
               units_sales/count_distinct_trips::float as units_per_trip,
               RANK() OVER(PARTITION BY analysis_periods ORDER BY cal_sun_wk_ending_dt ASC) as nbr_week
              
        FROM ord_trd_itm_cnsmr_fact_ne_v a 
        
            INNER JOIN VMR_{brand_nm}_date_filter b ON (a.ord_date_key = b.date_key)
            INNER JOIN VMR_{brand_nm}_upc_filter c ON (a.trade_item_key = c.trade_item_key)
            INNER JOIN VMR_{brand_nm}_printing_stores s ON (a.ord_touchpoint_key = s.touchpoint_key)
            
        WHERE brand_nbr IN ({brand_nbr_str})
        
        GROUP BY 1,2,3,4
        
        ORDER BY 1 DESC
        
    DISTRIBUTE REPLICATE;

''')


# In[108]:


vmr_trend_brand = pd.read_sql(f'''
    SELECT *
    FROM VMR_{brand_nm}_brand_trend
    ORDER BY 1 DESC
    ''',conn)

vmr_trend_brand


# In[109]:


vmr_trend_brand_sl = pd.read_sql(f'''
    SELECT cal_sun_wk_ending_dt, 
           trend_chart,
           SUM(dollar_sales) as dollar_sales,
           SUM(units_sales) as units_sales
    FROM VMR_{brand_nm}_brand_trend
    WHERE trend_chart = 'Recent Trend Period' OR trend_chart = 'YAGO Trend Period'
    GROUP BY 1,2
    ORDER BY 1 DESC
    
    ''',conn)

vmr_trend_brand_sl


vmr_trend_brand_sl_exc = pd.read_sql(f'''
    SELECT cal_sun_wk_ending_dt as week, 
           SUM(dollar_sales) as dollars,
           SUM(units_sales) as units,
           SUM(count_distinct_trips) as trips,
           analysis_periods_second_chart as period
    FROM VMR_{brand_nm}_brand_trend
    WHERE analysis_periods_second_chart != 'NA'
    GROUP BY 1,5
    ORDER BY 1 DESC
    
    ''',conn)

vmr_trend_brand_sl_exc = vmr_trend_brand_sl_exc.rename(columns = {'week':'Week', 'dollars':'Dollars', 'units':'Units', 'period':'Period'})
vmr_trend_brand_sl_exc


# In[76]:

# ## **16) CREATING EXCEL OUTPUT:**

print('Excel')
from pandas import ExcelWriter
from pandas import ExcelFile
import xlsxwriter
from openpyxl import Workbook
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.styles import NamedStyle, PatternFill, Border, Side, Alignment, Protection, Font
from openpyxl.drawing.image import Image
from datetime import date


if segment_type == 1:
    segment_type_def = "UPC"
elif segment_type == 2:
    segment_type_def = "BRAND"
elif segment_type == 3:
    segment_type_def = "CATEGORY"
elif segment_type == 4:
    segment_type_def = "CUSTOM BRAND DESCR"
elif segment_type == 5:
    segment_type_def = "RETAILER DESCR"
elif segment_type == 6:
    segment_type_def = "KROGER DESCR"


export_file_name = report_name_for_export + "_" + segment_type_def +  "_VMR_Scorecard"  # Name of the final Excel file 

head_color = "Accent1"

wb = Workbook()
ws1 = wb.active
ws1.sheet_view.showGridLines = False

# create easier to read named styles to control number format in Excel
format_comma_no_decimal = NamedStyle(name='format_comma_no_decimal')
format_comma_no_decimal.number_format = '#,##0'
wb.add_named_style(format_comma_no_decimal)

format_comma_one_decimal = NamedStyle(name='format_comma_one_decimal')
format_comma_one_decimal.number_format = '#,##0.0'
wb.add_named_style(format_comma_one_decimal)

format_short_date = NamedStyle(name='format_short_date')
format_short_date.number_format = 'm/d/yyyy'
wb.add_named_style(format_short_date)

format_dollars_two_decimals = NamedStyle(name='format_dollars_two_decimals')
format_dollars_two_decimals.number_format = '$#,##0.00'
wb.add_named_style(format_dollars_two_decimals)

format_dollars_no_decimals = NamedStyle(name='format_dollars_no_decimals')
format_dollars_no_decimals.number_format = '$#,##0'
wb.add_named_style(format_dollars_no_decimals)

format_percent_one_decimal = NamedStyle(name='format_percent_one_decimal')
format_percent_one_decimal.number_format = '0.0%'
wb.add_named_style(format_percent_one_decimal)

format_percent_no_decimal = NamedStyle(name='format_percent_no_decimal')
format_percent_no_decimal.number_format = '0%'
wb.add_named_style(format_percent_no_decimal)


##### WORKSHEET 1 ####

print("worksheet1")

ws1.title = 'PARAMETERS'

ws1['A1'] = 'Brand Name:'
ws1['B1'] = brand_nm

ws1['A3'] = 'UPC LMC List:'
ws1['B3'] = lmc_list_id

ws1['A5'] = 'Brands:'
ws1['B5'] = brand_nbr_str
ws1['A6'] = 'Categories:'
ws1['B6'] = cat_nbr_str

ws1['A8'] = 'Reward BL Codes:'
ws1['B8'] = BL_CODES

ws1['A10'] = 'Announcements:'
ws1['B10'] = Announcement

ws1['A12'] = 'Qualifying Reward Amount:'
ws1['B12'] = min_threshold_statement

ws1['A14'] = 'Redemption days:'
ws1['B14'] = redemption_days

ws1['B15'] = ''

ws1['A16'] = 'Static_1:'
ws1['B16'] = static_1
ws1['A17'] = 'Static_2:'
ws1['B17'] = static_2


ws1['B3'].alignment = Alignment(wrap_text=True, horizontal="left" )

for i in range(14,18):
    
    ws1['B'+str(i)].alignment = Alignment(wrap_text=True, horizontal="left" )
    
ws1['A18'] = ''

ws1['A19'] = 'Trackable ID:'
ws1['A20'] = 'Qualifying Trip:'
ws1['A21'] = 'Qualifying ID:'
ws1['A22'] = 'Reward ID:'
ws1['B19'] = 'ID that we can track purchases over time'
ws1['B20'] = 'Trip where the qualifying reward amount was purchased'
ws1['B21'] = 'Trackable ID that purchased the qualifying amount'
ws1['B22'] = 'Trackable ID that purchased the qualifying amount AND received a reward print'

ws1['A23'] = ''

ws1['A24'] = 'Reward BLs Analyzed'

for r in dataframe_to_rows(df_promo_summary, index=False, header=True):
    ws1.append(r)

for cell in ws1[25]:
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")

for i in range(1,len(df_promo_summary)+1):
    for cell in ws1[25+i]:
        cell.alignment=Alignment(wrap_text=True, horizontal="center")

        
ws1['A'+str(25 + len(df_promo_summary)+1)] = ''

n = 25 + len(df_promo_summary)+1 


#------------------------------------------------------------

ws1['A'+str(n+2)] = "Analysis Period Dates"
    
for r in dataframe_to_rows(df_date_check, index=False, header=True):
    ws1.append(r)
    
for cell in ws1[n+3]:
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center" ) 

    
#------------------------------------------------------------
ws1['A'+str(n+10)] = 'Note:  VMR Period = Campaign Period'

ws1['A'+str(n+12)] = 'UPC Hierarchy'
    
for r in dataframe_to_rows(df_upcs_quick[['brand_nbr', 'brand_desc', 'nbr_of_upcs']], index=False, header=True):
    ws1.append(r) 
    
for cell in ws1[n+13]:
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    


#------------------------------------------------------------

next_cell = n+12+len(df_upcs_quick[['brand_nbr', 'brand_desc', 'nbr_of_upcs']]) + 4


if segment_type_def != 'UPC':
    
    ws1['A'+str(next_cell)] = "Segment List based on " + segment_type_def

    for r in dataframe_to_rows(df_segments_2, index=False, header=True):
        ws1.append(r)

    for cell in ws1[next_cell+1]:
        cell.style = 'Accent1'
        cell.alignment=Alignment(wrap_text=True, horizontal="center" ) 

    
    
##############################################################################

for cell in ws1[1]:
    cell.font = Font(bold=True, color = '4F81BD', size = 15, name = 'Calibri' )

ws1.column_dimensions['A'].width = 35        

for col_name in ['B','C']:
    for cols in ws1[col_name]:
        ws1.column_dimensions[col_name].width = 23       

for col_name in ['D','E','F','G']:
    for cols in ws1[col_name]:
        ws1.column_dimensions[col_name].width = 18 

for i in range(n+13,n+16):
    for j in ['B', 'C']:
        ws1[j+str(i)].alignment=Alignment(wrap_text=True, horizontal="center")
        
for i in range(n+14,n+16):
    for j in ['C']:       
        ws1[j+str(i)].style = format_comma_no_decimal
        
for i in range(n+2,n+10):
    for j in ['B', 'C']:
        ws1[j+str(i)].alignment=Alignment(wrap_text=True, horizontal="center")

        
for i in range(n+13,next_cell-1):
    for j in ['A', 'B', 'C']:
        ws1[j+str(i)].alignment=Alignment(wrap_text=True, horizontal="center")
        
        
for i in range(next_cell+1,next_cell+1+len(df_segments_2)+10):
    for j in ['A', 'B', 'C']:
        ws1[j+str(i)].alignment=Alignment(wrap_text=True, horizontal="center")
        
        
for i in [19,20,21,22]:
    ws1['A'+str(i)].font = Font(bold=True, color = '000000', size = 11, name = 'Calibri' )

    
##### WORKSHEET 2 ####

print("worksheet2")

ws8 = wb.create_sheet(title='TTL trips-Weekly Metrics')
ws8.sheet_view.showGridLines = False

ws8['A1'] = 'Total trips - Weekly Metrics'

ws8['A2'] = 'Metrics based on all trips where the promoted product was purchased (reward level and non-reward level trips)'

vmr_trend_brand_sl_2 = vmr_trend_brand_sl[['cal_sun_wk_ending_dt', 'dollar_sales', 'units_sales']]

for r in dataframe_to_rows(vmr_trend_brand_sl_2, index=False, header=True):
    ws8.append(r)

for cell in ws8[3]:
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    
    

for cell in ws8[1]:
    cell.font = Font(bold=True, color = '4F81BD', size = 15, name = 'Calibri' )



for col_name in ['A','B','C','D','E']:
    for cols in ws8[col_name]:
        ws8.column_dimensions[col_name].width = 25 
    

        
for i in (range(4,len(vmr_trend_brand_sl)+4)):
    for j in ('B'):
        ws8[j+str(i)].style = format_dollars_two_decimals
        
        
for i in (range(4,len(vmr_trend_brand_sl)+4)):
    for j in ('C'):
        ws8[j+str(i)].style = format_comma_no_decimal


        #Center all rows

##### WORKSHEET 2 ####

print("worksheet3")

ws9 = wb.create_sheet(title='TTL trips-Metrics')
ws9.sheet_view.showGridLines = False

ws9['A1'] = 'Total Trips Metrics'

ws9['A2'] = 'Metrics based on all trips where the promoted product was purchased (reward level and non-reward level trips)'


df_transposed = vmr_trend_yago_summary[['Dollars per Trip', 'Units per Trip']].T
df_transposed = df_transposed.reset_index()
df_transposed = df_transposed.rename(columns = {0:'YAGO Period', 1:'VMR Period', 'index': ''})
new_row_1 = {'':'Dollars Share','YAGO Period': "{:.1%}".format(share_yago_dol), 'VMR Period': "{:.1%}".format(share_vmr_dol)}
new_row_2 = {'':'Units Share','YAGO Period': "{:.1%}".format(share_yago_unit), 'VMR Period': "{:.1%}".format(share_vmr_unit)}
df_transposed.loc[len(df_transposed)] = new_row_1
df_transposed.loc[len(df_transposed)] = new_row_2


for r in dataframe_to_rows(df_transposed, index=False, header=True):
    ws9.append(r)

for cell in ws9[3]:
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    
    

for cell in ws9[1]:
    cell.font = Font(bold=True, color = '4F81BD', size = 15, name = 'Calibri' )

    
ws9.column_dimensions['A'].width = 40     


ws9['B4'].style = format_dollars_two_decimals
ws9['C4'].style = format_dollars_two_decimals


#update
ws9['B6'].style = format_comma_one_decimal
ws9['B7'].style = format_comma_one_decimal
ws9['C6'].style = format_comma_one_decimal
ws9['C7'].style = format_comma_one_decimal

#update

for i in range(4,9):
    for j in ['B', 'C']:
        ws9[j+str(i)].alignment=Alignment(wrap_text=True, horizontal="center")



##### WORKSHEET 2 ####

ws14 = wb.create_sheet(title='Reward Trips')
ws14.sheet_view.showGridLines = False

ws14['A1'] = 'Campaign Reward Trips'

ws14['A2'] = 'Metrics based on qualifying trips'


df_rwd_trans  = {
    
    '': ['All Other Transactions', 'Reward Level Transactions'],
    '% trips': ["{:.0%}".format(1-pct_reward_trips), "{:.0%}".format(pct_reward_trips)]
    
}

df_rwd_trans = pd.DataFrame(df_rwd_trans)


for r in dataframe_to_rows(df_rwd_trans, index=False, header=True):
    ws14.append(r)

for cell in ws14[3]:
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    
    

for cell in ws14[1]:
    cell.font = Font(bold=True, color = '4F81BD', size = 15, name = 'Calibri' )

    
#update
ws14['A6'] = ''
ws14['A7'] = 'VMR Campaign Period:'


for r in dataframe_to_rows(level_ct, index=False, header=True):
    ws14.append(r)

for cell in ws14[8]: #update
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    
    

row_end = 8 + len(level_ct)    


ws14['A'+str(row_end+1)] = ''


if len(level_ct_yago)>1:

    ws14['A'+str(row_end+2)] = 'YAGO Period:'

    for r in dataframe_to_rows(level_ct_yago, index=False, header=True):
        ws14.append(r)

    for cell in ws14[row_end+3]: #update
        cell.style = 'Accent1'
        cell.alignment=Alignment(wrap_text=True, horizontal="center")
        
else:
    ws14['A'+str(row_end+2)] = 'YAGO Period:'
    ws14['A'+str(row_end+3)] = '[No data available]'


row_end_2 = row_end+ 3 + len(level_ct_yago)    


ws14['A'+str(row_end_2+1)] = ''


if len(level_ct_prep)>1:

    ws14['A'+str(row_end_2+2)] = 'Pre-Period:'

    for r in dataframe_to_rows(level_ct_prep, index=False, header=True):
        ws14.append(r)

    for cell in ws14[row_end_2+3]: #update
        cell.style = 'Accent1'
        cell.alignment=Alignment(wrap_text=True, horizontal="center")  
           
else:
    ws14['A'+str(row_end_2+2)] = 'Pre-Period:'
    ws14['A'+str(row_end_2+3)] = '[No data available]'



row_end_3 = row_end_2 + 3 + len(level_ct_prep)

ws14['A'+str(row_end_3+1)] = ''


if len(level_ct_prior)>1:

    ws14['A'+str(row_end_3+2)] = 'Prior 52 wk Period:'

    for r in dataframe_to_rows(level_ct_prior, index=False, header=True):
        ws14.append(r)

    for cell in ws14[row_end_3+3]: #update
        cell.style = 'Accent1'
        cell.alignment=Alignment(wrap_text=True, horizontal="center")    
else:
    ws14['A'+str(row_end_3+2)] = 'Prior 52 wk Period:'
    ws14['A'+str(row_end_3+3)] = '[No data available]'


    
final_row = row_end_3+3+len(level_ct_prior)


ws14.column_dimensions['A'].width = 40     
ws14.column_dimensions['B'].width = 20
ws14.column_dimensions['C'].width = 20
ws14.column_dimensions['D'].width = 20


#update
for i in range(8, final_row+2):
    for j in ('B'):
        ws14[j+str(i)].style = format_dollars_two_decimals
        
for i in range(8, final_row+2):
    for j in ('C','D'):
        ws14[j+str(i)].style = format_comma_no_decimal
        

for cell in ws14[8]: #update
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    

for cell in ws14[row_end+3]: #update
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")

for cell in ws14[row_end_2+3]: #update
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")  

for cell in ws14[row_end_3+3]: #update
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    
        




##### WORKSHEET 2 ####

print("worksheet5")

ws10 = wb.create_sheet(title='Campaign Trends')
ws10.sheet_view.showGridLines = False

ws10['A1'] = 'Campaign Trends'

ws10['A2'] = 'Metrics based on qualifying IDs'


results_participants_1 = results_participants_01 
results_participants_1.columns = results_participants_1.iloc[0]
results_participants_1 = results_participants_1[1:]


for r in dataframe_to_rows(results_participants_1, index=False, header=True):
    ws10.append(r)

for cell in ws10[3]:
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    
    

for cell in ws10[1]:
    cell.font = Font(bold=True, color = '4F81BD', size = 15, name = 'Calibri' )

    
ws10.column_dimensions['A'].width = 40       
       


for i in range(4,6):
    for j in ('B'):
        ws10[j+str(i)].style = format_comma_no_decimal

for i in range(4,6):
    for j in ('C'):
        ws10[j+str(i)].style = format_dollars_two_decimals

for i in range(4,6):
    for j in ('D'):
        ws10[j+str(i)].style = format_comma_no_decimal

for i in range(4,6):
    for j in ('E'):
        ws10[j+str(i)].style = format_dollars_two_decimals


for i in range(4,6):
    for j in ('F'):
        ws10[j+str(i)].style = format_comma_one_decimal


for i in range(4,6):
    for j in ('G'):
        ws10[j+str(i)].style = format_comma_one_decimal


for i in range(4,6):
    for j in ('H'):
        ws10[j+str(i)].style = format_dollars_two_decimals


for i in range(4,6):
    for j in ('I'):
        ws10[j+str(i)].style = format_comma_one_decimal

        
ws10['B6'].style = format_percent_no_decimal        
ws10['C6'].style = format_percent_no_decimal        
ws10['D6'].style = format_percent_no_decimal        
ws10['E6'].style = format_percent_no_decimal        
ws10['F6'].style = format_percent_no_decimal
ws10['G6'].style = format_percent_no_decimal
ws10['H6'].style = format_percent_no_decimal
ws10['I6'].style = format_percent_no_decimal


#update

ws10['A7'] = '** Per X days is meant to equivailize metrics for the same length as campaign period over prior 52 weeks'

ws10['A8'] = '** Brand dollars and brand units are aggregated for the entire pre 52 weeks'


results_participants_2 = results_participants_02 
results_participants_2.columns = results_participants_2.iloc[0]
results_participants_2 = results_participants_2[1:]


for r in dataframe_to_rows(results_participants_2, index=False, header=True):
    ws10.append(r)

for cell in ws10[9]:
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    
    
       

for i in range(10,12):
    for j in ('B'):
        ws10[j+str(i)].style = format_comma_no_decimal

for i in range(10,12):
    for j in ('C'):
        ws10[j+str(i)].style = format_dollars_two_decimals

for i in range(10,12):
    for j in ('D'):
        ws10[j+str(i)].style = format_comma_no_decimal

for i in range(10,12):
    for j in ('E'):
        ws10[j+str(i)].style = format_dollars_two_decimals
                
for i in range(10,12):
    for j in ('F'):
        ws10[j+str(i)].style = format_comma_one_decimal

for i in range(10,12):
    for j in ('G'):
        ws10[j+str(i)].style = format_comma_one_decimal

for i in range(10,12):
    for j in ('H'):
        ws10[j+str(i)].style = format_dollars_two_decimals

for i in range(10,12):
    for j in ('I'):
        ws10[j+str(i)].style = format_comma_one_decimal
        
        
ws10['B12'].style = format_percent_no_decimal        
ws10['C12'].style = format_percent_no_decimal        
ws10['D12'].style = format_percent_no_decimal        
ws10['E12'].style = format_percent_no_decimal        
ws10['F12'].style = format_percent_no_decimal
ws10['G12'].style = format_percent_no_decimal
ws10['H12'].style = format_percent_no_decimal
ws10['I12'].style = format_percent_no_decimal     

        
#-------

ws10['A14'] = ''

results_participants_3 = results_participants_03 
results_participants_3.columns = results_participants_3.iloc[0]
results_participants_3 = results_participants_3[1:]


for r in dataframe_to_rows(results_participants_3, index=False, header=True):
    ws10.append(r)

for cell in ws10[15]:
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    
    
       

for i in range(16,18):
    for j in ('B'):
        ws10[j+str(i)].style = format_comma_no_decimal

for i in range(16,18):
    for j in ('C'):
        ws10[j+str(i)].style = format_dollars_two_decimals

for i in range(16,18):
    for j in ('D'):
        ws10[j+str(i)].style = format_comma_no_decimal

for i in range(16,18):
    for j in ('E'):
        ws10[j+str(i)].style = format_dollars_two_decimals
                
for i in range(16,18):
    for j in ('F'):
        ws10[j+str(i)].style = format_comma_one_decimal

for i in range(16,18):
    for j in ('G'):
        ws10[j+str(i)].style = format_comma_one_decimal

for i in range(16,18):
    for j in ('H'):
        ws10[j+str(i)].style = format_dollars_two_decimals

for i in range(16,18):
    for j in ('I'):
        ws10[j+str(i)].style = format_comma_one_decimal

ws10['B18'].style = format_percent_no_decimal        
ws10['C18'].style = format_percent_no_decimal        
ws10['D18'].style = format_percent_no_decimal        
ws10['E18'].style = format_percent_no_decimal        
ws10['F18'].style = format_percent_no_decimal
ws10['G18'].style = format_percent_no_decimal
ws10['H18'].style = format_percent_no_decimal
ws10['I18'].style = format_percent_no_decimal     

        
        

##### WORKSHEET 2 ####

print("worksheet6")

if category_show_parameter == True:

    ws11 = wb.create_sheet(title='Category Share')
    ws11.sheet_view.showGridLines = False

    ws11['A1'] = 'Category Share'

    ws11['A2'] = 'Metrics based on qualifying IDs'


    df_dollar_share = {

        '': ['YAGO Period', 'VMR Period', 'Post Period'],
        'Dollar Share': [share_dollar_yago, share_dollar_vmr, share_dollar_post_4wk]

    }

    df_dollar_share = pd.DataFrame(df_dollar_share)


    df_unit_share = {

        '': ['YAGO Period', 'VMR Period', 'Post Period'],
        'Unit Share': [share_units_yago, share_units_vmr, share_units_post_4wk]

    }

    df_unit_share = pd.DataFrame(df_unit_share)



    for r in dataframe_to_rows(df_dollar_share, index=False, header=True):
        ws11.append(r)

    for cell in ws11[3]:
        cell.style = 'Accent1'
        cell.alignment=Alignment(wrap_text=True, horizontal="center")    


    for cell in ws11[1]:
        cell.font = Font(bold=True, color = '4F81BD', size = 15, name = 'Calibri' )


    ws11.column_dimensions['A'].width = 40        


    ws11['A7'] = ''
    ws11['A8'] = ''


    for r in dataframe_to_rows(df_unit_share, index=False, header=True):
        ws11.append(r)

    for cell in ws11[9]:
        cell.style = 'Accent1'
        cell.alignment=Alignment(wrap_text=True, horizontal="center")    



    #update

    for i in range(4,7):
        for j in ('B'):
            ws11[j+str(i)].style = format_percent_no_decimal


    for i in range(10,15):
        for j in ('B'):
            ws11[j+str(i)].style = format_percent_no_decimal #check this

       
##### WORKSHEET 3 ####

print("worksheet7")
f'''
ws3 = wb.create_sheet(title='Per Trip Metrics')
ws3.sheet_view.showGridLines = False

ws3['A1'] = 'Dollars per Trip and Units per Trip for each period by segment among qualified trackable IDs:'

ws3['A2'] = ''

if vmr_trend_yago_summary.loc[0, 'Analysis Periods'] == 'YAGO Period':

    ws3['A3'] = 'Comparison VMR period vs YAGO period:'

    for r in dataframe_to_rows(comparison_yago, index=False, header=True):
        ws3.append(r)

    for cell in ws3[4]:
        cell.style = 'Accent1'
        cell.alignment=Alignment(wrap_text=True, horizontal="center")    

    for i in (range(5,5+len(comparison_yago))):
        for j in ('B', 'C'):
            ws3[j+str(i)].style = format_dollars_two_decimals

    for i in (range(5,5+len(comparison_yago))):
        for j in ('D', 'G'):
            ws3[j+str(i)].style = format_percent_no_decimal

    for i in (range(5,5+len(comparison_yago))):
        for j in ('E', 'F'):
            ws3[j+str(i)].style = format_comma_one_decimal


    ################################################################################################    

    next_row = 4+len(comparison_yago)+1

    ws3['A'+str(next_row)] = 'Comparison VMR period vs Pre-period period:'

    for r in dataframe_to_rows(comparison_preperiod, index=False, header=True):
        ws3.append(r)

    for cell in ws3[next_row+1]:
        cell.style = 'Accent1'
        cell.alignment=Alignment(wrap_text=True, horizontal="center")   

    for i in (range(next_row+1,len(comparison_preperiod)+next_row+3)):
        for j in ('B', 'C'):
            ws3[j+str(i)].style = format_dollars_two_decimals

    for i in (range(next_row+1,len(comparison_preperiod)+next_row+3)):
        for j in ('D', 'G'):
            ws3[j+str(i)].style = format_percent_no_decimal

    for i in (range(next_row+1,len(comparison_preperiod)+next_row+3)):
        for j in ('E', 'F'):
            ws3[j+str(i)].style = format_comma_one_decimal
            
    for cell in ws3[next_row+1]:
        cell.style = 'Accent1'
        cell.alignment=Alignment(wrap_text=True, horizontal="center")   
    

            
else:
    
    ws3['A3'] = 'Comparison VMR period vs Pre-period period:'

    for r in dataframe_to_rows(comparison_preperiod, index=False, header=True):
        ws3.append(r)

    for cell in ws3[4]:
        cell.style = 'Accent1'
        cell.alignment=Alignment(wrap_text=True, horizontal="center")   

    for i in (range(5,5+len(comparison_preperiod))):
        for j in ('B', 'C'):
            ws3[j+str(i)].style = format_dollars_two_decimals

    for i in (range(5,5+len(comparison_preperiod))):
        for j in ('D', 'G'):
            ws3[j+str(i)].style = format_percent_no_decimal

    for i in (range(5,5+len(comparison_preperiod))):
        for j in ('E', 'F'):
            ws3[j+str(i)].style = format_comma_one_decimal

    
    
    ################################################################################################


for cell in ws3[1]:
    cell.font = Font(bold=True, color = '4F81BD', size = 15, name = 'Calibri' )

    
ws3.column_dimensions['A'].width = 40        

for col_name in ['B','C']:
    for cols in ws3[col_name]:
        ws3.column_dimensions[col_name].width = 23       

for col_name in ['D','E','F','G']:
    for cols in ws3[col_name]:
        ws3.column_dimensions[col_name].width = 18 
'''
    
##### WORKSHEET 4 ####

print("worksheet8")

ws4 = wb.create_sheet(title='Pre-Period Profile')
ws4.sheet_view.showGridLines = False

ws4['A1'] = 'Count of New Brand Buyers:'

ws4['B1'] = brand_consump.loc[0, '% Brand IDs'] * total_reward_ids
ws4['B1'].style = format_comma_no_decimal
ws4['B1'].alignment = Alignment(wrap_text=True, horizontal="left" )
ws4['B1'].font = Font(bold=True, color = '4F81BD', size = 15, name = 'Calibri' )

ws4['A2'] = '**New Brand Buyer Count is projected off total reward IDs'
ws4['A3'] = ''

ws4['A4'] = 'Pre-52 week brand and category profile among qualifying IDs'
ws4['A5'] = '**Profile groups based on DOLLAR spend'
ws4['A6'] = ''


for r in dataframe_to_rows(brand_consump[['Brand Group', '% Brand IDs']], index=False, header=True):
    ws4.append(r)

for cell in ws4[7]:
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    
    

for i in (range(8,14)):
    for j in ('B'):
        ws4[j+str(i)].style = format_percent_no_decimal

        
###################################################
        
if category_show_parameter == True:


    ws4['A14'] = ''

    for r in dataframe_to_rows(cat_consump[['Category Group', '% Category IDs']], index=False, header=True):
        ws4.append(r)

    for cell in ws4[15]:
        cell.style = 'Accent1'
        cell.alignment=Alignment(wrap_text=True, horizontal="center")    


    for i in (range(16,22)):
        for j in ('B'):
            ws4[j+str(i)].style = format_percent_no_decimal


###################################################################
        
ws4['A1'].font = Font(bold=True, color = '4F81BD', size = 15, name = 'Calibri' )
ws4['A4'].font = Font(bold=True, color = '4F81BD', size = 15, name = 'Calibri' )
    
ws4.column_dimensions['A'].width = 40        

for col_name in ['B','C']:
    for cols in ws4[col_name]:
        ws4.column_dimensions[col_name].width = 23       

for col_name in ['D','E','F','G']:
    for cols in ws4[col_name]:
        ws4.column_dimensions[col_name].width = 18 
    
    
##### WORKSHEET 2 ####

print("worksheet9")

ws12 = wb.create_sheet(title='Repeat Rate')
ws12.sheet_view.showGridLines = False

ws12['A1'] = 'Repeat Rate'

ws12['A2'] = 'Metrics based on qualifying IDs'


df_repeat_rate = {
    
    '': ['Total Participants', 'New Buyers Participants', 'Existing Buyers Participants'],
    'Repeat Rate': [pct_repurch, pct_repurch_new_buyers, pct_repurch_existing_buyers]
    
}

df_repeat_rate = pd.DataFrame(df_repeat_rate)


for r in dataframe_to_rows(df_repeat_rate, index=False, header=True):
    ws12.append(r)

for cell in ws12[3]:
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    
    

for cell in ws12[1]:
    cell.font = Font(bold=True, color = '4F81BD', size = 15, name = 'Calibri' )

    
ws12.column_dimensions['A'].width = 40        


for i in (range(4,len(df_repeat_rate)+4)):
    for j in ('B'):
        ws12[j+str(i)].style = format_percent_no_decimal
        
        
        
##### WORKSHEET 2 ####

print("worksheet10")

ws13 = wb.create_sheet(title='Trends by Segment')
ws13.sheet_view.showGridLines = False

ws13['A1'] = 'Trends by Segment'

ws13['A2'] = 'Metrics based on qualifying IDs'

ws13['A3'] = ''

ws13['A4'] = 'Campaign vs Prior 52 wks'


#first table: Prior 52 wk

df_trend_by_segment_1 = table_slide_1
df_trend_by_segment_1.columns = df_trend_by_segment_1.iloc[0]
df_trend_by_segment_1 = df_trend_by_segment_1[1:].reset_index(drop=True)


for r in dataframe_to_rows(df_trend_by_segment_1, index=False, header=True):
    ws13.append(r)

for cell in ws13[5]:
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    
    
    
for cell in ws13[1]:
    cell.font = Font(bold=True, color = '4F81BD', size = 15, name = 'Calibri' )

ws13.column_dimensions['A'].width = 40        


#formatting first table

for i in (range(6,len(df_trend_by_segment_1)+7)):
    for j in ('B'):
        ws13[j+str(i)].style = format_dollars_two_decimals
                
for i in (range(6,len(df_trend_by_segment_1)+7)):
    for j in ('C'):
        ws13[j+str(i)].style = format_dollars_two_decimals
                
for i in (range(6,len(df_trend_by_segment_1)+7)):
    for j in ('D'):
        ws13[j+str(i)].style = format_percent_no_decimal

for i in (range(6,len(df_trend_by_segment_1)+7)):
    for j in ('E'):
        ws13[j+str(i)].style = format_comma_one_decimal
        
for i in (range(6,len(df_trend_by_segment_1)+7)):
    for j in ('F'):
        ws13[j+str(i)].style = format_comma_one_decimal
        
for i in (range(6,len(df_trend_by_segment_1)+7)):
    for j in ('G'):
        ws13[j+str(i)].style = format_percent_no_decimal

        
#------------------------------------------------------------------------

next_cell = len(df_trend_by_segment_1)+7+1

ws13['A'+str(next_cell)] = 'Campaign vs Pre Period'


#Second table: Pre Period

df_trend_by_segment_2 = table_slide_2
df_trend_by_segment_2.columns = df_trend_by_segment_2.iloc[0]
df_trend_by_segment_2 = df_trend_by_segment_2[1:].reset_index(drop=True)


for r in dataframe_to_rows(df_trend_by_segment_2, index=False, header=True):
    ws13.append(r)

for cell in ws13[next_cell+1]:
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    

    

#formatting second table

for i in (range(next_cell+2,len(df_trend_by_segment_2)+next_cell+3)):
    for j in ('B'):
        ws13[j+str(i)].style = format_dollars_two_decimals
        
for i in (range(next_cell+2,len(df_trend_by_segment_2)+next_cell+3)):
    for j in ('C'):
        ws13[j+str(i)].style = format_dollars_two_decimals
                
for i in (range(next_cell+2,len(df_trend_by_segment_2)+next_cell+3)):
    for j in ('D'):
        ws13[j+str(i)].style = format_percent_no_decimal
                
for i in (range(next_cell+2,len(df_trend_by_segment_2)+next_cell+3)):
    for j in ('E'):
        ws13[j+str(i)].style = format_comma_one_decimal
        
for i in (range(next_cell+2,len(df_trend_by_segment_2)+next_cell+3)):
    for j in ('F'):
        ws13[j+str(i)].style = format_comma_one_decimal

for i in (range(next_cell+2,len(df_trend_by_segment_2)+next_cell+3)):
    for j in ('G'):
        ws13[j+str(i)].style = format_percent_no_decimal
        
        
        
#------------------------------------------------------------------------

#Third table: YAGO Period


next_cell_b = len(df_trend_by_segment_2)+next_cell+3+1

ws13['A'+str(next_cell_b)] = 'Campaign vs YAGO'


df_trend_by_segment = table_slide
df_trend_by_segment.columns = df_trend_by_segment.iloc[0]
df_trend_by_segment = df_trend_by_segment[1:].reset_index(drop=True)


for r in dataframe_to_rows(df_trend_by_segment, index=False, header=True):
    ws13.append(r)

for cell in ws13[next_cell_b+1]:
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    

        
#formatting third table

for i in (range(next_cell_b+2,len(df_trend_by_segment)+next_cell_b+4)):
    for j in ('B'):
        ws13[j+str(i)].style = format_dollars_two_decimals
        
for i in (range(next_cell_b+2,len(df_trend_by_segment)+next_cell_b+4)):
    for j in ('C'):
        ws13[j+str(i)].style = format_dollars_two_decimals
                
for i in (range(next_cell_b+2,len(df_trend_by_segment)+next_cell_b+4)):
    for j in ('D'):
        ws13[j+str(i)].style = format_percent_no_decimal
                
for i in (range(next_cell_b+2,len(df_trend_by_segment)+next_cell_b+4)):
    for j in ('E'):
        ws13[j+str(i)].style = format_comma_one_decimal
        
for i in (range(next_cell_b+2,len(df_trend_by_segment)+next_cell_b+4)):
    for j in ('F'):
        ws13[j+str(i)].style = format_comma_one_decimal

for i in (range(next_cell_b+2,len(df_trend_by_segment)+next_cell_b+4)):
    for j in ('G'):
        ws13[j+str(i)].style = format_percent_no_decimal
        
        
##### WORKSHEET 5 ####

print("worksheet11")

if nbr_segments > 1 :

    ws5 = wb.create_sheet(title='Trip Segment Combo')
    ws5.sheet_view.showGridLines = False

    ws5['A1'] = 'Brand combinations purchased on the VMR qualified Trip among qualifying IDs:'

    ws5['A2'] = ''

    ws5['A3'] = '% of Trips including 2 or more segments:'
    ws5['A3'].font = Font(bold=True, color = '000000', size = 11, name = 'Calibri' )
    

    nbr_trips_two_segments = (df_vmr_period_combos_all.loc[(df_vmr_period_combos_all['No. of Segments']>=2), ['Trips']].sum())
    nbr_total_trips = (df_vmr_period_combos_all['Trips'].sum())

    ws5['B3'] = float(nbr_trips_two_segments/nbr_total_trips)
    ws5['B3'].style = format_percent_no_decimal
    ws5['B3'].alignment = Alignment(wrap_text=True, horizontal="left" )

    ws5['A4'] = ''

    for r in dataframe_to_rows(df_vmr_period_combos_all_out[["Segment", "Pct Of Trips", "No. of Segments"]], index=False, header=True):
        ws5.append(r)

    for cell in ws5[5]:
        cell.style = 'Accent1'
        cell.alignment=Alignment(wrap_text=True, horizontal="center")    


    for i in (range(6,len(df_vmr_period_combos_all_out)+6)):
        for j in ('B'):
            ws5[j+str(i)].style = format_percent_no_decimal


    for cell in ws5[1]:
        cell.font = Font(bold=True, color = '4F81BD', size = 15, name = 'Calibri' )


    ws5.column_dimensions['A'].width = 40        

    for col_name in ['B','C']:
        for cols in ws5[col_name]:
            ws5.column_dimensions[col_name].width = 23       

    for col_name in ['D','E','F','G']:
        for cols in ws5[col_name]:
            ws5.column_dimensions[col_name].width = 18 

            
            
##### WORKSHEET 6 ####

print("worksheet12")

if nbr_segments > 1 :

    ws6 = wb.create_sheet(title='New-Existing Segments')
    ws6.sheet_view.showGridLines = False

    ws6['A1'] = 'Segment analysis of purchases on VMR qualifying trip among qualifying IDs'
    ws6['A2'] = ''


    for r in dataframe_to_rows(df_plus_one[['','Pct Of Shoppers']], index=False, header=True):
        ws6.append(r)

    for cell in ws6[3]:
        cell.style = 'Accent1'
        cell.alignment=Alignment(wrap_text=True, horizontal="center")    


    for i in (range(4,6)):
        for j in ('B'):
            ws6[j+str(i)].style = format_percent_no_decimal



    #################################################################
    ws6['A6'] = ''

    for r in dataframe_to_rows(df_new_existing_combos[['','Pct Of Shoppers']], index=False, header=True):
        ws6.append(r)

    for cell in ws6[7]:
        cell.style = 'Accent1'
        cell.alignment=Alignment(wrap_text=True, horizontal="center")    


    for i in (range(8,11)):
        for j in ('B'):
            ws6[j+str(i)].style = format_percent_no_decimal


    ###################################################################

    for cell in ws6[1]:
        cell.font = Font(bold=True, color = '4F81BD', size = 15, name = 'Calibri' )


    ws6.column_dimensions['A'].width = 40        

    for col_name in ['B','C']:
        for cols in ws6[col_name]:
            ws6.column_dimensions[col_name].width = 23       

    for col_name in ['D','E','F','G']:
        for cols in ws6[col_name]:
            ws6.column_dimensions[col_name].width = 18 

            
##### WORKSHEET 7 ####

print("worksheet13")

ws7 = wb.create_sheet(title='Baskets Sizes')
ws7.sheet_view.showGridLines = False

ws7['A1'] = 'Baskets sizes during VMR and redemption period (All trips - No ID Filter):'

ws7['A2'] = ''

for r in dataframe_to_rows(df_vmr_ttls_basket[['Analysis period','Average Basket Size']], index=False, header=True):
    ws7.append(r)

for cell in ws7[3]:
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    

for i in (range(4,6)):
    for j in ('B'):
        ws7[j+str(i)].style = format_dollars_two_decimals

for i in (range(7,9)):
    for j in ('B'):
        ws7[j+str(i)].style = format_dollars_two_decimals


ws7['B6'].style = format_percent_no_decimal
ws7['B9'].style = format_percent_no_decimal

####################################################

ws7['A10'] = ''


if nbr_redemp_promoted['trips'].values[0] != 0:

    ws7['A11'] = '% redemption trips including promoted product'
    ws7['A11'].font = Font(bold=True, color = '000000', size = 11, name = 'Calibri' )
    ws7['B11'] = str(round((nbr_redemp_promoted['trips'].values[0]/nbr_redemp_trips['trips'].values[0])*100,0)) + '%'
    ws7['B11'].alignment=Alignment(wrap_text=True, horizontal="right")    



for cell in ws7[1]:
    cell.font = Font(bold=True, color = '4F81BD', size = 15, name = 'Calibri' )


ws7.column_dimensions['A'].width = 41.09        

for col_name in ['B','C']:
    for cols in ws7[col_name]:
        ws7.column_dimensions[col_name].width = 23       

for col_name in ['D','E','F','G']:
    for cols in ws7[col_name]:
        ws7.column_dimensions[col_name].width = 18 

        
        
##### WORKSHEET 2 ####

print("worksheet14")

ws2 = wb.create_sheet(title='Sales by Segment')
ws2.sheet_view.showGridLines = False

ws2['A1'] = 'Total Dollars and Units purchased on qualifying trips  (All trips - No ID filter):'

ws2['A2'] = ''

for r in dataframe_to_rows(total_vmr_details_by_brand, index=False, header=True):
    ws2.append(r)

for cell in ws2[3]:
    cell.style = 'Accent1'
    cell.alignment=Alignment(wrap_text=True, horizontal="center")    
    

for cell in ws2[1]:
    cell.font = Font(bold=True, color = '4F81BD', size = 15, name = 'Calibri' )

    
ws2.column_dimensions['A'].width = 40        

for col_name in ['B','C']:
    for cols in ws2[col_name]:
        ws2.column_dimensions[col_name].width = 23       

for col_name in ['D','E','F','G']:
    for cols in ws2[col_name]:
        ws2.column_dimensions[col_name].width = 18 
    

    
for i in (range(4,len(total_vmr_details_by_brand)+4)):
    for j in ('B'):
        ws2[j+str(i)].style = format_comma_no_decimal
        
        
for i in (range(4,len(total_vmr_details_by_brand)+4)):
    for j in ('C'):
        ws2[j+str(i)].style = format_dollars_no_decimals

        
for i in (range(4,len(total_vmr_details_by_brand)+4)):
    for j in ('D','E'):
        ws2[j+str(i)].style = format_percent_no_decimal

        
   
print('Starting Power Point')

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_TICK_MARK
from pptx.chart.data import CategoryChartData
from pptx.util import Inches
from pptx.enum.dml import MSO_LINE


# In[111]:


#Loading Power Point template to build the final deck

prs = Presentation("/opt/airflow/src/templates/VMR_Scorecard_Template.pptx")


# In[112]:


#Adding Slide 0: Catalina Front Page

sld_0 = prs.slides.add_slide(prs.slide_layouts[0])


# In[113]:


#Adding Slide 1: General Metrics

sld_1 = prs.slides.add_slide(prs.slide_layouts[1])

sld_1.placeholders[10].text = f'{str(reward_print_start).replace("-","/")} - {str(reward_print_end).replace("-","/")}'


if nbr_retailers < 3:
    sld_1.placeholders[11].text = f'{retailers}'
else:
    sld_1.placeholders[11].text = f'National'

new_brand_buyers_proj = brand_consump.loc[0, '% Brand IDs'] * total_reward_ids
    
sld_1.placeholders[13].text = f'${"{:,}".format(round(dollars_moved_vmr))}'

sld_1.placeholders[14].text = f'{"{:.0%}".format(change_dpt)}'

sld_1.placeholders[15].text = f'{"{:,}".format(round(new_brand_buyers_proj))}'

sld_1.placeholders[16].text = f'{"{:.0%}".format(pct_repurch)}'

sld_1.placeholders[12].text = f'Drive sales of Promoted Brand with a spend {min_threshold_statement} threshold next shopping trip offer'


# In[114]:

if vmr_trend_yago_summary.loc[0, 'Analysis Periods'] == 'YAGO Period':


    #Adding Slide 2: VMR Campaign Trends DOLLARS WITH VERTICAL LINES 

    sld_2 = prs.slides.add_slide(prs.slide_layouts[2])

    #Chart 2.1: Trend

    # X and Y values for Current Dollars:

    reward_dates = [df_promo_summary['Actual Start'].min(), df_promo_summary['Actual Start'].min(), df_promo_summary['Actual End'].max(), df_promo_summary['Actual End'].max()]
    dates_trend  = list(vmr_trend_brand_sl[vmr_trend_brand_sl['trend_chart'] == 'Recent Trend Period'].sort_values(by=['cal_sun_wk_ending_dt'], ascending = True).cal_sun_wk_ending_dt)
    curr_dol_val = list(vmr_trend_brand_sl[vmr_trend_brand_sl['trend_chart'] == 'Recent Trend Period'].sort_values(by=['cal_sun_wk_ending_dt'], ascending = True).dollar_sales.values)
    yago_dol_val = vmr_trend_brand_sl[vmr_trend_brand_sl['trend_chart'] == 'YAGO Trend Period'].sort_values(by=['cal_sun_wk_ending_dt'], ascending = True).dollar_sales

    print(len(dates_trend))
    print(len(curr_dol_val))
    print(len(yago_dol_val))

    if len(list(yago_dol_val)) != len(curr_dol_val):
        if len(list(yago_dol_val)) > len(curr_dol_val):
            yago_dol_val = list(yago_dol_val.head(len(curr_dol_val)).values)

    
    #Creating Df with all the dates
    df_trend = pd.DataFrame({'cal_sun_wk_ending_dt': dates_trend, 'dollar_sales': curr_dol_val, 'dollar_sales_yago': yago_dol_val})
    #print(df_trend)

    #Adding reward period dates
    reward_dates = [df_promo_summary['Actual Start'].min(), df_promo_summary['Actual Start'].min(), df_promo_summary['Actual End'].max(), df_promo_summary['Actual End'].max()] 
    new_data = pd.DataFrame({'cal_sun_wk_ending_dt': reward_dates, 'dollar_sales': ['#N/A','#N/A','#N/A','#N/A'], 'dollar_sales_yago': ['#N/A','#N/A','#N/A','#N/A'] })

    # Insert the new line at the end of the DataFrame
    
    data_frames = [df_trend, new_data]
    df_trend = pd.concat(data_frames, ignore_index=True)

    df_trend = df_trend.sort_values(by=['cal_sun_wk_ending_dt'], ascending = True)
    df_trend = df_trend.reset_index(drop=True)
    #print(df_trend)



    #VERTICAL ONE LINE

    vertical_one = []
    counter = 0
    for i in range(len(df_trend['cal_sun_wk_ending_dt'])):
        if df_trend['cal_sun_wk_ending_dt'][i] == df_promo_summary['Actual Start'].min():
            if counter == 0:
                vertical_one.append(0)
                counter += 1
            else:
                vertical_one.append(max(curr_dol_val) + max(curr_dol_val)/4)
                counter = 0
        else:
            vertical_one.append('')

    #print(vertical_one)


    #VERTICAL SECOND LINE

    vertical_second = []
    counter = 0
    for i in range(len(df_trend['cal_sun_wk_ending_dt'])):
        if df_trend['cal_sun_wk_ending_dt'][i] == df_promo_summary['Actual End'].max():
            if counter == 0:
                vertical_second.append(0)
                counter += 1
            else:
                vertical_second.append(max(curr_dol_val) + max(curr_dol_val)/4)
                counter = 0
        else:
            vertical_second.append('')

    #print(vertical_second)



    chart_data = CategoryChartData()
    chart_data.categories = df_trend['cal_sun_wk_ending_dt']
    chart_data.add_series('YAGO Period', df_trend['dollar_sales_yago'], '$#,##0')
    chart_data.add_series('VMR Period', df_trend['dollar_sales'], '$#,##0')
    chart_data.add_series('Campaign Start', vertical_one, '$#,##0')
    chart_data.add_series('Campaign End', vertical_second, '$#,##0')


    chart2_1 = sld_2.placeholders[10].insert_chart(XL_CHART_TYPE.LINE, chart_data).chart


    chart2_1.font.size  = Pt(12)
    chart2_1.value_axis.has_major_gridlines = False
    chart2_1.has_legend = True
    chart2_1.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart2_1.legend.include_in_layout = False
    chart2_1.value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE

    widths_sizes = [0.04,0.04,0.02,0.02]

    series = chart2_1.series
    for i, serie in enumerate(chart2_1.series):
        line_format = serie.format.line
        line_format.width = Inches(widths_sizes[i])  

        if i!=0 and i!=1:
            line_format.dash_style = MSO_LINE.SQUARE_DOT


    #Chart 2.2: Bar Chart 1

    chart_data = ChartData()
    chart_data.categories = vmr_trend_yago_summary['Analysis Periods']
    chart_data.add_series('', vmr_trend_yago_summary['Dollars per Trip'], '$#,##0.00')


    chart2_2 = sld_2.placeholders[11].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
    chart2_2.font.size  = Pt(12)
    chart2_2.category_axis.tick_labels.font.bold = True
    chart2_2.value_axis.has_major_gridlines = False


    chart2_2.plots[0].has_data_labels = True
    data_labels_2 = chart2_2.plots[0].data_labels
    data_labels_2.font.size = Pt(14)
    data_labels_2.font.bold = True
    data_labels_2.number_format = '$#,##0.00'



    #Chart 2.3: Bar Chart 2

    pct_chg_dol_share = (share_vmr_dol - share_yago_dol)*100

    chart_data = ChartData()
    chart_data.categories = vmr_trend_yago_summary['Analysis Periods']
    chart_data.add_series('', [share_yago_dol, share_vmr_dol], '0.0%')


    chart2_3 = sld_2.placeholders[12].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
    chart2_3.font.size  = Pt(12)
    chart2_3.category_axis.tick_labels.font.bold = True
    chart2_3.value_axis.has_major_gridlines = False


    chart2_3.plots[0].has_data_labels = True
    data_labels_3 = chart2_3.plots[0].data_labels
    data_labels_3.font.size = Pt(14)
    data_labels_3.font.bold = True
    data_labels_3.number_format = '0.0%'


    sld_2.placeholders[13].text = f'VMR Period +{dollars_chg_yago}% chg vs YAGO'
    
    if dollars_per_trip_chg_yago > 0:    
        sld_2.placeholders[14].text = f'+{dollars_per_trip_chg_yago}%                       vs YAGO'
    else:
        sld_2.placeholders[14].text = f'{dollars_per_trip_chg_yago}%                       vs YAGO'
        
    if pct_chg_dol_share > 0:
        sld_2.placeholders[15].text = f'+{"{:.1f}".format(pct_chg_dol_share)} Pts vs YAGO' 
    else:
        sld_2.placeholders[15].text = f'{"{:.1f}".format(pct_chg_dol_share)} Pts vs YAGO' 
    
    sld_2.placeholders[16].text = f'Campaign Period: {str(reward_print_start).replace("-","/")} - {str(reward_print_end).replace("-","/")} \nYAGO Period: {str(yago_print_start).replace("-","/")} - {str(yago_print_end).replace("-","/")}' 

    phr = sld_2.placeholders[16]
    # Modify the font name and size
    font_name = 'Arial Black'
    font_size = Pt(8)
    # Access the text frame and text properties
    text_frame = phr.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = font_size


# In[115]:


#Double checking the content of the dollars trend chart

x = vmr_trend_brand_sl[vmr_trend_brand_sl['trend_chart'] == 'Recent Trend Period'].sort_values(by=['cal_sun_wk_ending_dt'], ascending = True).cal_sun_wk_ending_dt
y = vmr_trend_brand_sl[vmr_trend_brand_sl['trend_chart'] == 'YAGO Trend Period'].sort_values(by=['cal_sun_wk_ending_dt'], ascending = True).dollar_sales.values
z = vmr_trend_brand_sl[vmr_trend_brand_sl['trend_chart'] == 'Recent Trend Period'].sort_values(by=['cal_sun_wk_ending_dt'], ascending = True).dollar_sales.values

print("The amount of dates that will be plot is: " + str(len(x)))
print("The number of dollar points that will be plot is: " + str(len(y)))
print("The number of dollar points from YAGO that will be plot is: " + str(len(z)))



if vmr_trend_yago_summary.loc[0, 'Analysis Periods'] != 'YAGO Period':

    #Adding Slide 2: VMR Campaign Trends DOLLARS WITH VERTICAL LINES 

    sld_2 = prs.slides.add_slide(prs.slide_layouts[2])

    #Chart 2.1: Trend

    # X and Y values for Current Dollars:

    reward_dates = [df_promo_summary['Actual Start'].min(), df_promo_summary['Actual Start'].min(), df_promo_summary['Actual End'].max(), df_promo_summary['Actual End'].max()]
    dates_trend  = list(vmr_trend_brand_sl[vmr_trend_brand_sl['trend_chart'] == 'Recent Trend Period'].sort_values(by=['cal_sun_wk_ending_dt'], ascending = True).cal_sun_wk_ending_dt)
    curr_dol_val = list(vmr_trend_brand_sl[vmr_trend_brand_sl['trend_chart'] == 'Recent Trend Period'].sort_values(by=['cal_sun_wk_ending_dt'], ascending = True).dollar_sales.values)
    yago_dol_val = vmr_trend_brand_sl[vmr_trend_brand_sl['trend_chart'] == 'YAGO Trend Period'].sort_values(by=['cal_sun_wk_ending_dt'], ascending = True).dollar_sales

    print(len(dates_trend))
    print(len(curr_dol_val))
    print(len(yago_dol_val))

    if len(list(yago_dol_val)) != len(curr_dol_val):
        if len(list(yago_dol_val)) > len(curr_dol_val):
            yago_dol_val = list(yago_dol_val.head(len(curr_dol_val)).values)


    #Creating Df with all the dates
    #df_trend = pd.DataFrame({'cal_sun_wk_ending_dt': dates_trend, 'dollar_sales': curr_dol_val, 'dollar_sales_yago': yago_dol_val})
    df_trend = pd.DataFrame({'cal_sun_wk_ending_dt': dates_trend, 'dollar_sales': curr_dol_val})

    #print(df_trend)

    #Adding reward period dates
    reward_dates = [df_promo_summary['Actual Start'].min(), df_promo_summary['Actual Start'].min(), df_promo_summary['Actual End'].max(), df_promo_summary['Actual End'].max()] 
    new_data = pd.DataFrame({'cal_sun_wk_ending_dt': reward_dates, 'dollar_sales': ['#N/A','#N/A','#N/A','#N/A'] })

    # Insert the new line at the end of the DataFrame
    data_frames = [df_trend, new_data]
    df_trend = pd.concat(data_frames, ignore_index=True)

    df_trend = df_trend.sort_values(by=['cal_sun_wk_ending_dt'], ascending = True)
    df_trend = df_trend.reset_index(drop=True)
    #print(df_trend)



    #VERTICAL ONE LINE

    vertical_one = []
    counter = 0
    for i in range(len(df_trend['cal_sun_wk_ending_dt'])):
        if df_trend['cal_sun_wk_ending_dt'][i] == df_promo_summary['Actual Start'].min():
            if counter == 0:
                vertical_one.append(0)
                counter += 1
            else:
                vertical_one.append(max(curr_dol_val) + max(curr_dol_val)/4)
                counter = 0
        else:
            vertical_one.append('')

    #print(vertical_one)


    #VERTICAL SECOND LINE

    vertical_second = []
    counter = 0
    for i in range(len(df_trend['cal_sun_wk_ending_dt'])):
        if df_trend['cal_sun_wk_ending_dt'][i] == df_promo_summary['Actual End'].max():
            if counter == 0:
                vertical_second.append(0)
                counter += 1
            else:
                vertical_second.append(max(curr_dol_val) + max(curr_dol_val)/4)
                counter = 0
        else:
            vertical_second.append('')

    #print(vertical_second)



    chart_data = CategoryChartData()
    chart_data.categories = df_trend['cal_sun_wk_ending_dt']
    #chart_data.add_series('YAGO Period', df_trend['dollar_sales_yago'], '$#,##0')
    chart_data.add_series('VMR Period', df_trend['dollar_sales'], '$#,##0')
    chart_data.add_series('Campaign Start', vertical_one, '$#,##0')
    chart_data.add_series('Campaign End', vertical_second, '$#,##0')


    chart2_1 = sld_2.placeholders[10].insert_chart(XL_CHART_TYPE.LINE, chart_data).chart


    chart2_1.font.size  = Pt(12)
    chart2_1.value_axis.has_major_gridlines = False
    chart2_1.has_legend = True
    chart2_1.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart2_1.legend.include_in_layout = False
    chart2_1.value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE

    widths_sizes = [0.04,0.04,0.02,0.02]

    series = chart2_1.series
    for i, serie in enumerate(chart2_1.series):
        line_format = serie.format.line
        line_format.width = Inches(widths_sizes[i])  

        if i!=0 and i!=1:
            line_format.dash_style = MSO_LINE.SQUARE_DOT


    #Chart 2.2: Bar Chart 1

    chart_data = ChartData()
    chart_data.categories = vmr_trend_yago_summary['Analysis Periods']
    chart_data.add_series('', vmr_trend_yago_summary['Dollars per Trip'], '$#,##0.00')


    chart2_2 = sld_2.placeholders[11].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
    chart2_2.font.size  = Pt(12)
    chart2_2.category_axis.tick_labels.font.bold = True
    chart2_2.value_axis.has_major_gridlines = False


    chart2_2.plots[0].has_data_labels = True
    data_labels_2 = chart2_2.plots[0].data_labels
    data_labels_2.font.size = Pt(14)
    data_labels_2.font.bold = True
    data_labels_2.number_format = '$#,##0.00'



    #Chart 2.3: Bar Chart 2

    pct_chg_dol_share = (share_vmr_dol - share_yago_dol)*100

    chart_data = ChartData()
    chart_data.categories = vmr_trend_yago_summary['Analysis Periods']
    chart_data.add_series('', [share_yago_dol, share_vmr_dol], '0.0%')


    chart2_3 = sld_2.placeholders[12].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
    chart2_3.font.size  = Pt(12)
    chart2_3.category_axis.tick_labels.font.bold = True
    chart2_3.value_axis.has_major_gridlines = False


    chart2_3.plots[0].has_data_labels = True
    data_labels_3 = chart2_3.plots[0].data_labels
    data_labels_3.font.size = Pt(14)
    data_labels_3.font.bold = True
    data_labels_3.number_format = '0.0%'


    sld_2.placeholders[13].text = f'VMR Period +{dollars_chg_yago}% chg vs YAGO'
    
    if dollars_per_trip_chg_yago > 0:    
        sld_2.placeholders[14].text = f'+{dollars_per_trip_chg_yago}%                       vs YAGO'
    else:
        sld_2.placeholders[14].text = f'{dollars_per_trip_chg_yago}%                       vs YAGO'
        
    if pct_chg_dol_share > 0:
        sld_2.placeholders[15].text = f'+{"{:.1f}".format(pct_chg_dol_share)} Pts vs YAGO' 
    else:
        sld_2.placeholders[15].text = f'{"{:.1f}".format(pct_chg_dol_share)} Pts vs YAGO' 
    
    sld_2.placeholders[16].text = f'Campaign Period: {str(reward_print_start).replace("-","/")} - {str(reward_print_end).replace("-","/")} \nYAGO Period: {str(yago_print_start).replace("-","/")} - {str(yago_print_end).replace("-","/")}' 

    phr = sld_2.placeholders[16]
    # Modify the font name and size
    font_name = 'Arial Black'
    font_size = Pt(8)
    # Access the text frame and text properties
    text_frame = phr.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = font_size

print('slide2')       


# In[116]:


if vmr_trend_yago_summary.loc[0, 'Analysis Periods'] == 'YAGO Period':

    #Adding Slide 3: VMR Campaign UNITS

    sld_3 = prs.slides.add_slide(prs.slide_layouts[11])

    #Chart 2.1: Trend

    # X and Y values for Current Dollars:

    reward_dates = [df_promo_summary['Actual Start'].min(), df_promo_summary['Actual Start'].min(), df_promo_summary['Actual End'].max(), df_promo_summary['Actual End'].max()]
    dates_trend  = list(vmr_trend_brand_sl[vmr_trend_brand_sl['trend_chart'] == 'Recent Trend Period'].sort_values(by=['cal_sun_wk_ending_dt'], ascending = True).cal_sun_wk_ending_dt)
    curr_dol_val = list(vmr_trend_brand_sl[vmr_trend_brand_sl['trend_chart'] == 'Recent Trend Period'].sort_values(by=['cal_sun_wk_ending_dt'], ascending = True).units_sales.values)
    yago_dol_val = vmr_trend_brand_sl[vmr_trend_brand_sl['trend_chart'] == 'YAGO Trend Period'].sort_values(by=['cal_sun_wk_ending_dt'], ascending = True).units_sales

    len(dates_trend)
    len(curr_dol_val)
    len(yago_dol_val)

    if len(yago_dol_val) != len(curr_dol_val):
        yago_dol_val = list(yago_dol_val.head(len(curr_dol_val)).values)

        #Creating Df with all the dates
    df_trend = pd.DataFrame({'cal_sun_wk_ending_dt': dates_trend, 'units_sales': curr_dol_val, 'units_sales_yago': yago_dol_val})
    #print(df_trend)

    #Adding reward period dates
    reward_dates = [df_promo_summary['Actual Start'].min(), df_promo_summary['Actual Start'].min(), df_promo_summary['Actual End'].max(), df_promo_summary['Actual End'].max()] 
    new_data = pd.DataFrame({'cal_sun_wk_ending_dt': reward_dates, 'units_sales': ['#N/A','#N/A','#N/A','#N/A'], 'units_sales_yago': ['#N/A','#N/A','#N/A','#N/A'] })

    # Insert the new line at the end of the DataFrame
    data_frames = [df_trend, new_data]
    df_trend = pd.concat(data_frames, ignore_index=True)


    df_trend = df_trend.sort_values(by=['cal_sun_wk_ending_dt'], ascending = True)
    df_trend = df_trend.reset_index(drop=True)
    #print(df_trend)



    #VERTICAL ONE LINE

    vertical_one = []
    counter = 0
    for i in range(len(df_trend['cal_sun_wk_ending_dt'])):
        if df_trend['cal_sun_wk_ending_dt'][i] == df_promo_summary['Actual Start'].min():
            if counter == 0:
                vertical_one.append(0)
                counter += 1
            else:
                vertical_one.append(max(curr_dol_val) + max(curr_dol_val)/4)
                counter = 0
        else:
            vertical_one.append('')

    #print(vertical_one)


    #VERTICAL SECOND LINE

    vertical_second = []
    counter = 0
    for i in range(len(df_trend['cal_sun_wk_ending_dt'])):
        if df_trend['cal_sun_wk_ending_dt'][i] == df_promo_summary['Actual End'].max():
            if counter == 0:
                vertical_second.append(0)
                counter += 1
            else:
                vertical_second.append(max(curr_dol_val) + max(curr_dol_val)/4)
                counter = 0
        else:
            vertical_second.append('')

    #print(vertical_second)



    chart_data = CategoryChartData()
    chart_data.categories = df_trend['cal_sun_wk_ending_dt']
    chart_data.add_series('YAGO Period', df_trend['units_sales_yago'], '#,##0')
    chart_data.add_series('VMR Period', df_trend['units_sales'], '#,##0')
    chart_data.add_series('Campaign Start', vertical_one, '#,##0')
    chart_data.add_series('Campaign End', vertical_second, '#,##0')


    chart3_1 = sld_3.placeholders[10].insert_chart(XL_CHART_TYPE.LINE, chart_data).chart


    chart3_1.font.size  = Pt(12)
    chart3_1.value_axis.has_major_gridlines = False
    chart3_1.has_legend = True
    chart3_1.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart3_1.legend.include_in_layout = False
    chart3_1.value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE

    widths_sizes = [0.04,0.04,0.02,0.02]

    series = chart3_1.series
    for i, serie in enumerate(chart3_1.series):
        line_format = serie.format.line
        line_format.width = Inches(widths_sizes[i])  

        if i!=0 and i!=1:
            line_format.dash_style = MSO_LINE.SQUARE_DOT


    #Chart 3.2: Bar Chart 1

    chart_data = ChartData()
    chart_data.categories = vmr_trend_yago_summary['Analysis Periods']
    chart_data.add_series('', vmr_trend_yago_summary['Units per Trip'], '#,##0.0')


    chart3_2 = sld_3.placeholders[11].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
    chart3_2.font.size  = Pt(12)
    chart3_2.category_axis.tick_labels.font.bold = True
    chart3_2.value_axis.has_major_gridlines = False


    chart3_2.plots[0].has_data_labels = True
    data_labels_2 = chart3_2.plots[0].data_labels
    data_labels_2.font.size = Pt(14)
    data_labels_2.font.bold = True
    data_labels_2.number_format = '#,##0.0'



    #Chart 3.3: Bar Chart 2


    pct_chg_unit_share = (share_vmr_unit - share_yago_unit)*100

    chart_data = ChartData()
    chart_data.categories = vmr_trend_yago_summary['Analysis Periods']
    chart_data.add_series('', [share_yago_unit, share_vmr_unit], '0.0%')


    chart3_3 = sld_3.placeholders[12].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
    chart3_3.font.size  = Pt(12)
    chart3_3.category_axis.tick_labels.font.bold = True
    chart3_3.value_axis.has_major_gridlines = False


    chart3_3.plots[0].has_data_labels = True
    data_labels_3 = chart3_3.plots[0].data_labels
    data_labels_3.font.size = Pt(14)
    data_labels_3.font.bold = True
    data_labels_3.number_format = '0.0%'


    sld_3.placeholders[13].text = f'VMR Period +{units_chg_yago}% chg vs YAGO'
    
    if units_per_trip_chg_yago > 0:
        sld_3.placeholders[14].text = f'+{units_per_trip_chg_yago}%                     vs YAGO'
    else:
        sld_3.placeholders[14].text = f'{units_per_trip_chg_yago}%                     vs YAGO'
        
    if pct_chg_unit_share > 0:
        sld_3.placeholders[15].text = f'+{"{:.1f}".format(pct_chg_unit_share)} Pts vs YAGO' 
    else:
        sld_3.placeholders[15].text = f'{"{:.1f}".format(pct_chg_unit_share)} Pts vs YAGO' 
    
    sld_3.placeholders[16].text = f'Campaign Period: {str(reward_print_start).replace("-","/")} - {str(reward_print_end).replace("-","/")} \nYAGO Period: {str(yago_print_start).replace("-","/")} - {str(yago_print_end).replace("-","/")}' 

    phr = sld_3.placeholders[16]
    # Modify the font name and size
    font_name = 'Arial Black'
    font_size = Pt(8)
    # Access the text frame and text properties
    text_frame = phr.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = font_size



if vmr_trend_yago_summary.loc[0, 'Analysis Periods'] != 'YAGO Period':

    #Adding Slide 3: VMR Campaign UNITS

    sld_3 = prs.slides.add_slide(prs.slide_layouts[11])

    #Chart 2.1: Trend

    # X and Y values for Current Dollars:

    reward_dates = [df_promo_summary['Actual Start'].min(), df_promo_summary['Actual Start'].min(), df_promo_summary['Actual End'].max(), df_promo_summary['Actual End'].max()]
    dates_trend  = list(vmr_trend_brand_sl[vmr_trend_brand_sl['trend_chart'] == 'Recent Trend Period'].sort_values(by=['cal_sun_wk_ending_dt'], ascending = True).cal_sun_wk_ending_dt)
    curr_dol_val = list(vmr_trend_brand_sl[vmr_trend_brand_sl['trend_chart'] == 'Recent Trend Period'].sort_values(by=['cal_sun_wk_ending_dt'], ascending = True).units_sales.values)
    yago_dol_val = vmr_trend_brand_sl[vmr_trend_brand_sl['trend_chart'] == 'YAGO Trend Period'].sort_values(by=['cal_sun_wk_ending_dt'], ascending = True).units_sales

    len(dates_trend)
    len(curr_dol_val)
    len(yago_dol_val)

    if len(yago_dol_val) != len(curr_dol_val):
        if len(yago_dol_val) > len(curr_dol_val):
            yago_dol_val = list(yago_dol_val.head(len(curr_dol_val)).values)

        #Creating Df with all the dates
    df_trend = pd.DataFrame({'cal_sun_wk_ending_dt': dates_trend, 'units_sales': curr_dol_val})
    #print(df_trend)

    #Adding reward period dates
    reward_dates = [df_promo_summary['Actual Start'].min(), df_promo_summary['Actual Start'].min(), df_promo_summary['Actual End'].max(), df_promo_summary['Actual End'].max()] 
    new_data = pd.DataFrame({'cal_sun_wk_ending_dt': reward_dates, 'units_sales': ['#N/A','#N/A','#N/A','#N/A']})

    # Insert the new line at the end of the DataFrame
    data_frames = [df_trend, new_data]
    df_trend = pd.concat(data_frames, ignore_index=True)


    df_trend = df_trend.sort_values(by=['cal_sun_wk_ending_dt'], ascending = True)
    df_trend = df_trend.reset_index(drop=True)
    #print(df_trend)



    #VERTICAL ONE LINE

    vertical_one = []
    counter = 0
    for i in range(len(df_trend['cal_sun_wk_ending_dt'])):
        if df_trend['cal_sun_wk_ending_dt'][i] == df_promo_summary['Actual Start'].min():
            if counter == 0:
                vertical_one.append(0)
                counter += 1
            else:
                vertical_one.append(max(curr_dol_val) + max(curr_dol_val)/4)
                counter = 0
        else:
            vertical_one.append('')

    #print(vertical_one)


    #VERTICAL SECOND LINE

    vertical_second = []
    counter = 0
    for i in range(len(df_trend['cal_sun_wk_ending_dt'])):
        if df_trend['cal_sun_wk_ending_dt'][i] == df_promo_summary['Actual End'].max():
            if counter == 0:
                vertical_second.append(0)
                counter += 1
            else:
                vertical_second.append(max(curr_dol_val) + max(curr_dol_val)/4)
                counter = 0
        else:
            vertical_second.append('')

    #print(vertical_second)



    chart_data = CategoryChartData()
    chart_data.categories = df_trend['cal_sun_wk_ending_dt']
    #chart_data.add_series('YAGO Period', df_trend['units_sales_yago'], '#,##0')
    chart_data.add_series('VMR Period', df_trend['units_sales'], '#,##0')
    chart_data.add_series('Campaign Start', vertical_one, '#,##0')
    chart_data.add_series('Campaign End', vertical_second, '#,##0')


    chart3_1 = sld_3.placeholders[10].insert_chart(XL_CHART_TYPE.LINE, chart_data).chart


    chart3_1.font.size  = Pt(12)
    chart3_1.value_axis.has_major_gridlines = False
    chart3_1.has_legend = True
    chart3_1.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart3_1.legend.include_in_layout = False
    chart3_1.value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE

    widths_sizes = [0.04,0.04,0.02,0.02]

    series = chart3_1.series
    for i, serie in enumerate(chart3_1.series):
        line_format = serie.format.line
        line_format.width = Inches(widths_sizes[i])  

        if i!=0 and i!=1:
            line_format.dash_style = MSO_LINE.SQUARE_DOT


    #Chart 3.2: Bar Chart 1

    chart_data = ChartData()
    chart_data.categories = vmr_trend_yago_summary['Analysis Periods']
    chart_data.add_series('', vmr_trend_yago_summary['Units per Trip'], '#,##0.0')


    chart3_2 = sld_3.placeholders[11].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
    chart3_2.font.size  = Pt(12)
    chart3_2.category_axis.tick_labels.font.bold = True
    chart3_2.value_axis.has_major_gridlines = False


    chart3_2.plots[0].has_data_labels = True
    data_labels_2 = chart3_2.plots[0].data_labels
    data_labels_2.font.size = Pt(14)
    data_labels_2.font.bold = True
    data_labels_2.number_format = '#,##0.0'



    #Chart 3.3: Bar Chart 2


    pct_chg_unit_share = (share_vmr_unit - share_yago_unit)*100

    chart_data = ChartData()
    chart_data.categories = vmr_trend_yago_summary['Analysis Periods']
    chart_data.add_series('', [share_yago_unit, share_vmr_unit], '0.0%')


    chart3_3 = sld_3.placeholders[12].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
    chart3_3.font.size  = Pt(12)
    chart3_3.category_axis.tick_labels.font.bold = True
    chart3_3.value_axis.has_major_gridlines = False


    chart3_3.plots[0].has_data_labels = True
    data_labels_3 = chart3_3.plots[0].data_labels
    data_labels_3.font.size = Pt(14)
    data_labels_3.font.bold = True
    data_labels_3.number_format = '0.0%'


    sld_3.placeholders[13].text = f'VMR Period +{units_chg_yago}% chg vs YAGO'
    
    if units_per_trip_chg_yago > 0:
        sld_3.placeholders[14].text = f'+{units_per_trip_chg_yago}%                     vs YAGO'
    else:
        sld_3.placeholders[14].text = f'{units_per_trip_chg_yago}%                     vs YAGO'
        
    if pct_chg_unit_share > 0:
        sld_3.placeholders[15].text = f'+{"{:.1f}".format(pct_chg_unit_share)} Pts vs YAGO' 
    else:
        sld_3.placeholders[15].text = f'{"{:.1f}".format(pct_chg_unit_share)} Pts vs YAGO' 
    
    sld_3.placeholders[16].text = f'Campaign Period: {str(reward_print_start).replace("-","/")} - {str(reward_print_end).replace("-","/")} \nYAGO Period: {str(yago_print_start).replace("-","/")} - {str(yago_print_end).replace("-","/")}' 

    phr = sld_3.placeholders[16]
    # Modify the font name and size
    font_name = 'Arial Black'
    font_size = Pt(8)
    # Access the text frame and text properties
    text_frame = phr.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = font_size

            
            
print('slide3')       
# In[117]:


#Adding Slide 4: VMR Dollars Moved
sld_4 = prs.slides.add_slide(prs.slide_layouts[3])

sld_4.placeholders[10].text = f'VMR reward trips accounted for {round(pct_reward_trips*100)}% of total brand transactions during the VMR Campaign \n${"{:,}".format(round(dollars_moved_vmr))} total dollars moved among all campaign reward levels'
sld_4.placeholders[13].text = f'${"{:,}".format(round(dollars_moved_vmr))} TOTAL DOLLARS MOVED AT REWARD TRIP LEVELS'

#Chart 4.1: Pie Chart

chart_data = ChartData()
chart_data.categories = ['All Other Trips', 'Reward Level Trips']
chart_data.add_series('', (1-pct_reward_trips, pct_reward_trips), '0%')


chart4_1 = sld_4.placeholders[11].insert_chart(XL_CHART_TYPE.PIE, chart_data).chart
chart4_1.font.size  = Pt(12)
chart4_1.has_legend = True
chart4_1.legend.position = XL_LEGEND_POSITION.BOTTOM
chart4_1.legend.include_in_layout = False
chart4_1.has_title = False

leg = chart4_1.legend
if leg:
    leg.font.bold = True

chart4_1.plots[0].has_data_labels = True
data_labels_1 = chart4_1.plots[0].data_labels
data_labels_1.font.size = Pt(16)
data_labels_1.font.bold = True
data_labels_1.font.color.rgb = RGBColor(255, 255, 255)
data_labels_1.number_format = '0%'
data_labels_1.position = XL_LABEL_POSITION.CENTER


#Chart 4.2: Bar Chart

chart_data = ChartData()
chart_data.categories = level_ct.loc[(level_ct['level']!= 'All Other Transactions') & (level_ct['level']!= 'Grand Total')].level
chart_data.add_series('', level_ct.loc[(level_ct['level']!= 'All Other Transactions') & (level_ct['level']!= 'Grand Total')].dollars, '$#,##0')


chart4_2 = sld_4.placeholders[12].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
chart4_2.font.size  = Pt(11)
chart4_2.category_axis.tick_labels.font.bold = True
chart4_2.value_axis.has_major_gridlines = False


chart4_2.plots[0].has_data_labels = True
data_labels_1 = chart4_2.plots[0].data_labels
data_labels_1.font.size = Pt(14)
data_labels_1.font.bold = True
data_labels_1.number_format = '$#,##0'



bar_colors = [RGBColor(152, 162, 167), RGBColor(230, 103, 73), RGBColor(240, 192, 86), RGBColor(152, 162, 167), RGBColor(230, 103, 73), RGBColor(240, 192, 86), RGBColor(152, 162, 167), RGBColor(230, 103, 73), RGBColor(240, 192, 86) ]

# Iterate through the bars and set colors
for i, series in enumerate(chart4_2.series):
    for j, point in enumerate(series.points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = bar_colors[j]
        
        
sld_4.placeholders[16].text = f'Campaign Period: {str(reward_print_start).replace("-","/")} - {str(reward_print_end).replace("-","/")}' 

phr = sld_4.placeholders[16]
# Modify the font name and size
font_name = 'Arial Black'
font_size = Pt(10)
# Access the text frame and text properties
text_frame = phr.text_frame
for paragraph in text_frame.paragraphs:
    paragraph.font.name = font_name
    paragraph.font.size = font_size

print('slide4')

# In[118]:


#Adding Slide 5: CAMPAIGN TRENDS
sld_5 = prs.slides.add_slide(prs.slide_layouts[12])


#Table 5.1:  #11


tph1_1 = sld_5.placeholders[11].insert_table(rows = results_participants.shape[0], cols = results_participants.shape[1])
tph1_1._element.graphic.graphicData.tbl[0][-1].text = '{69012ECD-51FC-41F1-AA8D-1B2483CD663E}' 
table5_1 = tph1_1.table



#populate slide 1 table 1
for i in range(results_participants.shape[0]):
    for j in range(results_participants.shape[1]):
        cell = table5_1.cell(i, j)

        if (i!= 0 and i!=3) and (j == 1 or j== 4) :
            cell.text = "${:,.2f}".format(results_participants.iloc[i,j])
        elif (i!= 0 and i!=3) and (j == 2 or j==3 or j==5):
            cell.text = "{:,.1f}".format(results_participants.iloc[i,j])
        elif i==3 and (j==1 or j==2 or j==3 or j==4 or j==5):

            if results_participants.iloc[i,j] != '':
                cell.text = "{:.0%}".format(results_participants.iloc[i,j])
            else:
                cell.text = results_participants.iloc[i,j]
        else:
            cell.text = str(results_participants.iloc[i,j])

        cell.text_frame.paragraphs[0].font.size = Pt(14.5)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER



#Chart 5.1:   #12

chart_data = ChartData()
chart_data.categories = ['Dollars per Buyer', 'Units per Buyer', 'Trips per Buyer', 'Dollars per Trip', 'Units per Trip']

chart_data.add_series('', [change_dpb, change_upb, change_tpb, change_dpt, change_upt], '0%')


chart5_1 = sld_5.placeholders[12].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
chart5_1.font.size  = Pt(12)
chart5_1.category_axis.tick_labels.font.bold = True
chart5_1.has_legend = False

chart5_1.has_title = True
chart5_1.chart_title.text_frame.text = f"% CHANGE - Campaign Period vs YAGO Period"
chart_title = chart5_1.chart_title.text_frame
font = chart_title.paragraphs[0].runs[0].font
font.size = Pt(14) 
font.bold = True

chart5_1.value_axis.has_major_gridlines = False

chart5_1.plots[0].has_data_labels = True
data_labels_1 = chart5_1.plots[0].data_labels
data_labels_1.font.size = Pt(14)
data_labels_1.font.bold = True
data_labels_1.font.color.rgb = RGBColor(0, 32, 96)
data_labels_1.number_format = '0%'
data_labels_1.position = XL_LABEL_POSITION.OUTSIDE_END


#Applying conditional formatting to color the bars depending on the label dollars, units or trips

colors = [RGBColor(0, 56, 104),RGBColor(127, 169, 184),RGBColor(240, 192, 86),RGBColor(0, 56, 104),RGBColor(127, 169, 184)]


for index, series in enumerate(chart5_1.series):
    for point, color in zip(series.points, colors):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color
        
        

from pptx.oxml.xmlchemy import OxmlElement

# Iterate over each series and point in the chart
for serie in chart5_1.series:
    for point in serie.points:
        # Create the new XML element
        element = OxmlElement('c:invertIfNegative')
        # Set the attribute 'val' to '0'
        element.set('val', '0')
        # Append the new element to the point's format element
        point.format.element.append(element)


#Text placeholder 10

if change_dpb > 0:
    sld_5.placeholders[10].text = f'Campaign Participants avg dollars per buyer CHANGE +{"{:.0%}".format(change_dpb)} during the campaign period vs YAGO period'
else:
    sld_5.placeholders[10].text = f'Campaign Participants avg dollars per buyer CHANGE {"{:.0%}".format(change_dpb)} during the campaign period vs YAGO period'


sld_5.placeholders[16].text = f'Campaign Period: {str(reward_print_start).replace("-","/")} - {str(reward_print_end).replace("-","/")} \nYAGO Period: {str(yago_print_start).replace("-","/")} - {str(yago_print_end).replace("-","/")}' 

phr = sld_5.placeholders[16]
# Modify the font name and size
font_name = 'Arial Black'
font_size = Pt(8)
# Access the text frame and text properties
text_frame = phr.text_frame
for paragraph in text_frame.paragraphs:
    paragraph.font.name = font_name
    paragraph.font.size = font_size
print('slide5')

# In[119]:


#Adding Slide 6: CATEGORY SHARE
sld_6 = prs.slides.add_slide(prs.slide_layouts[13])

#Chart 6.1: Bar Chart

chart_data = ChartData()
chart_data.categories = ['YAGO Period', 'VMR Period', 'Post Period']
chart_data.add_series('', (share_dollar_yago, share_dollar_vmr, share_dollar_post_4wk), '0%')


chart6_1 = sld_6.placeholders[17].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
chart6_1.font.size  = Pt(12)
chart6_1.category_axis.tick_labels.font.bold = True
chart6_1.value_axis.has_major_gridlines = False
chart6_1.value_axis.visible = False


chart6_1.plots[0].has_data_labels = True
data_labels_1 = chart6_1.plots[0].data_labels
data_labels_1.font.size = Pt(14)
data_labels_1.font.bold = True
data_labels_1.number_format = '0%'


#Chart 6.2: Bar Chart

chart_data = ChartData()
chart_data.categories = ['YAGO Period', 'VMR Period', 'Post Period']
chart_data.add_series('', (share_units_yago, share_units_vmr, share_units_post_4wk), '0%')


chart6_2 = sld_6.placeholders[18].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
chart6_2.font.size  = Pt(12)
chart6_2.category_axis.tick_labels.font.bold = True
chart6_2.value_axis.has_major_gridlines = False
chart6_2.value_axis.visible = False


chart6_2.plots[0].has_data_labels = True
data_labels_1 = chart6_2.plots[0].data_labels
data_labels_1.font.size = Pt(14)
data_labels_1.font.bold = True
data_labels_1.number_format = '0%'


#Text

chg_cat_share_dol = ( share_dollar_vmr - share_dollar_yago ) *100
chg_cat_share_units = ( share_units_vmr - share_units_yago ) *100

if chg_cat_share_dol >0:
    sld_6.placeholders[16].text = f'Campaign Participants dollar share changed +{"{:.0f}".format(chg_cat_share_dol)} pts during the campaign period vs YAGO'
else:
    sld_6.placeholders[16].text = f'Campaign Participants dollar share changed {"{:.0f}".format(chg_cat_share_dol)} pts during the campaign period vs YAGO'
    
    
if chg_cat_share_dol > 0:
    sld_6.placeholders[14].text = f'VMR Period +{"{:.0f}".format(chg_cat_share_dol)} pts vs YAGO'
else:
    sld_6.placeholders[14].text = f'VMR Period {"{:.0f}".format(chg_cat_share_dol)} pts vs YAGO'
    
    
if chg_cat_share_units > 0:
    sld_6.placeholders[15].text = f'VMR Period +{"{:.0f}".format(chg_cat_share_units)} pts vs YAGO'
else:
    sld_6.placeholders[15].text = f'VMR Period {"{:.0f}".format(chg_cat_share_units)} pts vs YAGO'


sld_6.placeholders[19].text = f'Campaign Period: {str(reward_print_start).replace("-","/")} - {str(reward_print_end).replace("-","/")} \nYAGO Period: {str(yago_print_start).replace("-","/")} - {str(yago_print_end).replace("-","/")} \nPost Period: {str(post_4wk_start).replace("-","/")} - {str(post_4wk_end).replace("-","/")}' 


phr = sld_6.placeholders[19]
# Modify the font name and size
font_name = 'Arial Black'
font_size = Pt(8)
# Access the text frame and text properties
text_frame = phr.text_frame
for paragraph in text_frame.paragraphs:
    paragraph.font.name = font_name
    paragraph.font.size = font_size
    
print('slide6')

# In[120]:


#Adding Slide 7: PRE PERIOD PROFILE AND REPEAT
sld_7 = prs.slides.add_slide(prs.slide_layouts[6])

#Chart 7.1: Pie Chart

chart_data = ChartData()
chart_data.categories = ['New Buyers', 'Existing Buyers']


pct_new_buy = round(brand_consump.loc[0, '% Brand IDs'],2)

chart_data.add_series('', (pct_new_buy, 1- pct_new_buy), '0%')


chart7_1 = sld_7.placeholders[11].insert_chart(XL_CHART_TYPE.PIE, chart_data).chart
chart7_1.font.size  = Pt(12)
chart7_1.has_legend = True
chart7_1.legend.position = XL_LEGEND_POSITION.BOTTOM
chart7_1.legend.include_in_layout = False
chart7_1.has_title = False

leg = chart7_1.legend
if leg:
    leg.font.bold = True


chart7_1.plots[0].has_data_labels = True
data_labels_1 = chart7_1.plots[0].data_labels
data_labels_1.font.size = Pt(16)
data_labels_1.font.bold = True
data_labels_1.font.color.rgb = RGBColor(255, 255, 255)
data_labels_1.number_format = '0%'
data_labels_1.position = XL_LABEL_POSITION.CENTER



#Chart 7.2: Bar Chart

chart_data = ChartData()
chart_data.categories = ['Total Participants', 'New Buyers Participants', 'Existing Buyers Participants']

chart_data.add_series('', (pct_repurch, pct_repurch_new_buyers, pct_repurch_existing_buyers), '0%')


chart7_2 = sld_7.placeholders[12].insert_chart(XL_CHART_TYPE.BAR_CLUSTERED, chart_data).chart
chart7_2.font.size  = Pt(12)
chart7_2.category_axis.tick_labels.font.bold = True
chart7_2.value_axis.has_major_gridlines = False


chart7_2.plots[0].has_data_labels = True
data_labels_1 = chart7_2.plots[0].data_labels
data_labels_1.font.size = Pt(14)
data_labels_1.font.bold = True
data_labels_1.number_format = '0%'


bar_colors = [RGBColor(240, 192, 86), RGBColor(0, 56, 104), RGBColor(127, 169, 184)]

# Iterate through the bars and set colors
for i, series in enumerate(chart7_2.series):
    for j, point in enumerate(series.points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = bar_colors[j]




#Text placeholder

sld_7.placeholders[10].text = f'{"{:.0%}".format(pct_new_buy)} of Campaign Participants were NEW to the brand with average repeat of {"{:.0%}".format(pct_repurch_new_buyers)} \n{"{:.0%}".format(pct_repurch_existing_buyers)} of existing Campaign Participants repeated on promoted products post VMR reward trip'

sld_7.placeholders[16].text = f'Campaign Period: {str(reward_print_start).replace("-","/")} - {str(reward_print_end).replace("-","/")} \nPost Period: {str(post_4wk_start).replace("-","/")} - {str(post_4wk_end).replace("-","/")} \nRepeat Period = Campaign Period + Post Period' 

phr = sld_7.placeholders[16]
# Modify the font name and size
font_name = 'Arial Black'
font_size = Pt(8)
# Access the text frame and text properties
text_frame = phr.text_frame
for paragraph in text_frame.paragraphs:
    paragraph.font.name = font_name
    paragraph.font.size = font_size
    
print('slide7')

# In[121]:



#Adding Slide 8: VMR Campaigns Trends by Segment

if nbr_segments > 1:

    sld_8 = prs.slides.add_slide(prs.slide_layouts[7])

    #Chart 8.1: CLustered bar

    chart_data = ChartData()
    chart_data.categories = table_slide.loc[table_slide['Segment']!='Total'].Segment.iloc[1:len(table_slide)]

    chart_data.add_series('Units per Trip', table_slide.loc[table_slide['Segment']!='Total']['% Change Units'].iloc[1:len(table_slide)], '0%')
    chart_data.add_series('Dollars per Trip', table_slide.loc[table_slide['Segment']!='Total']['% Change Dollars'].iloc[1:len(table_slide)], '0%')

    chart8_1 = sld_8.placeholders[11].insert_chart(XL_CHART_TYPE.BAR_CLUSTERED, chart_data).chart
    chart8_1.font.size  = Pt(12)
    chart8_1.category_axis.tick_labels.font.bold = True
    chart8_1.has_legend = True
    chart8_1.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart8_1.legend.include_in_layout = False
    chart8_1.has_title = False
    chart8_1.value_axis.has_major_gridlines = False

    leg = chart8_1.legend
    if leg:
        leg.font.bold = True


    chart8_1.plots[0].has_data_labels = True
    data_labels_1 = chart8_1.plots[0].data_labels
    data_labels_1.font.size = Pt(14)
    data_labels_1.font.bold = True
    data_labels_1.font.color.rgb = RGBColor(0, 32, 96)
    data_labels_1.number_format = '0%'
    data_labels_1.position = XL_LABEL_POSITION.OUTSIDE_END


    #Replacing color for each independent serie:

    bar_colors = [RGBColor(0, 56, 104), RGBColor(240, 192, 86)]

    for i,serie in enumerate(chart8_1.series):
        chart8_1.series[i].format.fill.solid()
        chart8_1.series[i].format.fill.fore_color.rgb = bar_colors[i]



    #Other option to turn off the revert negative values option:

    from pptx.oxml.xmlchemy import OxmlElement

    # Iterate over each series and point in the chart
    for serie in chart8_1.series:
        for point in serie.points:
            # Create the new XML element
            element = OxmlElement('c:invertIfNegative')
            # Set the attribute 'val' to '0'
            element.set('val', '0')
            # Append the new element to the point's format element
            point.format.element.append(element)


    #Table 8.2:


    tph1_1 = sld_8.placeholders[12].insert_table(rows = table_slide_s.shape[0]-1, cols = table_slide_s.shape[1])
    tph1_1._element.graphic.graphicData.tbl[0][-1].text = '{69012ECD-51FC-41F1-AA8D-1B2483CD663E}' 
    table1_1 = tph1_1.table


    #populate slide 1 table 1
    for i in range(table_slide_s.shape[0]-1):
        for j in range(table_slide_s.shape[1]):
            cell = table1_1.cell(i, j)

            if i!= 0 and j == 1:
                cell.text = "${:,}".format(table_slide_s.iloc[i,j])
            elif i!= 0 and (j == 2 or j==4):
                cell.text = "{:.0%}".format(table_slide_s.iloc[i,j])
            else:
                cell.text = str(table_slide_s.iloc[i,j])

            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


    #Text placeholder

    highest = x_high[x_high['Segment']!='Total'].sort_values(by='Campaign Dollars per Trip', ascending=False)
    highest = highest['Segment'].values[0]

    sld_8.placeholders[10].text = f'{highest} segment had the highest dollars per trip during campaign period vs YAGO period'
    sld_8.placeholders[16].text = f'Campaign Period: {str(reward_print_start).replace("-","/")} - {str(reward_print_end).replace("-","/")} \nYAGO Period: {str(yago_print_start).replace("-","/")} - {str(yago_print_end).replace("-","/")}' 

    phr = sld_8.placeholders[16]
    # Modify the font name and size
    font_name = 'Arial Black'
    font_size = Pt(8)
    # Access the text frame and text properties
    text_frame = phr.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = font_size

    print('slide8')
# In[122]:



#Adding Slide 9: Combinations

if nbr_segments > 1:
    sld_9 = prs.slides.add_slide(prs.slide_layouts[8])


    #Chart 9.1: 

    combos_sum = []
    combos_single = []
    conditional_colors = []


    for i in range(len(combos_slide)):

        if '+' in combos_slide['Segment'][i]:
            combos_sum.append(combos_slide['Pct Of Trips'][i])
        else:
            combos_single.append(combos_slide['Pct Of Trips'][i])


    chart_data = ChartData()
    chart_data.categories = combos_slide['Segment']
    chart_data.add_series('', combos_slide['Pct Of Trips'], '0%')

    chart9_2 = sld_9.placeholders[11].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
    chart9_2.font.size  = Pt(12)
    chart9_2.category_axis.tick_labels.font.bold = True
    chart9_2.value_axis.has_major_gridlines = False


    chart9_2.plots[0].has_data_labels = True
    data_labels_1 = chart9_2.plots[0].data_labels
    data_labels_1.font.size = Pt(14)
    data_labels_1.font.bold = True
    data_labels_1.number_format = '0%'


    for combination in combos_slide['Segment']:   

        if '+' in combination:
            conditional_colors.append(RGBColor(240, 192, 86))
        else:
            conditional_colors.append(RGBColor(0, 51, 104))


    colors = conditional_colors


    for index, series in enumerate(chart9_2.series):
        for point, color in zip(series.points, colors):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = color


    #Text placeholder

    pct_new_slide = float(nbr_trips_two_segments/nbr_total_trips)

    sld_9.placeholders[10].text = f'{"{:.0%}".format(pct_new_slide)} of reward baskets included 2 or more segments'

    sld_9.placeholders[16].text = f'Campaign Period: {str(reward_print_start).replace("-","/")} - {str(reward_print_end).replace("-","/")}' 

    
    phr = sld_9.placeholders[16]
    # Modify the font name and size
    font_name = 'Arial Black'
    font_size = Pt(10)
    # Access the text frame and text properties
    text_frame = phr.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = font_size
        
        
    print('slide9')


# In[123]:



#Adding Slide 10: New Segment Purchase

if nbr_segments > 1:
    sld_10 = prs.slides.add_slide(prs.slide_layouts[9])

    #Chart 10.1: NEW SEGMENT

    chart_data = ChartData()
    chart_data.categories = ['Did Not Buy a New Segment', 'Tried a New Segment']

    chart_data.add_series('', (1-df_plus_one.loc[1, 'Pct Of Shoppers'], df_plus_one.loc[1, 'Pct Of Shoppers']), '0%')


    chart10_1 = sld_10.placeholders[11].insert_chart(XL_CHART_TYPE.PIE, chart_data).chart
    chart10_1.font.size  = Pt(12)
    chart10_1.has_legend = True
    chart10_1.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart10_1.legend.include_in_layout = False
    chart10_1.has_title = False
    
    leg = chart10_1.legend
    if leg:
        leg.font.bold = True


    chart10_1.plots[0].has_data_labels = True
    data_labels_1 = chart10_1.plots[0].data_labels
    data_labels_1.font.size = Pt(16)
    data_labels_1.font.bold = True
    data_labels_1.font.color.rgb = RGBColor(255, 255, 255)
    data_labels_1.number_format = '0%'
    data_labels_1.position = XL_LABEL_POSITION.CENTER


    #Chart 10.2: NEW SEGMENT 2


    pct_data_new_part = (new_seg_pct.new_seg_pct)*100
    segments_new_part =  new_seg_pct.segment

    chart_data = ChartData()
    chart_data.categories = segments_new_part

    chart_data.add_series('', pct_data_new_part/100, '0%')


    chart10_2 = sld_10.placeholders[12].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
    chart10_2.font.size  = Pt(12)
    chart10_2.category_axis.tick_labels.font.bold = True
    chart10_2.value_axis.has_major_gridlines = False
    chart10_2.value_axis.visible = False


    chart10_2.plots[0].has_data_labels = True
    data_labels_1 = chart10_2.plots[0].data_labels
    data_labels_1.font.size = Pt(14)
    data_labels_1.font.bold = True
    data_labels_1.number_format = '0%'

    for series in chart10_2.series:
        for point in series.points:
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor(240, 192, 86)


    #Text placeholder

    pct_new_slide = df_plus_one.loc[1, 'Pct Of Shoppers']

    sld_10.placeholders[10].text = f'{"{:.0%}".format(pct_new_slide)} of participants who were previous buyers of campaign promoted products tried a NEW segment'
    sld_10.placeholders[16].text = f'Campaign Period: {str(reward_print_start).replace("-","/")} - {str(reward_print_end).replace("-","/")}' 
    
    phr = sld_10.placeholders[16]
    # Modify the font name and size
    font_name = 'Arial Black'
    font_size = Pt(10)
    # Access the text frame and text properties
    text_frame = phr.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = font_size
        
        
    print('slide10')


# In[124]:



#Adding Slide 11: Basket Sizes
sld_11 = prs.slides.add_slide(prs.slide_layouts[5])

#Chart 11.1: Bar Chart

chart_data = ChartData()
chart_data.categories = ['All Other Baskets during VMR Reward Period', 'VMR Reward Baskets']

chart_data.add_series('', (round(df_vmr_ttls_basket.loc[1, 'Average Basket Size'],2), round(df_vmr_ttls_basket.loc[0, 'Average Basket Size'],2)), '$#,##0')


chart11_1 = sld_11.placeholders[11].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
chart11_1.font.size  = Pt(12)
chart11_1.category_axis.tick_labels.font.bold = True
chart11_1.value_axis.has_major_gridlines = False
#chart11_1.series[0].overlap = -30


chart11_1.plots[0].has_data_labels = True
data_labels_1 = chart11_1.plots[0].data_labels
data_labels_1.font.size = Pt(14)
data_labels_1.font.bold = True
data_labels_1.number_format = '$#,##0.00'


bar_colors = [RGBColor(127, 169, 184), RGBColor(0, 56, 104)]

# Iterate through the bars and set colors
for i, series in enumerate(chart11_1.series):
    for j, point in enumerate(series.points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = bar_colors[j]


#Text placeholder

pct_1 = "${:,.2f}".format(nbr_reward_promoted_trips['avg_basket_size'].values[0])
sld_11.placeholders[13].text = f'{pct_1}'




#Chart 11.2: Bar Chart

if nbr_redemp_promoted['trips'].values[0] != 0:

    chart_data = ChartData()
    chart_data.categories = ['All Other Baskets during VMR Redemption Period', 'VMR Redemption Baskets']

    chart_data.add_series('', (round(df_vmr_ttls_basket.loc[4, 'Average Basket Size'],2), round(df_vmr_ttls_basket.loc[3, 'Average Basket Size'],2)), '$#,##0')


    chart11_2 = sld_11.placeholders[12].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data).chart
    chart11_2.font.size  = Pt(12)
    chart11_2.category_axis.tick_labels.font.bold = True
    chart11_2.value_axis.has_major_gridlines = False


    chart11_2.plots[0].has_data_labels = True
    data_labels_1 = chart11_2.plots[0].data_labels
    data_labels_1.font.size = Pt(14)
    data_labels_1.font.bold = True
    data_labels_1.number_format = '$#,##0.00'


    bar_colors = [RGBColor(127, 169, 184), RGBColor(0, 56, 104)]

    # Iterate through the bars and set colors
    for i, series in enumerate(chart11_2.series):
        for j, point in enumerate(series.points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = bar_colors[j]


    #Text placeholder

    pct_2 = round(nbr_redemp_promoted['trips'].values[0]/nbr_redemp_trips['trips'].values[0]*100)
    sld_11.placeholders[14].text = f'{pct_2}%'


#Text placeholder at bottom

pct_rew_bask = (round(df_vmr_ttls_basket.loc[0, 'Average Basket Size'],2) - round(df_vmr_ttls_basket.loc[1, 'Average Basket Size'],2))/ round(df_vmr_ttls_basket.loc[1, 'Average Basket Size'],2)

if nbr_redemp_promoted['trips'].values[0] != 0:
    pct_red_bask = (round(df_vmr_ttls_basket.loc[3, 'Average Basket Size'],2) - round(df_vmr_ttls_basket.loc[4, 'Average Basket Size'],2))/ round(df_vmr_ttls_basket.loc[4, 'Average Basket Size'],2)

#Text placeholder at bottom

if nbr_redemp_promoted['trips'].values[0] != 0:
    sld_11.placeholders[10].text = f'Reward Baskets were {"{:.0%}".format(pct_rew_bask)} larger than All Other baskets during same time period \nRedemption Baskets were {"{:.0%}".format(pct_red_bask)} larger than All Other baskets during same time period'
else:
    sld_11.placeholders[10].text = f'Reward Baskets were {"{:.0%}".format(pct_rew_bask)} larger than All Other baskets during same time period'
    
campaign_dts = f'''{str(reward_print_start).replace("-","/")} - {str(reward_print_end).replace("-","/")}'''

sld_11.placeholders[16].text = f'Campaign Period: {campaign_dts}' 

phr = sld_11.placeholders[16]
# Modify the font name and size
font_name = 'Arial Black'
font_size = Pt(10)
# Access the text frame and text properties
text_frame = phr.text_frame
for paragraph in text_frame.paragraphs:
    paragraph.font.name = font_name
    paragraph.font.size = font_size
    
    
print('slide11')


# In[125]:


sld_12 = prs.slides.add_slide(prs.slide_layouts[10])

sld_13 = prs.slides.add_slide(prs.slide_layouts[14])

# ### Creating a folder on SharePoint and saving both ouputs in there:

# In[ ]:

# ============================================
# SAVE OUTPUTS TO LOCAL FOLDER
# ============================================
import os
import datetime as dt

program_nm = program_nm[:25]
file_name = program_nm[:25]
last_data_dt = (dt.date.today() - dt.timedelta(days=2))

# Map segment_type to report type prefix
segment_prefix_map = {
    1: "UPC",
    2: "BRAND",
    3: "CATEGORY",
    4: "RETAIL",
    5: "RETAIL",
    6: "Custom"
}
report_type_prefix = segment_prefix_map.get(segment_type, "OTHER")

# Create output directory (inside Docker: /opt/airflow/outputs, maps to project outputs/ folder)
output_base_dir = "/opt/airflow/outputs"
output_dir = os.path.join(output_base_dir, program_nm)
os.makedirs(output_dir, exist_ok=True)

print(f"Output folder: {output_dir}")

# Save Excel workbook (includes report type prefix and run ID for uniqueness)
excel_filename = f'{file_name}_{report_type_prefix}_ID{id_param}_VMR_{str(last_data_dt)}.xlsx'
excel_path = os.path.join(output_dir, excel_filename)
wb.save(excel_path)
print(f'Excel saved: {excel_filename}')

# In[ ]:

print('If this is printed, The Excel output has been saved successfully!')

# In[ ]:

# Save PowerPoint presentation (includes report type prefix and run ID for uniqueness)
pptx_filename = f'{file_name}_{report_type_prefix}_ID{id_param}_VMR_{str(last_data_dt)}.pptx'
pptx_path = os.path.join(output_dir, pptx_filename)
prs.save(pptx_path)
print(f'PPTX saved: {pptx_filename}')

# In[ ]:

print('If this is printed, The PPTX output has been saved successfully!')
